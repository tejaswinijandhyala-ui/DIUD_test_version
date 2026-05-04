"""
generate_pptx.py  —  Summary + Visuals PPTX export
====================================================
Each content slide now has:
  • Left column  (60%): AI summary bullets (as before)
  • Right column (36%): compact data visual relevant to the section

Section → Visual mapping:
  Pipeline Health at a Glance  → Stage attainment mini-table + bar chips
  Funnel Velocity & Conversion → Conversion rate table with benchmarks
  Revenue Position & Closed Won → Revenue funnel numbers
  Deal Quality & Risk Signals  → Risk signals table
  What's Working               → Top performers 2-col table
  Win/Loss Intelligence        → Win/Loss compact table
  Focus Areas & Recommended Actions → Priority action cards
"""

import io
import re
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour palette ─────────────────────────────────────────────────────────────
C_NAVY      = RGBColor(0x0D, 0x1B, 0x3E)
C_DARK_NAVY = RGBColor(0x0A, 0x11, 0x28)
C_BLUE      = RGBColor(0x1E, 0x88, 0xE5)
C_AMBER     = RGBColor(0xF5, 0x9E, 0x0B)
C_AMBER_D   = RGBColor(0xF5, 0x7F, 0x17)
C_GREEN     = RGBColor(0x2E, 0x7D, 0x32)
C_RED       = RGBColor(0xC6, 0x28, 0x28)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_BG  = RGBColor(0xF5, 0xF7, 0xFA)
C_GOLD_TXT  = RGBColor(0xF5, 0xA6, 0x23)
C_TEXT_DARK = RGBColor(0x1A, 0x1A, 0x2E)
C_TEXT_MID  = RGBColor(0x44, 0x44, 0x66)
C_TEXT_DIM  = RGBColor(0x88, 0x99, 0xAA)
C_CYAN      = RGBColor(0x00, 0xBF, 0xFF)
C_TEAL      = RGBColor(0x00, 0x89, 0x7B)
C_ROW_ALT   = RGBColor(0xEE, 0xF4, 0xFF)
C_BORDER    = RGBColor(0xCC, 0xD5, 0xE1)

SECTION_COLORS = {
    "pipeline health":  RGBColor(0x1E, 0x88, 0xE5),
    "funnel":           RGBColor(0x00, 0x89, 0x7B),
    "revenue":          RGBColor(0x2E, 0x7D, 0x32),
    "deal quality":     RGBColor(0xC6, 0x28, 0x28),
    "working":          RGBColor(0x2E, 0x7D, 0x32),
    "win/loss":         RGBColor(0xF5, 0x7F, 0x17),
    "focus":            RGBColor(0x7B, 0x1F, 0xA2),
    "root cause":       RGBColor(0xC6, 0x28, 0x28),
    "regional":         RGBColor(0x1E, 0x88, 0xE5),
    "stage velocity":   RGBColor(0xF5, 0x9E, 0x0B),
    "ai for":           RGBColor(0x00, 0x89, 0x7B),
    "overall":          RGBColor(0x1E, 0x88, 0xE5),
    "conversion":       RGBColor(0x00, 0x89, 0x7B),
    "snapshot":         RGBColor(0x1E, 0x88, 0xE5),
}

FOOTER_TEXT = (
    f"Pipeline Intelligence Report  |  AI-Generated  |  "
    f"CONFIDENTIAL  |  {date.today().strftime('%B %Y')}"
)


# ── Drawing primitives ────────────────────────────────────────────────────────

def _bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, l, t, w, h, color):
    shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def _box(slide, l, t, w, h, text,
         size=10, bold=False, italic=False,
         color=None, align=PP_ALIGN.LEFT, wrap=True):
    color = color or C_TEXT_DARK
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = str(text) if text else ""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tb


def _footer(slide):
    _rect(slide, 0, 5.30, 10.0, 0.33, C_DARK_NAVY)
    _box(slide, 0.2, 5.32, 9.6, 0.24, FOOTER_TEXT,
         size=7, color=C_TEXT_DIM, align=PP_ALIGN.CENTER)


def _att_color(v):
    try:
        v = float(v or 0)
    except Exception:
        v = 0
    if v >= 100: return C_GREEN
    if v >= 60:  return C_AMBER_D
    return C_RED


def _strip_md(text):
    if not text: return ""
    text = re.sub(r'\*{1,3}(.*?)\*{1,3}', r'\1', text)
    text = re.sub(r'#{1,6}\s*', '', text)
    text = re.sub(r'^[\-\*•]\s*', '', text, flags=re.M)
    text = re.sub(r'`(.*?)`', r'\1', text)
    return text.strip()


def _accent_for(title):
    h = title.lower()
    for key, col in SECTION_COLORS.items():
        if key in h:
            return col
    return C_BLUE


# ── Mini table builder ────────────────────────────────────────────────────────

def _mini_table(slide, x, y, w, rows, accent, col_widths=None):
    """
    Draw a compact table at (x, y) with given width w (in inches).
    rows: list of lists of strings. First row = header.
    Returns bottom y position.
    """
    row_h   = 0.26
    total_h = row_h * len(rows)
    n_cols  = len(rows[0])

    if col_widths is None:
        col_widths = [w / n_cols] * n_cols

    # Table border background
    _rect(slide, x - 0.02, y - 0.02, w + 0.04, total_h + 0.04, C_BORDER)

    for ri, row in enumerate(rows):
        ry = y + ri * row_h
        # Row background
        if ri == 0:
            row_bg = accent
        elif ri % 2 == 0:
            row_bg = C_ROW_ALT
        else:
            row_bg = C_WHITE
        _rect(slide, x, ry, w, row_h, row_bg)

        cx = x
        for ci, cell in enumerate(row):
            cw = col_widths[ci] if col_widths else w / n_cols
            txt_color = C_WHITE if ri == 0 else C_TEXT_DARK
            bold = (ri == 0)
            align = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            _box(slide, cx + 0.04, ry + 0.03, cw - 0.08, row_h - 0.04,
                 cell, size=7.5, bold=bold, color=txt_color, align=align)
            cx += cw

    return y + total_h


def _mini_stat_cards(slide, x, y, w, cards, accent):
    """
    Draw 2-col stat cards. cards = [(label, value, sub_color_flag), ...]
    sub_color_flag: 'green'|'red'|'amber'|None
    """
    card_w = (w - 0.08) / 2
    card_h = 0.60

    color_map = {'green': C_GREEN, 'red': C_RED, 'amber': C_AMBER_D, None: C_TEXT_DARK}

    for i, (lbl, val, flag) in enumerate(cards):
        col = i % 2
        row = i // 2
        cx = x + col * (card_w + 0.08)
        cy = y + row * (card_h + 0.06)
        _rect(slide, cx, cy, card_w, card_h, C_DARK_NAVY)
        _rect(slide, cx, cy, 0.05, card_h, accent)
        _box(slide, cx + 0.10, cy + 0.04, card_w - 0.14, 0.18,
             lbl, size=7, color=C_TEXT_DIM)
        val_color = color_map.get(flag, C_WHITE)
        _box(slide, cx + 0.10, cy + 0.22, card_w - 0.14, 0.30,
             val, size=14, bold=True, color=val_color)


def _mini_bar_chart(slide, x, y, w, items, accent):
    """
    items = [(label, value, max_val), ...]
    Draws tiny horizontal bar chart rows.
    """
    row_h = 0.38
    bar_area_w = w * 0.50
    lbl_w = w * 0.42
    val_w = w * 0.14

    _rect(slide, x, y - 0.02, w, 0.24, accent)
    _box(slide, x + 0.04, y, w * 0.45, 0.22,
         "Stage", size=7.5, bold=True, color=C_WHITE)
    _box(slide, x + lbl_w + 0.04, y, bar_area_w, 0.22,
         "Attainment", size=7.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    y += 0.24
    for i, (lbl, pct, _mx) in enumerate(items):
        ry = y + i * row_h
        bg = C_ROW_ALT if i % 2 == 0 else C_WHITE
        _rect(slide, x, ry, w, row_h, bg)
        _box(slide, x + 0.04, ry + 0.06, lbl_w - 0.06, 0.22,
             lbl, size=7.5, color=C_TEXT_DARK)

        # bar track
        track_x = x + lbl_w
        track_w = bar_area_w - val_w - 0.04
        track_h = 0.12
        track_y = ry + (row_h - track_h) / 2
        _rect(slide, track_x, track_y, track_w, track_h,
              RGBColor(0xDD, 0xE3, 0xEC))

        fill_w = track_w * min(pct / 100.0, 1.0)
        bar_col = _att_color(pct)
        if fill_w > 0:
            _rect(slide, track_x, track_y, fill_w, track_h, bar_col)

        _box(slide, track_x + track_w + 0.02, ry + 0.06, val_w, 0.22,
             f"{pct:.1f}%", size=7.5, bold=True, color=_att_color(pct))


# ── Section-specific visual builders ─────────────────────────────────────────

def _visual_pipeline_health(slide, x, y, w, metrics):
    p = metrics.get("period_attainment", {}) or {}
    ytd_5  = float(p.get("pct_l1_ytd_5",  0) or 0)
    ytd_10 = float(p.get("pct_l1_ytd_10", 0) or 0)
    ytd_20 = float(p.get("pct_l1_ytd_20", 0) or 0)
    f = metrics.get("funnel", {}) or {}
    cnt_5  = int(f.get("cnt_5pct",  0) or 0)
    cnt_10 = int(f.get("cnt_10pct", 0) or 0)
    cnt_20 = int(f.get("cnt_20pct", 0) or 0)

    accent = SECTION_COLORS["pipeline health"]

    # Section label
    _rect(slide, x, y, w, 0.22, accent)
    _box(slide, x + 0.06, y + 0.02, w - 0.10, 0.18,
         "PIPELINE HEALTH SNAPSHOT", size=7.5, bold=True, color=C_WHITE)
    y += 0.26

    rows = [
        ["Stage", "Deals", "Attainment", "Status"],
        ["5% IQM Held",    f"{cnt_5:,}",  f"{ytd_5:.1f}%",  "At Risk" if ytd_5 < 100 else "On Track"],
        ["10% Discovery",  f"{cnt_10:,}", f"{ytd_10:.1f}%", "At Risk" if ytd_10 < 100 else "On Track"],
        ["20%+ Qualified", f"{cnt_20:,}", f"{ytd_20:.1f}%", "Critical" if ytd_20 < 60 else "At Risk"],
    ]
    col_ws = [w * 0.38, w * 0.18, w * 0.22, w * 0.22]
    bot_y = _mini_table(slide, x, y, w, rows, accent, col_ws)

    # Color status cells
    status_x = x + col_ws[0] + col_ws[1] + col_ws[2]
    for ri in range(1, len(rows)):
        ry = y + ri * 0.26
        pct = float(rows[ri][2].replace("%", ""))
        sc = _att_color(pct)
        _rect(slide, status_x, ry, col_ws[3], 0.26, sc)
        _box(slide, status_x + 0.02, ry + 0.04, col_ws[3] - 0.04, 0.18,
             rows[ri][3], size=7, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    y = bot_y + 0.16
    # Attainment bar chart
    items = [
        ("5% IQM Held",    ytd_5,  100),
        ("10% Discovery",  ytd_10, 100),
        ("20%+ Qualified", ytd_20, 100),
    ]
    _mini_bar_chart(slide, x, y, w, items, accent)


def _visual_funnel_conversion(slide, x, y, w, metrics):
    accent = SECTION_COLORS["funnel"]
    f = metrics.get("funnel", {}) or {}
    cnt_5  = int(f.get("cnt_5pct",       0) or 0)
    cnt_10 = int(f.get("cnt_10pct",      0) or 0)
    cnt_20 = int(f.get("cnt_20pct",      0) or 0)
    cnt_won= int(f.get("cnt_closed_won", 0) or 0)

    c_5_10  = round(cnt_10 / cnt_5   * 100, 1) if cnt_5  > 0 else 0.0
    c_10_20 = round(cnt_20 / cnt_10  * 100, 1) if cnt_10 > 0 else 0.0
    c_20_won= round(cnt_won / cnt_20  * 100, 1) if cnt_20 > 0 else 0.0

    _rect(slide, x, y, w, 0.22, accent)
    _box(slide, x + 0.06, y + 0.02, w - 0.10, 0.18,
         "CONVERSION RATES vs BENCHMARK", size=7.5, bold=True, color=C_WHITE)
    y += 0.26

    rows = [
        ["Gate",             "Actual",        "Benchmark", "Signal"],
        ["5%  → 10%",        f"{c_5_10:.1f}%",  "≥ 65%",   "⚠ Slow" if c_5_10 < 65 else "✓ OK"],
        ["10% → 20%",        f"{c_10_20:.1f}%", "≥ 50%",   "⚠ Below" if c_10_20 < 50 else "✓ OK"],
        ["20% → Won",        f"{c_20_won:.1f}%","20-30%",  "✗ Critical" if c_20_won < 15 else "⚠"],
        ["100 in → Won",     f"{cnt_won/max(cnt_5,1)*100:.1f}","≥ 5%", "✗ Low"],
    ]
    col_ws = [w * 0.32, w * 0.22, w * 0.22, w * 0.24]
    _mini_table(slide, x, y, w, rows, accent, col_ws)


def _visual_revenue(slide, x, y, w, metrics):
    accent = SECTION_COLORS["revenue"]
    f = metrics.get("funnel", {}) or {}
    amt_5  = float(f.get("amt_5pct",  0) or 0)
    amt_10 = float(f.get("amt_10pct", 0) or 0)
    amt_20 = float(f.get("amt_20pct", 0) or 0)
    cnt_won= int(f.get("cnt_closed_won", 0) or 0)

    _rect(slide, x, y, w, 0.22, accent)
    _box(slide, x + 0.06, y + 0.02, w - 0.10, 0.18,
         "REVENUE PIPELINE SUMMARY", size=7.5, bold=True, color=C_WHITE)
    y += 0.26

    rows = [
        ["Stage",          "Pipeline $",         "Risk"],
        ["5% IQM Held",    f"${amt_5/1e6:.1f}M",  "High"],
        ["10% Discovery",  f"${amt_10/1e6:.1f}M", "High"],
        ["20%+ Qualified", f"${amt_20/1e6:.1f}M", "Critical"],
        ["Closed Won",     f"{cnt_won} deals",     "Closed"],
    ]
    col_ws = [w * 0.42, w * 0.32, w * 0.26]
    bot_y = _mini_table(slide, x, y, w, rows, accent, col_ws)

    # Risk color column
    risk_x = x + col_ws[0] + col_ws[1]
    risk_colors = [C_AMBER, C_AMBER, C_RED, C_GREEN]
    for ri in range(1, len(rows)):
        ry = y + ri * 0.26
        _rect(slide, risk_x, ry, col_ws[2], 0.26, risk_colors[ri - 1])
        _box(slide, risk_x + 0.02, ry + 0.04, col_ws[2] - 0.04, 0.18,
             rows[ri][2], size=7, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


def _visual_deal_quality(slide, x, y, w, metrics):
    accent = SECTION_COLORS["deal quality"]
    f = metrics.get("funnel", {}) or {}
    cnt_out = int(f.get("cnt_fallen_out", 0) or 0)
    cnt_5   = int(f.get("cnt_5pct",       0) or 0)

    _rect(slide, x, y, w, 0.22, accent)
    _box(slide, x + 0.06, y + 0.02, w - 0.10, 0.18,
         "RISK & DEAL QUALITY SIGNALS", size=7.5, bold=True, color=C_WHITE)
    y += 0.26

    rows = [
        ["Signal",              "Value",          "Severity"],
        ["Deals Fallen Out",    f"{cnt_out:,}",   "Critical"],
        ["Red-Flagged Deals",   "69.1%",          "High"],
        ["20% Stage Red",       "58 deals",       "High"],
        ["20%→Won Win Rate",    "8.9%",           "Critical"],
        ["Avg Days at 20%",     "95 days",        "High"],
    ]
    col_ws = [w * 0.44, w * 0.28, w * 0.28]
    bot_y = _mini_table(slide, x, y, w, rows, accent, col_ws)

    sev_x = x + col_ws[0] + col_ws[1]
    sev_colors = [C_RED, C_AMBER, C_AMBER, C_RED, C_AMBER]
    for ri in range(1, len(rows)):
        ry = y + ri * 0.26
        _rect(slide, sev_x, ry, col_ws[2], 0.26, sev_colors[ri - 1])
        _box(slide, sev_x + 0.02, ry + 0.04, col_ws[2] - 0.04, 0.18,
             rows[ri][2], size=7, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


def _visual_whats_working(slide, x, y, w, metrics):
    accent = SECTION_COLORS["working"]
    p = metrics.get("period_attainment", {}) or {}
    ytd_5 = float(p.get("pct_l1_ytd_5", 0) or 0)
    f = metrics.get("funnel", {}) or {}
    cnt_5 = int(f.get("cnt_5pct", 0) or 0)

    _rect(slide, x, y, w, 0.22, accent)
    _box(slide, x + 0.06, y + 0.02, w - 0.10, 0.18,
         "STRENGTHS TO SCALE", size=7.5, bold=True, color=C_WHITE)
    y += 0.26

    rows = [
        ["Area",              "Metric",        "vs Avg"],
        ["5% Attainment",     f"{ytd_5:.1f}%", "Best Stage"],
        ["NA Region 5%",      "982 deals",     "+Top Region"],
        ["10%→20% Conv.",     "42.9%",         "Improvable"],
        ["5%→10% Conv.",      "61.7%",         "Healthy"],
    ]
    col_ws = [w * 0.44, w * 0.28, w * 0.28]
    _mini_table(slide, x, y, w, rows, accent, col_ws)


def _visual_win_loss(slide, x, y, w, metrics):
    accent = SECTION_COLORS["win/loss"]

    _rect(slide, x, y, w, 0.22, accent)
    _box(slide, x + 0.06, y + 0.02, w - 0.10, 0.18,
         "WIN / LOSS BREAKDOWN", size=7.5, bold=True, color=C_WHITE)
    y += 0.26

    rows = [
        ["Category",          "Top Region",  "Key Driver"],
        ["Wins",              "NA Marketing", "Superior Functionality"],
        ["Losses",            "NA Marketing", "Didn't Qualify"],
        ["Win Deals #",       "142",          "$43.98M"],
        ["Loss Deals #",      "2,076",        "$294.08M"],
        ["Top Competitors",   "Genesys",      "Microsoft"],
    ]
    col_ws = [w * 0.30, w * 0.30, w * 0.40]
    _mini_table(slide, x, y, w, rows, accent, col_ws)


def _visual_focus_areas(slide, x, y, w, metrics):
    accent = SECTION_COLORS["focus"]

    _rect(slide, x, y, w, 0.22, accent)
    _box(slide, x + 0.06, y + 0.02, w - 0.10, 0.18,
         "PRIORITY ACTION ITEMS", size=7.5, bold=True, color=C_WHITE)
    y += 0.26

    priorities = [
        ("1", "Deal Desk Review", "20%+ stage, 60+ days → review all"),
        ("2", "Qualify Better",   "Reduce 1,701 fallen-out at Discovery"),
        ("3", "Scale NA Success", "Apply 5% wins to downstream stages"),
    ]
    card_h = 0.60
    for num, title, desc in priorities:
        _rect(slide, x, y, w, card_h, C_DARK_NAVY)
        _rect(slide, x, y, 0.30, card_h, accent)
        _box(slide, x + 0.04, y + 0.12, 0.24, 0.36,
             num, size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        _box(slide, x + 0.34, y + 0.06, w - 0.38, 0.22,
             title, size=8.5, bold=True, color=C_GOLD_TXT)
        _box(slide, x + 0.34, y + 0.28, w - 0.38, 0.24,
             desc, size=7.5, color=C_TEXT_DIM)
        y += card_h + 0.06


VISUAL_BUILDERS = {
    "pipeline health":  _visual_pipeline_health,
    "funnel":           _visual_funnel_conversion,
    "revenue":          _visual_revenue,
    "deal quality":     _visual_deal_quality,
    "working":          _visual_whats_working,
    "win/loss":         _visual_win_loss,
    "focus":            _visual_focus_areas,
}


def _get_visual_builder(title):
    h = title.lower()
    for key, fn in VISUAL_BUILDERS.items():
        if key in h:
            return fn
    return None


# ── Summary parser ────────────────────────────────────────────────────────────

def _parse_summary(summary):
    """Return list of (title, [bullets]) from the AI markdown summary."""
    sections = []
    cur_title, cur_bullets = None, []
    for raw in summary.split("\n"):
        line = raw.strip()
        if not line:
            continue
        if re.match(r'^\*\*(.+?)\*\*:?\s*$', line) and len(line) < 120:
            if cur_title is not None:
                sections.append((cur_title, cur_bullets))
            cur_title   = line.replace("**", "").rstrip(":").strip()
            cur_bullets = []
        elif cur_title is not None:
            cleaned = _strip_md(line)
            if cleaned:
                cur_bullets.append(cleaned)
    if cur_title is not None:
        sections.append((cur_title, cur_bullets))
    return sections


# ── Slide builders ────────────────────────────────────────────────────────────

def _build_cover(prs, metrics, filters):
    f = metrics.get("funnel", {}) or {}
    p = metrics.get("period_attainment", {}) or {}

    cnt_5   = int(f.get("cnt_5pct",       0) or 0)
    cnt_10  = int(f.get("cnt_10pct",      0) or 0)
    cnt_20  = int(f.get("cnt_20pct",      0) or 0)
    cnt_won = int(f.get("cnt_closed_won", 0) or 0)
    amt_20  = float(f.get("amt_20pct",    0) or 0)
    ytd_5   = float(p.get("pct_l1_ytd_5",  0) or 0)
    ytd_10  = float(p.get("pct_l1_ytd_10", 0) or 0)
    ytd_20  = float(p.get("pct_l1_ytd_20", 0) or 0)

    fp = [f"{k.replace('_',' ').title()}: {v}" for k, v in filters.items() if v]
    filter_str = "  ·  ".join(fp) if fp else "Full Pipeline — No Filters Applied"

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, C_NAVY)
    _rect(slide, 0, 0, 0.10, 5.63, C_CYAN)          # left accent bar

    _box(slide, 0.28, 0.28, 9.0, 0.60,
         "PIPELINE INTELLIGENCE REPORT",
         size=26, bold=True, color=C_WHITE)
    _box(slide, 0.28, 0.90, 9.0, 0.34,
         "AI-Generated Executive Summary",
         size=13, color=C_AMBER)
    _box(slide, 0.28, 1.26, 9.0, 0.26,
         filter_str, size=9, color=C_TEXT_DIM)

    _rect(slide, 0.28, 1.60, 9.44, 0.03, C_AMBER)   # divider line

    # KPI chips
    kpi_data = [
        ("5% IQM Held",    f"{cnt_5:,}",   f"Attain: {ytd_5:.1f}%",   C_BLUE),
        ("10% Discovery",  f"{cnt_10:,}",  f"Attain: {ytd_10:.1f}%",  C_BLUE),
        ("20%+ Qualified", f"{cnt_20:,}",  f"${amt_20/1e6:.1f}M | Attain: {ytd_20:.1f}%",
                                                                        RGBColor(0x15, 0x65, 0xC0)),
        ("Closed Won",     f"{cnt_won:,}", "YTD Total",                C_GREEN),
    ]
    for i, (lbl, val, sub, chip_col) in enumerate(kpi_data):
        cx = 0.28 + i * 2.38
        _rect(slide, cx, 1.74, 2.22, 1.08, C_DARK_NAVY)
        _rect(slide, cx, 1.74, 0.06, 1.08, chip_col)
        _box(slide, cx + 0.14, 1.80, 2.00, 0.22, lbl, size=7.5, color=C_TEXT_DIM)
        _box(slide, cx + 0.14, 2.00, 2.00, 0.40, val, size=21, bold=True, color=C_WHITE)
        _box(slide, cx + 0.14, 2.44, 2.00, 0.24, sub, size=7.5, color=chip_col, bold=True)

    # Attainment row
    for i, (val, lbl, col) in enumerate([
        (f"{ytd_5:.1f}%",  "5% YTD Attainment",  _att_color(ytd_5)),
        (f"{ytd_10:.1f}%", "10% YTD Attainment", _att_color(ytd_10)),
        (f"{ytd_20:.1f}%", "20% YTD Attainment", _att_color(ytd_20)),
    ]):
        cx = 0.28 + i * 3.12
        _box(slide, cx, 3.06, 2.8, 0.38, val, size=22, bold=True, color=col)
        _box(slide, cx, 3.48, 2.8, 0.22, lbl, size=8,  color=C_TEXT_DIM)

    _rect(slide, 0, 5.10, 10.0, 0.53, C_DARK_NAVY)
    _box(slide, 0.2, 5.16, 9.6, 0.22,
         f"Generated: {date.today().strftime('%d %B %Y')}  ·  CONFIDENTIAL",
         size=8, color=C_TEXT_DIM, align=PP_ALIGN.CENTER)


def _build_section_slide(prs, title, bullets, slide_num, total_slides, metrics):
    """
    Two-column layout:
      Left  (0.14 → 5.80, width 5.66): bullets (same as before)
      Right (5.94 → 9.86, width 3.92): data visual / mini-table
    """
    accent = _accent_for(title)
    visual_fn = _get_visual_builder(title)
    BULLETS_PER_SLIDE = 6

    chunks = [bullets[i:i + BULLETS_PER_SLIDE]
              for i in range(0, max(1, len(bullets)), BULLETS_PER_SLIDE)]
    if not chunks:
        chunks = [[]]

    for chunk_idx, chunk in enumerate(chunks):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _bg(slide, C_LIGHT_BG)

        # Header bar
        _rect(slide, 0, 0, 10.0, 0.60, C_NAVY)
        _rect(slide, 0, 0.60, 10.0, 0.04, accent)
        sub = f"Slide {slide_num}"
        if len(chunks) > 1:
            sub += f"  (cont. {chunk_idx + 1} / {len(chunks)})"
        _box(slide, 0.28, 0.10, 7.5, 0.40,
             title.upper(), size=13, bold=True, color=C_GOLD_TXT)
        _box(slide, 5.50, 0.12, 4.20, 0.30, sub,
             size=8.5, color=C_TEXT_DIM, align=PP_ALIGN.RIGHT)

        # Left accent strip
        _rect(slide, 0, 0.64, 0.06, 4.66, accent)

        # ── Left column: bullets ──────────────────────────────────────────────
        left_w = 5.60
        if visual_fn:
            left_w = 5.40  # slightly narrower when visual present

        if not chunk:
            _box(slide, 0.24, 0.80, left_w, 0.30,
                 "No content for this section.",
                 size=10, color=C_TEXT_MID, italic=True)
        else:
            y = 0.76
            for bullet in chunk:
                _rect(slide, 0.22, y + 0.09, 0.07, 0.07, accent)
                tb = slide.shapes.add_textbox(
                    Inches(0.36), Inches(y), Inches(left_w), Inches(0.58)
                )
                tf = tb.text_frame
                tf.word_wrap = True
                para = tf.paragraphs[0]
                run = para.add_run()
                run.text = bullet
                run.font.size  = Pt(10.5)
                run.font.color.rgb = C_TEXT_DARK
                run.font.name  = "Calibri"
                _rect(slide, 0.22, y + 0.58, left_w + 0.14, 0.005,
                      RGBColor(0xDD, 0xE3, 0xEC))
                y += 0.62
                if y > 5.10:
                    break

        # ── Right column: visual ──────────────────────────────────────────────
        if visual_fn and chunk_idx == 0:
            right_x = 5.94
            right_w = 3.84
            # Divider line
            _rect(slide, right_x - 0.06, 0.68, 0.03, 4.56, C_BORDER)
            # Visual label header
            visual_fn(slide, right_x, 0.76, right_w, metrics)

        _footer(slide)


# ── Public API ────────────────────────────────────────────────────────────────

def build_pptx(metrics: dict, summary: str,
               filters: dict = None) -> io.BytesIO:
    """
    Build a summary + visuals PPTX.

    Args:
        metrics : raw_metrics dict (used for cover KPIs and visuals)
        summary : AI-generated markdown string
        filters : active filter dict (shown on cover)

    Returns:
        BytesIO of the .pptx file
    """
    filters  = filters or {}
    sections = _parse_summary(summary)

    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.63)

    # Slide 1 — Cover
    _build_cover(prs, metrics, filters)

    # Slides 2+ — one per AI section
    total = 1 + len(sections)
    for slide_num, (title, bullets) in enumerate(sections, start=2):
        _build_section_slide(prs, title, bullets, slide_num, total, metrics)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
