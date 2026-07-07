import csv
import io
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment
from openpyxl.utils import get_column_letter
import json
import os
import re
import traceback
import uuid
from rules import (
    validate_sql_against_rules,
    validate_result_against_rules,
    get_intent,
    get_rulebook_entry,
    validate_summary_against_facts,
)
from datetime import date, datetime, timedelta
from typing import Dict, List, Literal, Optional, Callable, Tuple
from charts import build_chart_html, reply_already_has_chart

import httpx
import anthropic
from fastapi import FastAPI, Header, HTTPException, Depends
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, HTMLResponse, StreamingResponse
from pydantic import BaseModel
from dotenv import load_dotenv

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import (
    BaseDocTemplate, Frame, PageTemplate,
    PageBreak, Paragraph, Spacer, Table, TableStyle,
)

# =============================================================================
# Load ENV
# =============================================================================
load_dotenv()

# =============================================================================
# FastAPI App
# =============================================================================
app = FastAPI(title="DIUD", description="Decision Intelligence Using Data", version="4.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# =============================================================================
# Claude client
# =============================================================================
_ai_client    = anthropic.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))
_DEV_ADMIN_TOKEN = os.getenv("DEV_ADMIN_TOKEN")  # set this in your deployment env


def _require_admin(x_admin_token: Optional[str] = Header(default=None)):
    """
    Real, server-side gate for every developer-only debug endpoint —
    metrics, feedback, alerts, and the whole proposal system. Unlike the
    chat login (a plain client-side JS check with zero backend
    enforcement), this actually rejects the request on the server if the
    token doesn't match. Fails closed if DEV_ADMIN_TOKEN isn't configured
    at all, rather than silently leaving everything open.
    """
    if not _DEV_ADMIN_TOKEN:
        raise HTTPException(
            status_code=503,
            detail="Admin surface is not configured. Set DEV_ADMIN_TOKEN in the environment to enable it.",
        )
    if not x_admin_token or x_admin_token != _DEV_ADMIN_TOKEN:
        raise HTTPException(status_code=401, detail="Invalid or missing admin token.")


_CLAUDE_MODEL = "claude-sonnet-4-6"
ALLOWED_MODELS = {
    "sonnet": "claude-sonnet-4-6",
    "opus":   "claude-opus-4-6",
}

# FIX 1: Increased max_tokens across all call sites
# Chat responses: 4096 → handles long analytical responses without mid-sentence cut
# Export generation: 6144 → handles full reports with large deal tables
CHAT_MAX_TOKENS   = 6144
EXPORT_MAX_TOKENS = 6144

# =============================================================================
# SERVER-SIDE SESSION STORE
# =============================================================================
class QueryResult:
    def __init__(self, sql: str, columns: List[str], rows: List[dict],
                 total_rows: int, captured_at: str, filters_applied: str = ""):
        self.sql             = sql
        self.columns         = columns
        self.rows            = rows
        self.total_rows      = total_rows
        self.captured_at     = captured_at
        self.filters_applied = filters_applied


import threading

_SESSION_STORE:      Dict[str, QueryResult] = {}
_SESSION_TIMESTAMPS: Dict[str, datetime]    = {}

def _store_result(session_id: str, result: QueryResult):
    _SESSION_STORE[session_id]      = result
    _SESSION_TIMESTAMPS[session_id] = datetime.utcnow()
    
def _get_result(session_id: str) -> Optional[QueryResult]:
    """Retrieve stored query result for a session. Returns None if not found."""
    return _SESSION_STORE.get(session_id)

def _cleanup_sessions():
    cutoff  = datetime.utcnow() - timedelta(hours=4)
    expired = [sid for sid, ts in list(_SESSION_TIMESTAMPS.items()) if ts < cutoff]
    for sid in expired:
        _SESSION_STORE.pop(sid, None)
        _SESSION_TIMESTAMPS.pop(sid, None)
    if expired:
        print(f"🧹 Cleaned {len(expired)} expired sessions.")
    threading.Timer(3600, _cleanup_sessions).start()


_ALLOWED_CHART_TYPES = {"funnel", "attainment", "donut", "bar_h"}
_MAX_INSIGHT_NOTE_LEN = 160
_MAX_TITLE_LEN = 80


def _validate_chart_spec(spec: dict, columns: List[str]) -> Tuple[Optional[dict], Optional[str]]:
    """
    Validates a chart spec Claude proposed via choose_chart_spec against the
    ACTUAL columns of the result it's describing. Returns (clean_spec, error).
    clean_spec is None if invalid — callers must fall back to auto-detected
    charting rather than trust anything unvalidated. This is the only gate
    between "Claude describes a chart" and "a chart gets rendered" — nothing
    here is allowed to affect actual chart NUMBERS, only which columns/type
    are used, so there is no path for Claude to inject a wrong figure.
    """
    if not isinstance(spec, dict):
        return None, "spec is not an object"

    chart_type = spec.get("chart_type")
    if chart_type not in _ALLOWED_CHART_TYPES:
        return None, f"chart_type '{chart_type}' is not one of {sorted(_ALLOWED_CHART_TYPES)}"

    label_col = spec.get("label_column")
    if label_col is not None and label_col not in columns:
        return None, f"label_column '{label_col}' is not a real column in this result: {columns}"

    value_col = spec.get("value_column")
    if value_col is not None and value_col not in columns:
        return None, f"value_column '{value_col}' is not a real column in this result: {columns}"

    exclude_values = spec.get("exclude_values")
    if exclude_values is not None:
        if not isinstance(exclude_values, list) or not all(isinstance(v, str) for v in exclude_values):
            return None, "exclude_values must be a list of strings"
        exclude_values = exclude_values[:20]

    title = spec.get("title")
    if title is not None:
        if not isinstance(title, str):
            return None, "title must be a string"
        title = title[:_MAX_TITLE_LEN]

    insight_note = spec.get("insight_note")
    if insight_note is not None:
        if not isinstance(insight_note, str):
            return None, "insight_note must be a string"
        insight_note = insight_note[:_MAX_INSIGHT_NOTE_LEN]

    clean = {
        "chart_type": chart_type,
        "label_column": label_col,
        "value_column": value_col,
        "exclude_values": exclude_values,
        "title": title,
        "insight_note": insight_note,
    }
    return clean, None

_RULE_AUDIT_LOG: List[dict] = []   # most-recent-first, capped
_FEEDBACK_LOG: List[dict] = []     # most-recent-first, capped — thumbs up/down from users

# ── Human-approved-fix middle ground ─────────────────────────────────────
# DIUD can DRAFT a proposed fix for a flagged issue (via _draft_fix_proposal,
# a separate Claude call), but a proposal only ever takes effect after a
# human calls /debug/proposals/{id}/approve. Scoped deliberately narrow:
# a proposal can ONLY be an additive plain-English clarification appended
# to an agent's prompt — never a change to rules.py's actual validation
# logic. A bad sentence is easy to spot and remove; a bad autonomous edit
# to executable guardrail code is a different category of risk.
_PENDING_PROPOSALS: List[dict] = []
_PROPOSAL_HISTORY: List[dict] = []   # approved + rejected, most-recent-first
_SQL_AGENT_PATCHES: List[str] = []       # approved clarifications, SQL Agent prompt
_NARRATOR_AGENT_PATCHES: List[str] = []  # approved clarifications, Narrator prompt

def _log_feedback(session_id: Optional[str], user_question: str,
                   assistant_reply: str, rating: str,
                   issue_type: Optional[str], comment: Optional[str]):
    """
    Records a user's thumbs up/down. Tagged with the same get_intent
    classification used for rule violations, so negative feedback can be
    grouped by question type right alongside rule-violation failures —
    two different signals (did our own rules catch a problem vs. did a
    human actually think the answer was wrong) feeding the same review
    surface, not two disconnected systems.

    This ONLY records the signal. Nothing here changes DIUD's behavior
    automatically — see get_feedback_metrics() / /debug/feedback, which
    surface this for a person to review, deliberately not an auto-tuning
    loop. An agent that quietly rewrites its own guardrails based on
    unverified user reactions (which can be wrong, sarcastic, or testing
    edge cases on purpose) is not something to build silently as a
    side effect of "make it adaptive."
    """
    try:
        intent = get_intent(user_question, sql="")
        pattern_tag = (
            intent.get("pattern_hint")
            or intent.get("pattern")
            or intent.get("metric")
            or (f"stage_{intent['stage']}" if intent.get("stage") else None)
            or (f"cohort_{intent['cohort_stage']}" if intent.get("cohort_stage") else None)
            or "unclassified"
        )
    except Exception:
        pattern_tag = "unclassified"

    entry = {
        "ts": datetime.utcnow().isoformat(),
        "session_id": session_id,
        "pattern": pattern_tag,
        "rating": rating,                      # "up" | "down"
        "issue_type": issue_type or "",         # only meaningful when rating == "down"
        "comment": (comment or "")[:500],
        "user_question": user_question[:200],
        "assistant_reply_preview": assistant_reply[:300],
    }
    _FEEDBACK_LOG.insert(0, entry)
    del _FEEDBACK_LOG[500:]
    print(f"[FEEDBACK][{rating}][{pattern_tag}]"
          + (f"[{issue_type}]" if issue_type else "")
          + f" session={session_id}"
          + (f" comment={comment[:100]!r}" if comment else ""))


def get_feedback_metrics() -> dict:
    """
    Aggregates _FEEDBACK_LOG the same way get_audit_metrics() aggregates
    the rule-violation log: overall up/down rate, broken down by question
    pattern AND by reported issue type, so both "which type of QUESTION do
    users dislike" and "which KIND of mistake keeps happening" are
    answerable without reading raw log entries by eye.
    """
    total = len(_FEEDBACK_LOG)
    if total == 0:
        return {"total_feedback": 0, "thumbs_up": 0, "thumbs_down": 0,
                "positive_rate_pct": None, "by_pattern": {}, "by_issue_type": {},
                "recent_comments": []}

    up = sum(1 for e in _FEEDBACK_LOG if e["rating"] == "up")
    down = total - up

    by_pattern: Dict[str, Dict[str, int]] = {}
    by_issue_type: Dict[str, int] = {}
    for e in _FEEDBACK_LOG:
        p = by_pattern.setdefault(e["pattern"], {"up": 0, "down": 0})
        p[e["rating"]] = p.get(e["rating"], 0) + 1
        if e["rating"] == "down" and e["issue_type"]:
            by_issue_type[e["issue_type"]] = by_issue_type.get(e["issue_type"], 0) + 1

    return {
        "total_feedback": total,
        "thumbs_up": up,
        "thumbs_down": down,
        "positive_rate_pct": round(up / total * 100, 1),
        "by_pattern": dict(sorted(
            by_pattern.items(), key=lambda kv: kv[1].get("down", 0), reverse=True
        )),
        "by_issue_type": dict(sorted(by_issue_type.items(), key=lambda x: -x[1])),
        "recent_comments": [
            {"pattern": e["pattern"], "rating": e["rating"],
             "issue_type": e["issue_type"], "comment": e["comment"]}
            for e in _FEEDBACK_LOG if e["comment"] or e["issue_type"]
        ][:20],
    }


def _log_rule_audit(session_id: Optional[str], sql: str, violations: List[str],
                     stage: str, user_message: str):
    # Tag every entry with the classified query pattern/metric type (not just
    # which rule fired) so failure analysis can answer "which KIND of
    # question fails most" not only "which RULE fires most". Best-effort:
    # classification failing must never block logging the actual violation.
    try:
        intent = get_intent(user_message, sql=sql if sql != "summary_check" else "")
        pattern_tag = (
            intent.get("pattern_hint")
            or intent.get("pattern")
            or intent.get("metric")
            or (f"stage_{intent['stage']}" if intent.get("stage") else None)
            or (f"cohort_{intent['cohort_stage']}" if intent.get("cohort_stage") else None)
            or "unclassified"
        )
    except Exception:
        pattern_tag = "unclassified"

    entry = {
        "ts": datetime.utcnow().isoformat(),
        "session_id": session_id,
        "stage": stage,                # "pre_execute" | "post_execute" | "post_summary"
        "pattern": pattern_tag,
        "violations": violations,
        "sql_preview": sql[:300],
        "user_message": user_message[:200],
    }
    _RULE_AUDIT_LOG.insert(0, entry)
    del _RULE_AUDIT_LOG[200:]          # cap log size
    if violations:
        print(f"[RULE-AUDIT][{stage}][{pattern_tag}] session={session_id} violations={violations}")


def get_audit_metrics() -> dict:
    """
    Aggregates _RULE_AUDIT_LOG into a queryable success/failure report —
    the "Tool Success Rate" and "Failure Analysis" metrics that were
    previously entirely missing. The log was already capturing everything
    needed; nothing ever added it up. This reads the log fresh each call,
    so it always reflects the current in-memory window (last 200 events).
    """
    total = len(_RULE_AUDIT_LOG)
    if total == 0:
        return {
            "total_events": 0, "clean": 0, "with_violations": 0,
            "success_rate_pct": None, "by_rule_id": {}, "by_pattern": {},
            "by_stage": {}, "pattern_rates": {},
        }

    clean = sum(1 for e in _RULE_AUDIT_LOG if not e["violations"])
    with_violations = total - clean

    by_rule_id: Dict[str, int] = {}
    by_pattern: Dict[str, int] = {}
    by_stage: Dict[str, int] = {}
    # Per-pattern TOTALS (attempts + violations), not just violation counts —
    # a pattern with 8 violations out of 10 attempts is a real problem; a
    # pattern with 8 violations out of 200 attempts mostly just gets asked
    # a lot. by_pattern above can't distinguish these; pattern_rates can.
    pattern_totals: Dict[str, Dict[str, int]] = {}

    for e in _RULE_AUDIT_LOG:
        by_stage[e["stage"]] = by_stage.get(e["stage"], 0) + 1
        pt = pattern_totals.setdefault(e["pattern"], {"total": 0, "violations": 0})
        pt["total"] += 1
        if e["violations"]:
            pt["violations"] += 1
            by_pattern[e["pattern"]] = by_pattern.get(e["pattern"], 0) + 1
            for v in e["violations"]:
                # violation strings are formatted "[rule_id] §n ... : message"
                rule_id = v.split("]")[0].lstrip("[") if v.startswith("[") else "unknown"
                by_rule_id[rule_id] = by_rule_id.get(rule_id, 0) + 1

    return {
        "total_events": total,
        "clean": clean,
        "with_violations": with_violations,
        "success_rate_pct": round(clean / total * 100, 1),
        "pattern_rates": pattern_totals,
        "by_rule_id": dict(sorted(by_rule_id.items(), key=lambda x: -x[1])),
        "by_pattern": dict(sorted(by_pattern.items(), key=lambda x: -x[1])),
        "by_stage": by_stage,
    }


def get_flagged_issues(
    rule_violation_threshold: int = 3,
    pattern_failure_rate_threshold: float = 30.0,
    pattern_min_sample: int = 5,
    feedback_negative_rate_threshold: float = 40.0,
    feedback_min_sample: int = 3,
    issue_type_threshold: int = 3,
) -> dict:
    """
    THE CONNECTING PIECE between "we log everything" and "a human notices
    a repeated pattern" — cross-references the rule-violation audit log
    and the user-feedback log, and surfaces anything crossing a repeat
    threshold as a flagged item for review.

    This function NEVER changes DIUD's behavior. It does not touch
    rules.py, does not touch either agent's prompt, does not retry
    anything. It only produces a sorted list of "this looks worth a
    look" — matching the human-supervised design we agreed on: an agent
    that quietly rewrites its own guardrails based on failure counts is
    a real risk for a system whose numbers drive business decisions, so
    the decision of what to actually change stays with a person, every
    time. See /debug/alerts.

    Flags come from five checks, escalating in what they mean:
      1. A specific rule firing repeatedly (raw count) — the rule itself
         may be too strict, or the model keeps making the same mistake.
      2. A question pattern with a high RULE FAILURE RATE (not just a
         raw count — a pattern that's simply asked often will naturally
         accumulate more violations without this being unusual).
      3. A question pattern with a high NEGATIVE FEEDBACK RATE.
      4. A specific reported issue type recurring.
      5. CROSS-SIGNAL: a pattern flagged by BOTH the rule engine AND user
         feedback independently — two different signals agreeing is much
         stronger evidence than either alone, so this is always "critical".
    """
    audit = get_audit_metrics()
    feedback = get_feedback_metrics()
    flags: List[dict] = []

    for rule_id, count in audit["by_rule_id"].items():
        if count >= rule_violation_threshold:
            flags.append({
                "severity": "high" if count >= rule_violation_threshold * 2 else "medium",
                "source": "rule_engine",
                "subject": rule_id,
                "detail": f"Rule '{rule_id}' has fired {count} times in the current audit window.",
                "suggested_action": "Check whether this rule is too strict, or the model keeps making the same mistake.",
            })

    rule_flagged_patterns = set()
    for pattern, stats in audit.get("pattern_rates", {}).items():
        total = stats["total"]
        violations = stats["violations"]
        if total >= pattern_min_sample:
            rate = violations / total * 100
            if rate >= pattern_failure_rate_threshold:
                rule_flagged_patterns.add(pattern)
                flags.append({
                    "severity": "high" if rate >= 50 else "medium",
                    "source": "rule_engine",
                    "subject": pattern,
                    "detail": f"Question pattern '{pattern}' fails our own rules {rate:.0f}% of the time ({violations}/{total} attempts).",
                    "suggested_action": "This question type may need clearer prompt guidance or a rule review.",
                })

    feedback_flagged_patterns = set()
    for pattern, counts in feedback.get("by_pattern", {}).items():
        total = counts.get("up", 0) + counts.get("down", 0)
        if total >= feedback_min_sample:
            neg_rate = counts.get("down", 0) / total * 100
            if neg_rate >= feedback_negative_rate_threshold:
                feedback_flagged_patterns.add(pattern)
                flags.append({
                    "severity": "high" if neg_rate >= 60 else "medium",
                    "source": "user_feedback",
                    "subject": pattern,
                    "detail": f"Users rated '{pattern}' answers negatively {neg_rate:.0f}% of the time ({counts.get('down', 0)}/{total} ratings).",
                    "suggested_action": "Check recent comments for this pattern in /debug/feedback.",
                })

    for issue_type, count in feedback.get("by_issue_type", {}).items():
        if count >= issue_type_threshold:
            flags.append({
                "severity": "medium",
                "source": "user_feedback",
                "subject": issue_type,
                "detail": f"Users have reported '{issue_type}' {count} times.",
                "suggested_action": "Review recent comments tagged with this issue type.",
            })

    for pattern in (rule_flagged_patterns & feedback_flagged_patterns):
        flags.append({
            "severity": "critical",
            "source": "cross_signal",
            "subject": pattern,
            "detail": f"'{pattern}' is failing BOTH our own rules AND getting negative user "
                      f"feedback — two independent signals agree.",
            "suggested_action": "Prioritize this one — it's the strongest signal the system can produce.",
        })

    severity_order = {"critical": 0, "high": 1, "medium": 2, "low": 3}
    flags.sort(key=lambda f: severity_order.get(f["severity"], 4))

    return {
        "flag_count": len(flags),
        "flags": flags,
        "note": ("Review surface only — nothing here changes DIUD's behavior automatically. "
                 "A person decides what, if anything, to act on."),
    }


def _gather_evidence_for_flag(flag: dict) -> str:
    """
    Pulls a few concrete recent real examples backing a flag, so the Rule
    Reviewer drafts a fix grounded in actual cases rather than reasoning
    from a bare count. Deliberately capped at 4 examples per source —
    enough to see a pattern, not so much that it dominates the prompt.
    """
    subject = flag["subject"]
    source = flag["source"]
    lines: List[str] = []

    if source in ("rule_engine", "cross_signal"):
        rule_matches = [
            e for e in _RULE_AUDIT_LOG
            if any(v.startswith(f"[{subject}]") for v in e.get("violations", []))
        ]
        pattern_matches = [
            e for e in _RULE_AUDIT_LOG
            if e.get("pattern") == subject and e.get("violations")
        ]
        for e in (rule_matches or pattern_matches)[:4]:
            lines.append(
                f"- Question: {e['user_message']!r}\n"
                f"  SQL: {e['sql_preview'][:200]!r}\n"
                f"  Violations: {e['violations']}"
            )

    if source in ("user_feedback", "cross_signal"):
        fb_matches = [
            e for e in _FEEDBACK_LOG
            if e.get("pattern") == subject or e.get("issue_type") == subject
        ]
        for e in fb_matches[:4]:
            lines.append(
                f"- Question: {e['user_question']!r}\n"
                f"  Reply preview: {e['assistant_reply_preview'][:200]!r}\n"
                f"  Rating: {e['rating']}"
                + (f" ({e['issue_type']})" if e.get("issue_type") else "")
                + (f"\n  Comment: {e['comment']}" if e.get("comment") else "")
            )

    return "\n".join(lines) if lines else "No detailed examples available — flag is based on aggregate counts only."


_RULE_REVIEWER_SYSTEM_PROMPT = """
You are DIUD's Rule Reviewer — an internal tool, never user-facing.

You are given ONE flagged recurring problem from DIUD's own audit and
feedback logs, plus a few real examples backing it. Your job:

1. Diagnose the likely root cause in plain English (2-3 sentences).
2. Decide which agent's prompt the fix belongs in: "sql_agent" (if this
   is about how SQL gets written) or "narrator_agent" (if this is about
   how the final written answer gets formatted/phrased).
3. Draft ONE short, clear, ADDITIVE clarification (2-4 sentences) to
   append to that agent's prompt to prevent this specific mistake going
   forward — written in the same plain-English instructional style the
   rest of that agent's prompt already uses.

HARD RULES — do not violate these:
- You may only ADD a clarifying instruction. NEVER propose removing,
  weakening, loosening, or contradicting an existing rule — especially
  anything about mandatory filters, casting, or data correctness.
- If you cannot draft a safe, additive fix with real confidence from the
  evidence given, say so honestly in the diagnosis and leave
  PROPOSED_ADDITION empty rather than forcing a suggestion.
- You are not writing code. Do not reference rules.py, Python, or write
  SQL syntax rules directly — write a plain-English instruction the way
  a person would explain it to a colleague.

Respond in EXACTLY this format and nothing else:

DIAGNOSIS: <your diagnosis>
TARGET_PROMPT: sql_agent OR narrator_agent
PROPOSED_ADDITION:
<the clarification text, or leave this blank if you have no safe fix to propose>
"""


def _draft_fix_proposal(flag: dict) -> dict:
    """
    Calls a separate, tightly-scoped Claude call to draft a proposed fix
    for one flagged issue. Returns a proposal dict with status "pending"
    (has a real proposed addition) or "no_safe_fix" (the reviewer itself
    declined to guess). Either way, NOTHING is applied here — see
    /debug/proposals/{id}/approve, which is the only path that ever
    changes a live prompt, and only after a human calls it.
    """
    evidence = _gather_evidence_for_flag(flag)
    user_msg = (
        f"FLAGGED ISSUE:\n"
        f"Source: {flag['source']}\n"
        f"Subject: {flag['subject']}\n"
        f"Detail: {flag['detail']}\n\n"
        f"EVIDENCE:\n{evidence}"
    )
    response = _ai_client.messages.create(
        model=ALLOWED_MODELS["sonnet"],
        system=_RULE_REVIEWER_SYSTEM_PROMPT,
        messages=[{"role": "user", "content": user_msg}],
        temperature=0,
        max_tokens=800,
    )
    text = _extract_text(response.content)

    diagnosis_m = re.search(r"DIAGNOSIS:\s*(.*?)(?=\nTARGET_PROMPT:)", text, re.S)
    target_m = re.search(r"TARGET_PROMPT:\s*(sql_agent|narrator_agent)", text, re.I)
    addition_m = re.search(r"PROPOSED_ADDITION:\s*(.*)", text, re.S)

    diagnosis = diagnosis_m.group(1).strip() if diagnosis_m else text[:400]
    target = target_m.group(1).lower() if target_m else "sql_agent"
    addition = addition_m.group(1).strip() if addition_m else ""

    return {
        "id": str(uuid.uuid4())[:8],
        "created_at": datetime.utcnow().isoformat(),
        "flag": flag,
        "diagnosis": diagnosis,
        "target_prompt": target,
        "proposed_addition": addition,
        "status": "pending" if addition else "no_safe_fix",
    }

# =============================================================================
# ClickHouse HTTP proxy — base helpers
# =============================================================================
def _base_url() -> str:
    return (os.getenv("CLICKHOUSE_API_URL") or "").rstrip("/")

def _token() -> str:
    return os.getenv("CLICKHOUSE_API_TOKEN") or ""

def _auth_headers() -> dict:
    return {
        "Authorization": f"Bearer {_token()}",
        "Content-Type":  "application/json",
    }

FORBIDDEN_KEYWORDS = ["INSERT", "UPDATE", "DELETE", "DROP", "ALTER", "TRUNCATE", "CREATE"]

# =============================================================================
# Schema discovery
# =============================================================================
_LIVE_SCHEMA: dict = {}
_SCHEMA_BLOCK: str = "Schema not yet loaded."


def _proxy_get(path: str) -> dict | list | None:
    base_url = _base_url()
    token    = _token()
    if not base_url or not token:
        return None
    try:
        r = httpx.get(f"{base_url}{path}", headers=_auth_headers(), timeout=20)
        if r.status_code == 200:
            return r.json()
        print(f"   ⚠️  GET {path} → HTTP {r.status_code}: {r.text[:200]}")
        return None
    except Exception as e:
        print(f"   ⚠️  GET {path} → {e}")
        return None


def discover_schema() -> str:
    global _LIVE_SCHEMA, _SCHEMA_BLOCK

    print("🔎 Discovering schema from ClickHouse proxy…")
    databases_raw = _proxy_get("/databases")
    if not databases_raw:
        msg = "⚠️  Could not fetch databases — check CLICKHOUSE_API_URL / CLICKHOUSE_API_TOKEN."
        print(msg); _SCHEMA_BLOCK = msg; return msg

    if isinstance(databases_raw, list) and databases_raw:
        if isinstance(databases_raw[0], str):
            databases = databases_raw
        elif isinstance(databases_raw[0], dict):
            databases = [d.get("name") or d.get("database") or list(d.values())[0] for d in databases_raw]
        else:
            databases = [str(d) for d in databases_raw]
    elif isinstance(databases_raw, dict):
        databases = databases_raw.get("data") or databases_raw.get("databases") or list(databases_raw.values())[0] if databases_raw else []
    else:
        databases = []

    SKIP_DBS = {"system", "information_schema", "INFORMATION_SCHEMA"}
    databases = [d for d in databases if d not in SKIP_DBS]
    print(f"   Databases found: {databases}")

    schema_lines, schema_dict = [], {}

    for db in databases:
        tables_raw = _proxy_get(f"/tables/{db}")
        if not tables_raw:
            continue
        if isinstance(tables_raw, list) and tables_raw:
            tables = tables_raw if isinstance(tables_raw[0], str) else [t.get("name") or t.get("table") or list(t.values())[0] for t in tables_raw]
        elif isinstance(tables_raw, dict):
            tables = tables_raw.get("data") or tables_raw.get("tables") or list(tables_raw.values())[0] if tables_raw else []
        else:
            tables = []

        print(f"   {db}: tables = {tables}")
        for tbl in tables:
            schema_raw = _proxy_get(f"/schema/{db}/{tbl}")
            if not schema_raw:
                schema_lines.append(f"\nTABLE: {db}.{tbl}\n  (schema unavailable)")
                continue
            if isinstance(schema_raw, list):
                cols = schema_raw
            elif isinstance(schema_raw, dict):
                cols = schema_raw.get("columns") or schema_raw.get("data") or schema_raw.get("schema") or [schema_raw]
            else:
                cols = []

            schema_dict[f"{db}.{tbl}"] = cols
            col_lines = []
            for col in cols:
                if isinstance(col, dict):
                    col_name    = col.get("name") or col.get("column_name") or col.get("Field") or list(col.keys())[0]
                    col_type    = col.get("type") or col.get("data_type") or col.get("Type") or ""
                    col_comment = col.get("comment") or col.get("Comment") or ""
                    col_lines.append(f"  {col_name:<35} {col_type}" + (f"  — {col_comment}" if col_comment else ""))
                else:
                    col_lines.append(f"  {col}")
            schema_lines.append(f"\nTABLE: {db}.{tbl}")
            schema_lines.extend(col_lines)

    _LIVE_SCHEMA  = schema_dict
    _SCHEMA_BLOCK = "\n".join(schema_lines) if schema_lines else "No tables found."
    print(f"✅ Schema loaded: {list(schema_dict.keys())}")
    return _SCHEMA_BLOCK


# =============================================================================
# System prompt
# =============================================================================
def _build_system_prompt() -> str:
    compact_lines = []
    for table_key, cols in _LIVE_SCHEMA.items():
        col_parts = []
        for col in cols:
            if isinstance(col, dict):
                name = col.get("name") or col.get("column_name") or col.get("Field") or list(col.keys())[0]
                typ  = col.get("type") or col.get("data_type") or col.get("Type") or ""
                col_parts.append(f"{name}:{typ}")
            else:
                col_parts.append(str(col))
        compact_lines.append(f"{table_key}({', '.join(col_parts)})")

    schema = "\n".join(compact_lines) or "Schema not yet loaded."
    if len(schema) > 20000:
        schema = schema[:20000] + "\n[schema truncated]"

    return f"""
You are DIUD (Decision Intelligence Using Data) — a conversational business intelligence assistant.
You convert natural language questions into ClickHouse SQL queries, retrieve live data, and deliver
executive-grade insights. You never fabricate numbers and never run destructive SQL.

RULE PRIORITY ORDER (highest → lowest):
  1. Greeting Rule (§1)
  2. Safety — SELECT/WITH only, no destructive SQL ever
  3. MANDATORY_BASE_FILTERS (§3)
  4. Tool Usage (§2)
  5. Business & SQL Rules (§4–§13)
  6. Formatting & Chart Rules (§14–§15)

═══════════════════════════════════════════════════════════════
§1  GREETING RULE — HIGHEST PRIORITY
═══════════════════════════════════════════════════════════════
If the user's message is ONLY a greeting (hi, hey, hello, good morning, etc.),
respond with EXACTLY this text and nothing else:

  "Hey, I'm DIUD, your data intelligence agent to help you analyse
  the live ClickHouse data. How may I help you?"

No bullet points, no extras. This rule overrides everything below.

═══════════════════════════════════════════════════════════════
§2  TOOL USAGE
═══════════════════════════════════════════════════════════════
You have LIVE access to a ClickHouse database via the query_clickhouse tool.
Use it for any question about pipeline deals, AEs, regions, industries, stages,
win/loss, competitors, conversions, or any metric not already confirmed in the
current conversation context.

NEVER use numbers from memory or a previous conversation turn.
Every metric must be queried live from the database.

DATABASE CONNECTION FAILURE:
If the tool returns "DATABASE CONNECTION FAILED", relay that error directly
to the user without attempting to answer from memory.

EXPORT INTENT RULE:
When the user asks to export, download, create a report/presentation, or
retrieve data in any format, output this marker on its own line:

  __EXPORT_INTENT__

Then on the next line write a brief confirmation such as:
  "Sure! Opening the export panel — select your format and detail level, then click Generate Preview."

This applies to ALL of these phrasings (and similar ones):
  - "export this conversation as a PDF report"
  - "export as PDF" / "export as report"
  - "create a presentation" / "create a deck"
  - "export the list" / "give me those N deals"
  - "download this as CSV" / "save this as a report"

⚠️  NEVER say exporting is "outside your capabilities".
⚠️  NEVER suggest Ctrl+P, screenshots, copy-paste, or contacting an admin.
⚠️  NEVER treat export requests as questions about your features.
The export panel is fully functional — always emit __EXPORT_INTENT__ for any export request.
Do NOT re-run the query. The export panel handles format selection and download.

CHART SPEC RULE:
After query_clickhouse returns a result you're about to summarize, if a
chart would genuinely help (funnel/stage breakdowns, actual-vs-target,
top-N comparisons, breakdowns by region/source/competitor/etc — NOT raw
deal lists or single-number answers), call choose_chart_spec once, using
the EXACT column names from the result you just got back. You are only
choosing chart type, which columns, a title, and an optional one-line
insight already visible in the data — you are NOT computing any numbers,
percentages, or chart proportions yourself; those are always calculated
independently from the real query result. If a category is unlogged/blank
(e.g. "N/A") and would dominate the scale, list it in exclude_values
rather than silently ignoring the issue. Do not call this tool for deal
lists or single-value answers — there's nothing to chart.

═══════════════════════════════════════════════════════════════
§3  SCHEMA, DUPLICATE EXCLUSION, AND MANDATORY BASE FILTERS
═══════════════════════════════════════════════════════════════

── DUPLICATE RECORD EXCLUSION — ALWAYS APPLY ──────────────────
Rule                           | Apply when
-------------------------------|------------------------------------------
FINAL keyword                  | Every query against hs_analytics.* tables
countDistinct(), never count() | All deal/contact aggregations
DISTINCT in subquery           | Association table subqueries
GROUP BY + SUM                 | All target table queries

── TABLE 1: hs_analytics.deals (always use FINAL) ─────────────
PURPOSE: Core deals fact table.
Key columns:
  deal_id, deal_name, deal_owner, deal_stage, deal_type, pipeline,
  amount, region, deal_source_rollup, kore_primary_industry,
  account_priority_level, create_date, close_date,
  became_5_deal_date, became_10_deal_date, became_20_deal_date,
  became_30_deal_date, became_40_deal_date, became_60_deal_date,
  became_75_deal_date

── TABLE 2: hs_analytics.owners (always use FINAL) ────────────
PURPOSE: AE/owner master data.
Columns: id, firstName, lastName, email

── TABLE 3: hs_analytics.companies (always use FINAL) ─────────
PURPOSE: Company/account master data.
Columns: company_id, name, domain, industry, country, city

── TABLE 4: hs_analytics.contacts (always use FINAL) ──────────
PURPOSE: Contact/lead master data.
Columns:
  contact_id, email, first_name, last_name, company_name,
  company_priority, region, original_source, lead_status,
  lifecycle_stage,
  date_entered_marketing_qualified_lead_lifecycle_stage_pipeline

── TABLE 5: kore_ai_hubspot.gs_DealContactAssociation ─────────
PURPOSE: Many-to-many link between contacts and deals.
Columns: contact_id, deal_id
Usage: Always use DISTINCT in subqueries against this table.

── TABLE 6: kore_ai_hubspot.gs_marketing_targets ──────────────
PURPOSE: Marketing MQL and pipeline targets by source.
Columns: fy, quarter, month, region, original_source, mql_target

── TABLE 7: kore_ai_hubspot.gs_deal_ids_hs ────────────────────
PURPOSE: Allowlist of valid deal IDs — used in MANDATORY_BASE_FILTERS.
Columns: deal_id_hs

── MANDATORY BASE FILTERS (apply to EVERY deals query) ────────
Every SQL query against hs_analytics.deals MUST include ALL of these:

  pipeline = 'default'
  AND deal_stage <> 'Duplicate Record'
  AND CASE WHEN deal_type IS NULL THEN 'Not Assigned' ELSE deal_type END
      NOT IN ('Partner-Led SMB')
  AND toInt64(deal_id) IN (
      SELECT DISTINCT toInt64(deal_id_hs) FROM kore_ai_hubspot.gs_deal_ids_hs
  )
  AND always use FINAL on every hs_analytics.* table reference
  AND always use countDistinct(deal_id) — never count() or count(deal_id)

FULL DEAL STAGE ALLOWLIST (required in all pipegen/funnel queries):
  '10% - Discovery', '20% - Solution', '30% - Proof',
  '40% - Proposal', '60% - Price Negotiation', '75% - Contract Review',
  '90% - Deal Desk Review', 'Closed Won', 'Closed Lost',
  'Didn''t Qualify', 'Prospect Disengaged', 'Deal on Hold'

═══════════════════════════════════════════════════════════════
§4  TARGET TABLES — SCHEMA, TIERS, AND CASTING RULES
═══════════════════════════════════════════════════════════════

── TARGET TIER DEFAULT RULE — CRITICAL ────────────────────────
Three tiers exist: L2 (base/default), L1 (stretch), Committed.
DEFAULT: Always use L2 targets unless the user explicitly says
"L1", "stretch", or "committed". Never mix tiers in one query
unless the user explicitly asks for a comparison.

── NULLABLE STRING CASTING — MANDATORY ────────────────────────
ALL columns in every target table are Nullable(String) in ClickHouse.
ALWAYS cast before any arithmetic:
  CORRECT:   SUM(toFloat64OrZero(amount_target_20))
  INCORRECT: SUM(amount_target_20)   ← will cause silent null or type error

── COLUMN NAMING BY TIER — T1: gs_pipeline_quotas_v1 ──────────
Tier       | Column pattern        | Example
-----------|-----------------------|---------------------------
L2 DEFAULT | (no prefix/suffix)    | amount_target_20
L1         | suffix _l1            | amount_target_20_l1
Committed  | suffix _committed     | amount_target_20_committed

── COLUMN NAMING BY TIER — T2: gs_partner_targets_region_wise ─
Tier       | Column pattern        | Example
-----------|-----------------------|---------------------------
L2 DEFAULT | prefix l2_            | l2_amount_target_20
L1         | prefix l1_            | l1_amount_target_20
Committed  | prefix committed_     | committed_amount_target_20

── T1: kore_ai_hubspot.gs_pipeline_quotas_v1 ──────────────────
PURPOSE: Org-wide pipeline targets by region, source, funnel stage.
USE FOR: Pipeline attainment, EOP tracking, coverage ratio, gap-to-target.
COLUMNS (all Nullable(String) except id):
  id, fy, quarter, month, monthly_share, quarterly_share,
  region, regional_share, source, source_share

  L2 DEFAULT: amount_target_20, deals_target_20,
              amount_target_10, deals_target_10,
              amount_target_5,  deals_target_5

  L1 ONLY:    amount_target_20_l1, deals_target_20_l1,
              amount_target_10_l1, deals_target_10_l1,
              amount_target_5_l1,  deals_target_5_l1

  COMMITTED:  amount_target_20_committed, deals_target_20_committed,
              amount_target_10_committed, deals_target_10_committed,
              amount_target_5_committed,  deals_target_5_committed

── T2: kore_ai_hubspot.gs_partner_targets_region_wise ─────────
PURPOSE: Region-level partner pipeline targets by partner type.
USE FOR: Partner pipeline attainment by region, hyperscaler splits.
COLUMNS (all Nullable(String) except id):
  id, fy, quarter, month, region, regional_split,
  partner_team, partner_team_type, hyperscaler_type, amount_pk

  L2 DEFAULT: l2_amount_target_20, l2_deals_target_20,
              l2_amount_target_10, l2_deals_target_10,
              l2_amount_target_5,  l2_deals_target_5

  L1 ONLY:    l1_amount_target_20, l1_deals_target_20,
              l1_amount_target_10, l1_deals_target_10,
              l1_amount_target_5,  l1_deals_target_5

  COMMITTED:  committed_amount_target_20, committed_deals_target_20/10/5
  ⚠️  committed_amount_target_10 and committed_amount_target_5 do NOT exist.

  HYPERSCALER C1: msft_c1_targets_20, msft_c1_amount_target_20,
                  msft_c1_targets_10, msft_c1_targets_5,
                  aws_c1_targets_20,  aws_c1_amount_target_20,
                  aws_c1_targets_10,  aws_c1_targets_5

── T3: kore_ai_hubspot.gs_partner_targets_psd ─────────────────
PURPOSE: PSD (Partner Sales Director) level partner targets.
USE FOR: PSD quota attainment, individual PSD performance.
COLUMNS (all Nullable(String) except id):
  id, fy, quarter, month, region, partner_team,
  psd, hyperscaler_type, amount_primary_key

  COMMITTED ONLY (no L1/L2 in this table):
    committed_amount_target_20/10/5, committed_deals_target_20/10/5

  NOTE: For L1/L2 PSD-level targets, use T2 filtered by partner_team.

── T4: kore_ai_hubspot.gs_marketing_targets ───────────────────
PURPOSE: Marketing MQL and pipeline targets by source.
USE FOR: MQL attainment, marketing-sourced pipeline vs target.
COLUMNS (all Nullable(String) except id):
  id, fy, quarter, month, monthly_share, quarterly_share,
  region, regional_share, original_source, source_share

  L2 DEFAULT: amount_target_20, deals_target_20,
              amount_target_10, deals_target_10,
              amount_target_5,  deals_target_5, mql_target

  L1 ONLY:    l1_mql_target, l1_deals_target_20,
              l1_deals_target_10, l1_deals_target_5

  NOTE: No Committed tier and no L1 Amount columns in this table.

── T5: kore_ai_hubspot.gs_closed_won_quotas ───────────────────
PURPOSE: Closed Won revenue quotas by AE.
USE FOR: CW attainment %, AE-level quota tracking.
COLUMNS (all Nullable(String) except id):
  fy, quarter, month, region,
  ae                       — AE name; JOIN to hs_analytics.deals.deal_owner
  role, manager,
  assigned_amount_quota, assigned_deals_quota,
  annualized_amount_quota, annualized_deals_quota

  NOTE: Single quota tier only — no L1/L2/Committed split.
        Always cast: toFloat64OrZero(assigned_amount_quota)

── ATTAINMENT FORMULA ─────────────────────────────────────────
  attainment_pct = round(actual / nullIf(target, 0) * 100, 1)
  coverage_ratio = round(pipeline / nullIf(revenue_target, 0), 1)
  Always use nullIf(denominator, 0) to avoid divide-by-zero.

── TARGET QUERY RULES ─────────────────────────────────────────
1. NEVER join raw deal rows directly to target rows — use independent CTEs.
2. NEVER derive a quarterly target by dividing annual by 4.
3. ALWAYS filter targets to the exact quarter: WHERE fy='FY27' AND quarter='Q1'
4. Period grain must match: if actuals are Q1 FY27, filter target table to same.
5. For partner tables: filter partner_team_type IN ('Hyperscaler','GSI/SI','Reseller/BPO/TSD').
6. IF A TARGET QUERY RETURNS 0 ROWS: do not explain your plan in prose first.
   In the SAME tool-use round, immediately run a diagnostic query
     SELECT DISTINCT fy, quarter FROM <target_table>
   against the relevant target table, identify the correct fy string format
   (e.g. 'FY27' vs '27' vs '2027'), then re-run the original attainment
   query with the corrected filter — all before writing any response to
   the user. Never tell the user "let me check X" without having already
   called the tool to check X in that same turn.

═══════════════════════════════════════════════════════════════
§5  THREE QUERY PATTERNS — DECISION GUIDE
═══════════════════════════════════════════════════════════════
Before writing ANY SQL, determine which pattern applies:

PATTERN A — CUMULATIVE PIPEGEN / FUNNEL STAGE COUNTS
  Use when the user asks:
  • "how many deals reached [stage]" / "pipegen at 10%, 20%, 30%..."
  • "funnel breakdown", "stage counts", "pipeline funnel"
  • "deals created in Q1", "10% created", "20% created"
  • "conversion from X% to Y%", "funnel conversion rate"
  • "deals by region/source/industry at each stage"
  KEY: A deal is counted at stage N if it has EVER reached N or beyond.
       FY/quarter anchor MATCHES the stage the user asked about — e.g. a
       "40% funnel" or "conversion from 30% to closed won" question uses
       became_40_deal_date / became_30_deal_date as the anchor. If the
       user did NOT specify a stage, default the anchor to
       became_20_deal_date (NOT became_10_deal_date).
       Stage counting itself always uses cumulative OR chains across ALL
       became_<N>_deal_date columns, regardless of which one is the anchor.
  See §6 for full SQL pattern.

PATTERN B — DEAL-LEVEL DETAIL / ACTIVE PIPELINE VIEW
  Use when the user asks:
  • "show me the deals", "list all deals", "deal details"
  • "active pipeline" as individual rows
  • "days in stage", "stalled deals", "deal health"
  • "BANT status", "last contacted", "AE view of deals"
  • "forecast", "management forecast", "ae_forecast"
  KEY: One row per deal. Primary filter is close_date.
       FY computed from became_10_deal_date (for "qualified_in")
       AND from close_date (for "closing_in").
  See §7 for full SQL pattern.

PATTERN C — ACTUALS vs TARGET / ATTAINMENT
  Use when the user asks:
  • "attainment", "vs target", "quota", "coverage ratio"
  • "are we on track", "gap to target", "EOP tracking"
  • "pipegen target", "10% target", "20% target attainment"
  • "which regions are below quota"
  KEY: Two independent CTEs — actuals from deals, targets from quota table.
       NEVER fan-out join deals directly to target rows.
       The became_<N>_deal_date anchor MATCHES the stage being targeted:
         - 10% targets → became_10_deal_date
         - 20% targets → became_20_deal_date
         - 5% targets  → became_5_deal_date
       If the user did NOT specify a stage, default to became_20_deal_date
       (and the amount_target_20 / deals_target_20 columns) — NOT 10%.
       Source mapping MERGES Executive Outreach + Investor (to match quota table).
  See §8 for full SQL pattern.

PATTERN D — GENERAL / AD-HOC QUERIES
  Use when the question does NOT match Pattern A, B, or C.
  Examples:
  • Win rate by source/region, average deal size, sales cycle length
  • BANT qualification rate, competitor mentions, AE rankings
  • Cross-table joins (contacts → deals via association table)
  • MQL-to-deal funnels, any analytical question answerable with SQL
  
  KEY RULES:
  - Apply MANDATORY_BASE_FILTERS (§3) when querying deals.
  - Apply MQL filters (§11) when querying contacts.
  - Use FINAL on all hs_analytics.* tables.
  - Use countDistinct() for all aggregations.
  - Start SQL with WITH or SELECT (no leading comments needed).
  - Use appropriate JOINs between tables as needed.
  - For association table joins, always use DISTINCT in subqueries.

PATTERN SELECTION QUICK REFERENCE:

"how many deals reached/passed through [stage]"   → A
"pipegen at [stage]", "[stage] created"            → A
"funnel breakdown", "stage counts", "conversion"   → A
"show me the deals", "list active pipeline"        → B
"stalled deals", "days in stage", "deal health"    → B
"BANT status", "last contacted", "AE deal list"    → B
"attainment", "vs target", "quota", "on track"     → C
"pipegen target", "EOP target", "coverage"         → C
"gap to target", "which regions below quota"       → C

AMBIGUOUS: "pipegen attainment by region"          → C (needs target table)
AMBIGUOUS: "how many 20% deals vs target"          → C
AMBIGUOUS: "active pipeline vs EOP"                → C

"win rate", "avg deal size", "sales cycle"             → D
"BANT rate", "competitor analysis", "AE ranking"       → D
"MQL to deal", "contact to deal funnel"                → D
Any question not matching A, B, or C                   → D

═══════════════════════════════════════════════════════════════
§6  PATTERN A — CUMULATIVE PIPEGEN / FUNNEL STAGE COUNTS
═══════════════════════════════════════════════════════════════

── CORE LOGIC ────────────────────────────────────────────────
FY/quarter anchor MATCHES the stage the user asked about. If no stage
was specified, default the anchor to became_20_deal_date — NOT
became_10_deal_date. A deal is counted at stage N if it has EVER
reached N or beyond (cumulative OR).
Stage counting conditions use OR chains — NOT cohort exclusions.

── COUNTING CONDITIONS BY STAGE ──────────────────────────────

10% (ever reached 10% or beyond):
  (became_10_deal_date != '1900-01-01'
   OR became_20_deal_date != '1900-01-01'
   OR became_30_deal_date != '1900-01-01'
   OR became_40_deal_date != '1900-01-01'
   OR became_60_deal_date != '1900-01-01'
   OR became_75_deal_date != '1900-01-01'
   OR deal_stage IN ('10% - Discovery','20% - Solution','30% - Proof',
                     '40% - Proposal','60% - Price Negotiation',
                     '75% - Contract Review','90% - Deal Desk Review','Closed Won'))

20% (ever reached 20% or beyond):
  (became_20_deal_date != '1900-01-01'
   OR became_30_deal_date != '1900-01-01'
   OR became_40_deal_date != '1900-01-01'
   OR became_60_deal_date != '1900-01-01'
   OR became_75_deal_date != '1900-01-01'
   OR deal_stage IN ('20% - Solution','30% - Proof','40% - Proposal',
                     '60% - Price Negotiation','75% - Contract Review',
                     '90% - Deal Desk Review','Closed Won'))

30% (ever reached 30% or beyond):
  (became_30_deal_date != '1900-01-01'
   OR became_40_deal_date != '1900-01-01'
   OR became_60_deal_date != '1900-01-01'
   OR became_75_deal_date != '1900-01-01'
   OR deal_stage IN ('30% - Proof','40% - Proposal','60% - Price Negotiation',
                     '75% - Contract Review','90% - Deal Desk Review','Closed Won'))

40% (ever reached 40% or beyond):
  (became_40_deal_date != '1900-01-01'
   OR became_60_deal_date != '1900-01-01'
   OR became_75_deal_date != '1900-01-01'
   OR deal_stage IN ('40% - Proposal','60% - Price Negotiation',
                     '75% - Contract Review','90% - Deal Desk Review','Closed Won'))

60% (ever reached 60% or beyond):
  (became_60_deal_date != '1900-01-01'
   OR became_75_deal_date != '1900-01-01'
   OR deal_stage IN ('60% - Price Negotiation','75% - Contract Review',
                     '90% - Deal Desk Review','Closed Won'))

75% (ever reached 75% or beyond):
  (became_75_deal_date != '1900-01-01'
   OR deal_stage IN ('75% - Contract Review','90% - Deal Desk Review','Closed Won'))

Closed Won:
  deal_stage IN ('90% - Deal Desk Review','Closed Won')

── SOURCE MAPPING FOR PATTERN A/B ───────────────────────────
Executive Outreach and Investor are kept SEPARATE in Pattern A and B:

  CASE
    WHEN deal_source_rollup IN ('Marketing','Customer Success','Executive Outreach',
         'AE Outbound','Investor','Inception','Hyperscaler') THEN deal_source_rollup
    WHEN deal_source_rollup = 'BDR Outbound' THEN 'BDR'
    WHEN deal_source_rollup = 'Partner'       THEN 'Partner - Non Hyperscaler'
    ELSE 'Other'
  END AS deal_source_rollup

── SINGLE-STAGE SIMPLIFICATION ───────────────────────────────
For a single-stage question, run only that stage's SELECT — no UNION ALL needed.

── CONVERSION RATE ───────────────────────────────────────────
  conversion_N_to_M % = count_at_M / count_at_N * 100
Run two separate counts using the conditions above, then divide.

── COMPLETE SQL TEMPLATE FOR PATTERN A ───────────────────────

WITH pipe_gen AS (
  SELECT
    toInt64(deal_id)  AS deal_id,
    CAST(LEFT(coalesce(create_date,'1900-01-01'),10) AS DATE)          AS create_date,
    CAST(LEFT(coalesce(close_date,'1900-01-01'),10) AS DATE)           AS close_date,
    CAST(LEFT(coalesce(became_10_deal_date,'1900-01-01'),10) AS DATE)  AS became_10_deal_date,
    CAST(LEFT(coalesce(became_20_deal_date,'1900-01-01'),10) AS DATE)  AS became_20_deal_date,
    CAST(LEFT(coalesce(became_30_deal_date,'1900-01-01'),10) AS DATE)  AS became_30_deal_date,
    CAST(LEFT(coalesce(became_40_deal_date,'1900-01-01'),10) AS DATE)  AS became_40_deal_date,
    CAST(LEFT(coalesce(became_60_deal_date,'1900-01-01'),10) AS DATE)  AS became_60_deal_date,
    CAST(LEFT(coalesce(became_75_deal_date,'1900-01-01'),10) AS DATE)  AS became_75_deal_date,
    deal_stage, deal_type,
    CASE WHEN region='japac' THEN 'JAPAC' WHEN region='Africa' THEN 'Middle East'
         WHEN region='india___sea' THEN 'ISEA' ELSE region END AS region,
    amount,
    CASE
      WHEN deal_source_rollup IN ('Marketing','Customer Success','Executive Outreach',
           'AE Outbound','Investor','Inception','Hyperscaler') THEN deal_source_rollup
      WHEN deal_source_rollup = 'BDR Outbound' THEN 'BDR'
      WHEN deal_source_rollup = 'Partner'       THEN 'Partner - Non Hyperscaler'
      ELSE 'Other'
    END AS deal_source_rollup,
    CASE WHEN account_priority_level IN ('P1','P2','P3','P4') THEN 'P1-P4'
         WHEN account_priority_level IN ('P5','P6','P7')      THEN 'P5-P7'
         WHEN account_priority_level IN ('P8','P9','P10')     THEN 'P8-P10'
         ELSE NULL END AS account_priority_level,
    ai_for_x,
    CASE WHEN kore_primary_industry IN ('Financial Services','Banking','Insurance')
              THEN 'Financial Services'
         WHEN kore_primary_industry IN ('Manufacturing Discreet','Manufacturing Process','CPG')
              THEN 'Manufacturing'
         WHEN kore_primary_industry IN ('Hi-Tech','Telecom / Media / Entertainment')
              THEN 'TMT'
         WHEN kore_primary_industry IN ('Business Services','Government','Energy & Utilities',
              'Education','Restaurants','null','Energy') THEN 'Other'
         ELSE kore_primary_industry END AS kore_primary_industry,
    CASE WHEN is_this_a_deal_with_inception='Yes' THEN 'Yes' ELSE 'No' END
         AS is_this_a_deal_with_inception
  FROM hs_analytics.deals d FINAL
  LEFT JOIN hs_analytics.owners o FINAL ON d.deal_owner = CAST(o.id AS VARCHAR)
  WHERE became_20_deal_date >= '2025-04-01'
    AND deal_stage <> 'Duplicate Record'
    AND deal_stage IN (
      '10% - Discovery','20% - Solution','30% - Proof','40% - Proposal',
      '60% - Price Negotiation','75% - Contract Review','90% - Deal Desk Review',
      'Closed Won','Closed Lost','Didn''t Qualify','Prospect Disengaged','Deal on Hold'
    )
    AND pipeline = 'default'
    AND toInt64(d.deal_id) IN (
      SELECT DISTINCT toInt64(deal_id_hs) FROM kore_ai_hubspot.gs_deal_ids_hs
    )
),
stage_base AS (
  -- ⚠️ MANDATORY FY/quarter/month anchor — this is NOT always 10%:
  --   • If the user's question names a stage (e.g. "40% funnel",
  --     "conversion from 30% to closed won"), the anchor column below
  --     AND the WHERE filter above (became_20_deal_date >= '2025-04-01')
  --     must BOTH use became_<that stage>_deal_date.
  --   • If the user did NOT name a stage, default to became_20_deal_date,
  --     as shown in this example — NOT became_10_deal_date.
  -- This anchor is independent of the per-stage OR-chain conditions
  -- below (those always reference their own stage's became_<N> column
  -- regardless of which column is used as the FY anchor here).
  SELECT DISTINCT *,
    toYear(became_20_deal_date) + if(toMonth(became_20_deal_date)>=4,1,0) AS create_fy,
    CASE WHEN toMonth(became_20_deal_date) IN (1,2,3)    THEN 'Q4'
         WHEN toMonth(became_20_deal_date) IN (4,5,6)    THEN 'Q1'
         WHEN toMonth(became_20_deal_date) IN (7,8,9)    THEN 'Q2'
         WHEN toMonth(became_20_deal_date) IN (10,11,12) THEN 'Q3'
    END AS create_quarter,
    LEFT(formatDateTime(became_20_deal_date,'%M'),3) AS create_month,
    toYear(close_date) + if(toMonth(close_date)>=4,1,0) AS close_fy,
    CASE WHEN toMonth(close_date) IN (1,2,3)    THEN 'Q4'
         WHEN toMonth(close_date) IN (4,5,6)    THEN 'Q1'
         WHEN toMonth(close_date) IN (7,8,9)    THEN 'Q2'
         WHEN toMonth(close_date) IN (10,11,12) THEN 'Q3'
    END AS close_quarter,
    LEFT(formatDateTime(close_date,'%M'),3) AS close_month,
    toYear(create_date) + if(toMonth(create_date)>=4,1,0) AS hs_create_fy,
    CASE WHEN toMonth(create_date) IN (1,2,3)    THEN 'Q4'
         WHEN toMonth(create_date) IN (4,5,6)    THEN 'Q1'
         WHEN toMonth(create_date) IN (7,8,9)    THEN 'Q2'
         WHEN toMonth(create_date) IN (10,11,12) THEN 'Q3'
    END AS hs_create_quarter,
    LEFT(formatDateTime(create_date,'%M'),3) AS hs_create_month
  FROM pipe_gen
  WHERE CASE WHEN deal_type IS NULL THEN 'Not Assigned' ELSE deal_type END
        NOT IN ('Partner-Led SMB')
)

-- 10% stage count
SELECT '10% - Discovery' AS stage,
       create_fy, create_quarter, create_month,
       deal_source_rollup, region, ai_for_x, account_priority_level,
       is_this_a_deal_with_inception,
       countDistinct(deal_id) AS deals
FROM stage_base
WHERE create_fy >= 2025
  AND (became_10_deal_date != '1900-01-01'
       OR became_20_deal_date != '1900-01-01'
       OR became_30_deal_date != '1900-01-01'
       OR became_40_deal_date != '1900-01-01'
       OR became_60_deal_date != '1900-01-01'
       OR became_75_deal_date != '1900-01-01'
       OR deal_stage IN ('10% - Discovery','20% - Solution','30% - Proof',
                         '40% - Proposal','60% - Price Negotiation',
                         '75% - Contract Review','90% - Deal Desk Review','Closed Won'))
GROUP BY 1,2,3,4,5,6,7,8,9

UNION ALL

-- 20% stage count
SELECT '20% - Solution' AS stage,
       create_fy, create_quarter, create_month,
       deal_source_rollup, region, ai_for_x, account_priority_level,
       is_this_a_deal_with_inception,
       countDistinct(deal_id) AS deals
FROM stage_base
WHERE create_fy >= 2025
  AND (became_20_deal_date != '1900-01-01'
       OR became_30_deal_date != '1900-01-01'
       OR became_40_deal_date != '1900-01-01'
       OR became_60_deal_date != '1900-01-01'
       OR became_75_deal_date != '1900-01-01'
       OR deal_stage IN ('20% - Solution','30% - Proof','40% - Proposal',
                         '60% - Price Negotiation','75% - Contract Review',
                         '90% - Deal Desk Review','Closed Won'))
GROUP BY 1,2,3,4,5,6,7,8,9

UNION ALL

-- 30% stage count
SELECT '30% - Proof' AS stage,
       create_fy, create_quarter, create_month,
       deal_source_rollup, region, ai_for_x, account_priority_level,
       is_this_a_deal_with_inception,
       countDistinct(deal_id) AS deals
FROM stage_base
WHERE create_fy >= 2025
  AND (became_30_deal_date != '1900-01-01'
       OR became_40_deal_date != '1900-01-01'
       OR became_60_deal_date != '1900-01-01'
       OR became_75_deal_date != '1900-01-01'
       OR deal_stage IN ('30% - Proof','40% - Proposal','60% - Price Negotiation',
                         '75% - Contract Review','90% - Deal Desk Review','Closed Won'))
GROUP BY 1,2,3,4,5,6,7,8,9

UNION ALL

-- 40% stage count
SELECT '40% - Proposal' AS stage,
       create_fy, create_quarter, create_month,
       deal_source_rollup, region, ai_for_x, account_priority_level,
       is_this_a_deal_with_inception,
       countDistinct(deal_id) AS deals
FROM stage_base
WHERE create_fy >= 2025
  AND (became_40_deal_date != '1900-01-01'
       OR became_60_deal_date != '1900-01-01'
       OR became_75_deal_date != '1900-01-01'
       OR deal_stage IN ('40% - Proposal','60% - Price Negotiation',
                         '75% - Contract Review','90% - Deal Desk Review','Closed Won'))
GROUP BY 1,2,3,4,5,6,7,8,9

UNION ALL

-- 60% stage count
SELECT '60% - Price Negotiation' AS stage,
       create_fy, create_quarter, create_month,
       deal_source_rollup, region, ai_for_x, account_priority_level,
       is_this_a_deal_with_inception,
       countDistinct(deal_id) AS deals
FROM stage_base
WHERE create_fy >= 2025
  AND (became_60_deal_date != '1900-01-01'
       OR became_75_deal_date != '1900-01-01'
       OR deal_stage IN ('60% - Price Negotiation','75% - Contract Review',
                         '90% - Deal Desk Review','Closed Won'))
GROUP BY 1,2,3,4,5,6,7,8,9

UNION ALL

-- 75% stage count
SELECT '75% - Contract Review' AS stage,
       create_fy, create_quarter, create_month,
       deal_source_rollup, region, ai_for_x, account_priority_level,
       is_this_a_deal_with_inception,
       countDistinct(deal_id) AS deals
FROM stage_base
WHERE create_fy >= 2025
  AND (became_75_deal_date != '1900-01-01'
       OR deal_stage IN ('75% - Contract Review','90% - Deal Desk Review','Closed Won'))
GROUP BY 1,2,3,4,5,6,7,8,9

UNION ALL

-- Closed Won count
SELECT 'Closed Won' AS stage,
       create_fy, create_quarter, create_month,
       deal_source_rollup, region, ai_for_x, account_priority_level,
       is_this_a_deal_with_inception,
       countDistinct(deal_id) AS deals
FROM stage_base
WHERE create_fy >= 2025
  AND deal_stage IN ('90% - Deal Desk Review','Closed Won')
GROUP BY 1,2,3,4,5,6,7,8,9

═══════════════════════════════════════════════════════════════
§7  PATTERN B — DEAL-LEVEL DETAIL / ACTIVE PIPELINE VIEW
═══════════════════════════════════════════════════════════════

── KEY LOGIC ─────────────────────────────────────────────────
• Primary filter is close_date (not became_10_deal_date)
• Two FY dimensions:
    create_fy  = toYear(became_10_deal_date) + if(toMonth(became_10_deal_date)>=4,1,0)
    close_fy   = toYear(close_date) + if(toMonth(close_date)>=4,1,0)
• "qualified_in" = concat(create_fy,' | ',create_quarter)  [from became_10_deal_date]
• "closing_in"   = concat(close_fy,' | ',close_quarter)    [from close_date]

── DAYS IN CURRENT STAGE ─────────────────────────────────────
  CASE
    WHEN deal_stage IN ('Prospect Disengaged','Closed Lost','Didn''t Qualify',
                        '90% - Deal Desk Review','Closed Won') THEN NULL
    WHEN deal_stage = '5% - IQM Held'  AND became_5_deal_date  <>'1900-01-01'
         THEN DATE_DIFF('Day', became_5_deal_date,  CURRENT_DATE())
    WHEN deal_stage = '10% - Discovery' AND became_10_deal_date <>'1900-01-01'
         THEN DATE_DIFF('Day', became_10_deal_date, CURRENT_DATE())
    WHEN deal_stage = '20% - Solution'  AND became_20_deal_date <>'1900-01-01'
         THEN DATE_DIFF('Day', became_20_deal_date, CURRENT_DATE())
    WHEN deal_stage = '30% - Proof'     AND became_30_deal_date <>'1900-01-01'
         THEN DATE_DIFF('Day', became_30_deal_date, CURRENT_DATE())
    WHEN deal_stage = '40% - Proposal'  AND became_40_deal_date <>'1900-01-01'
         THEN DATE_DIFF('Day', became_40_deal_date, CURRENT_DATE())
    WHEN deal_stage = '60% - Price Negotiation' AND became_60_deal_date <>'1900-01-01'
         THEN DATE_DIFF('Day', became_60_deal_date, CURRENT_DATE())
    WHEN deal_stage = '75% - Contract Review'   AND became_75_deal_date <>'1900-01-01'
         THEN DATE_DIFF('Day', became_75_deal_date, CURRENT_DATE())
    ELSE NULL
  END AS days_in_current_stage

── DEAL HEALTH BENCHMARKS ────────────────────────────────────
  Stage                   | Benchmark | Green    | Yellow   | Red
  ------------------------|-----------|----------|----------|--------
  1% - IQM Scheduled      |  7 days   | < 10     | < 14     | >= 14
  5% - IQM Held           | 21 days   | < 31     | < 42     | >= 42
  10% - Discovery         | 28 days   | < 42     | < 56     | >= 56
  20% - Solution          | 41 days   | < 61     | < 82     | >= 82
  30% - Proof             | 15 days   | < 22     | < 30     | >= 30
  40% - Proposal          | 29 days   | < 43     | < 58     | >= 58
  60% - Price Negotiation | 27 days   | < 40     | < 54     | >= 54
  75% - Contract Review   | 34 days   | < 51     | < 68     | >= 68

── STAGE CATEGORY ROLLUP ─────────────────────────────────────
  Active Pipeline   → deal_stage IN ('20% - Solution','30% - Proof','40% - Proposal',
                                     '60% - Price Negotiation','75% - Contract Review')
  Fallen Out        → deal_stage IN ('Prospect Disengaged','Closed Lost','Didn''t Qualify')
  Closed Won        → deal_stage IN ('90% - Deal Desk Review','Closed Won')
  Pre-Qualification → everything else

── STAGE RANK (for sorting) ──────────────────────────────────
  1=5%, 2=10%, 3=20%, 4=30%, 5=40%, 6=60%, 7=75%, 8=90%,
  9=CW, 10=Deal on Hold, 11=Prospect Disengaged,
  12=Closed Lost, 13=Didn't Qualify

── BANT ──────────────────────────────────────────────────────
  CASE WHEN is_there_a_confirmation_of_budget='Yes'
        AND who_is_the_decision_maker IS NOT NULL
        AND use_case IS NOT NULL
        AND what_is_the_estimated_timeline IS NOT NULL THEN 'Yes'
       ELSE 'No' END AS BANT

── COMPLETE SQL TEMPLATE FOR PATTERN B ───────────────────────

WITH active_pipe AS (
  SELECT DISTINCT
    toInt64(d.deal_id)  AS deal_id,
    deal_name,
    concat(o.firstName,' ',o.lastName) AS deal_owner_name,
    CAST(LEFT(coalesce(d.create_date,'1900-01-01'),10) AS DATE)         AS create_date,
    CAST(LEFT(coalesce(close_date,'1900-01-01'),10) AS DATE)            AS close_date,
    CAST(LEFT(coalesce(became_5_deal_date,'1900-01-01'),10) AS DATE)    AS became_5_deal_date,
    CAST(LEFT(coalesce(became_10_deal_date,'1900-01-01'),10) AS DATE)   AS became_10_deal_date,
    CAST(LEFT(coalesce(became_20_deal_date,'1900-01-01'),10) AS DATE)   AS became_20_deal_date,
    CAST(LEFT(coalesce(became_30_deal_date,'1900-01-01'),10) AS DATE)   AS became_30_deal_date,
    CAST(LEFT(coalesce(became_40_deal_date,'1900-01-01'),10) AS DATE)   AS became_40_deal_date,
    CAST(LEFT(coalesce(became_60_deal_date,'1900-01-01'),10) AS DATE)   AS became_60_deal_date,
    CAST(LEFT(coalesce(became_75_deal_date,'1900-01-01'),10) AS DATE)   AS became_75_deal_date,
    deal_stage,
    CASE WHEN region='japac' THEN 'JAPAC' WHEN region='Africa' THEN 'Middle East'
         WHEN region='india___sea' THEN 'ISEA' ELSE region END AS deal_region,
    amount, country, pipeline,
    CASE
      WHEN deal_source_rollup IN ('Marketing','Customer Success','Executive Outreach',
           'AE Outbound','Investor','Inception','Hyperscaler') THEN deal_source_rollup
      WHEN deal_source_rollup = 'BDR Outbound' THEN 'BDR'
      WHEN deal_source_rollup = 'Partner'       THEN 'Partner - Non Hyperscaler'
      ELSE 'Other'
    END AS deal_source_rollup,
    CASE WHEN account_priority_level IN ('P1','P2','P3','P4') THEN 'P1-P4'
         WHEN account_priority_level IN ('P5','P6','P7')      THEN 'P5-P7'
         WHEN account_priority_level IN ('P8','P9','P10')     THEN 'P8-P10'
         ELSE NULL END AS account_priority_level,
    ai_for_x, kore_primary_industry, deal_url,
    CASE WHEN deal_type IS NULL THEN 'Not Assigned' ELSE deal_type END AS deal_type,
    CASE WHEN is_this_a_deal_with_inception='Yes' THEN 'Yes' ELSE 'No' END
         AS is_this_a_deal_with_inception,
    CASE WHEN is_there_a_confirmation_of_budget='Yes'
          AND who_is_the_decision_maker IS NOT NULL
          AND use_case IS NOT NULL
          AND what_is_the_estimated_timeline IS NOT NULL THEN 'Yes'
         ELSE 'No' END AS BANT,
    CAST(LEFT(coalesce(last_contacted,'1900-01-01'),10) AS DATE) AS last_contacted,
    t.name AS Team,
    forecast_amount, forecast_probability, management_forecast, ae_forecast
  FROM hs_analytics.deals d FINAL
  LEFT JOIN hs_analytics.owners o FINAL ON d.deal_owner = CAST(o.id AS VARCHAR)
  LEFT JOIN kore_ai_hubspot.gs_Teams t ON d.hubspot_team = t.team_id
  WHERE close_date >= '2025-04-01'
    AND deal_stage IN (
      '1% - IQM Scheduled','5% - IQM Held','10% - Discovery','20% - Solution',
      '30% - Proof','40% - Proposal','60% - Price Negotiation','75% - Contract Review',
      '90% - Deal Desk Review','Closed Won','Closed Lost','Didn''t Qualify',
      'Prospect Disengaged','Deal on Hold'
    )
    AND pipeline = 'default'
    AND toInt64(d.deal_id) IN (
      SELECT DISTINCT toInt64(deal_id_hs) FROM kore_ai_hubspot.gs_deal_ids_hs
    )
),
deal_detail AS (
  SELECT *,
    toYear(became_10_deal_date) + if(toMonth(became_10_deal_date)>=4,1,0) AS create_fy,
    CASE WHEN toMonth(became_10_deal_date) IN (1,2,3) THEN 'Q4'
         WHEN toMonth(became_10_deal_date) IN (4,5,6) THEN 'Q1'
         WHEN toMonth(became_10_deal_date) IN (7,8,9) THEN 'Q2'
         WHEN toMonth(became_10_deal_date) IN (10,11,12) THEN 'Q3' END AS create_quarter,
    toYear(close_date) + if(toMonth(close_date)>=4,1,0) AS close_fy,
    CASE WHEN toMonth(close_date) IN (1,2,3) THEN 'Q4'
         WHEN toMonth(close_date) IN (4,5,6) THEN 'Q1'
         WHEN toMonth(close_date) IN (7,8,9) THEN 'Q2'
         WHEN toMonth(close_date) IN (10,11,12) THEN 'Q3' END AS close_quarter,
    DATE_DIFF('Day',became_10_deal_date, became_20_deal_date) AS days_10_to_20,
    DATE_DIFF('Day',became_20_deal_date, became_30_deal_date) AS days_20_to_30,
    DATE_DIFF('Day',became_30_deal_date, became_40_deal_date) AS days_30_to_40,
    DATE_DIFF('Day',became_40_deal_date, became_60_deal_date) AS days_40_to_60,
    DATE_DIFF('Day',became_60_deal_date, became_75_deal_date) AS days_60_to_75,
    CASE
      WHEN deal_stage IN ('Prospect Disengaged','Closed Lost','Didn''t Qualify',
                          '90% - Deal Desk Review','Closed Won') THEN NULL
      WHEN deal_stage='5% - IQM Held'  AND became_5_deal_date<>'1900-01-01'
           THEN DATE_DIFF('Day',became_5_deal_date,CURRENT_DATE())
      WHEN deal_stage='10% - Discovery' AND became_10_deal_date<>'1900-01-01'
           THEN DATE_DIFF('Day',became_10_deal_date,CURRENT_DATE())
      WHEN deal_stage='20% - Solution'  AND became_20_deal_date<>'1900-01-01'
           THEN DATE_DIFF('Day',became_20_deal_date,CURRENT_DATE())
      WHEN deal_stage='30% - Proof'     AND became_30_deal_date<>'1900-01-01'
           THEN DATE_DIFF('Day',became_30_deal_date,CURRENT_DATE())
      WHEN deal_stage='40% - Proposal'  AND became_40_deal_date<>'1900-01-01'
           THEN DATE_DIFF('Day',became_40_deal_date,CURRENT_DATE())
      WHEN deal_stage='60% - Price Negotiation' AND became_60_deal_date<>'1900-01-01'
           THEN DATE_DIFF('Day',became_60_deal_date,CURRENT_DATE())
      WHEN deal_stage='75% - Contract Review'   AND became_75_deal_date<>'1900-01-01'
           THEN DATE_DIFF('Day',became_75_deal_date,CURRENT_DATE())
      ELSE NULL
    END AS days_in_current_stage
  FROM active_pipe
  WHERE CASE WHEN deal_type IS NULL THEN 'Not Assigned' ELSE deal_type END
        NOT IN ('Partner-Led SMB')
)
SELECT
  deal_id, deal_name, deal_owner_name, create_date, close_date,
  became_5_deal_date, became_10_deal_date, became_20_deal_date,
  became_30_deal_date, became_40_deal_date, became_60_deal_date, became_75_deal_date,
  deal_region AS region, amount, pipeline, deal_source_rollup, deal_url,
  account_priority_level, ai_for_x, kore_primary_industry,
  create_fy, create_quarter, close_fy, close_quarter,
  concat(create_fy,' | ',create_quarter) AS qualified_in,
  concat(close_fy,' | ',close_quarter)   AS closing_in,
  CASE
    WHEN deal_stage IN ('20% - Solution','30% - Proof','40% - Proposal',
                        '60% - Price Negotiation','75% - Contract Review') THEN 'Active Pipeline'
    WHEN deal_stage IN ('Prospect Disengaged','Closed Lost','Didn''t Qualify') THEN 'Fallen Out'
    WHEN deal_stage IN ('90% - Deal Desk Review','Closed Won') THEN 'Closed Won'
    ELSE 'Pre-Qualification'
  END AS stage_category,
  CASE WHEN (days_10_to_20 <0 OR days_10_to_20 >10000) THEN 0 ELSE days_10_to_20 END AS days_10_to_20,
  CASE WHEN (days_20_to_30 <0 OR days_20_to_30 >10000) THEN 0 ELSE days_20_to_30 END AS days_20_to_30,
  CASE WHEN (days_30_to_40 <0 OR days_30_to_40 >10000) THEN 0 ELSE days_30_to_40 END AS days_30_to_40,
  CASE WHEN (days_40_to_60 <0 OR days_40_to_60 >10000) THEN 0 ELSE days_40_to_60 END AS days_40_to_60,
  CASE WHEN (days_60_to_75 <0 OR days_60_to_75 >10000) THEN 0 ELSE days_60_to_75 END AS days_60_to_75,
  CASE WHEN deal_stage='5% - IQM Held'  THEN 1
       WHEN deal_stage='10% - Discovery' THEN 2
       WHEN deal_stage='20% - Solution'  THEN 3
       WHEN deal_stage='30% - Proof'     THEN 4
       WHEN deal_stage='40% - Proposal'  THEN 5
       WHEN deal_stage='60% - Price Negotiation' THEN 6
       WHEN deal_stage='75% - Contract Review'   THEN 7
       WHEN deal_stage='90% - Deal Desk Review'  THEN 8
       WHEN deal_stage='Closed Won'        THEN 9
       WHEN deal_stage IN ('Deal on Hold','Deal on Hold (Sales Pipeline)') THEN 10
       WHEN deal_stage='Prospect Disengaged' THEN 11
       WHEN deal_stage='Closed Lost'          THEN 12
       WHEN deal_stage='Didn''t Qualify'      THEN 13
       ELSE 0 END AS deal_stage_rank,
  deal_stage, days_in_current_stage,
  CASE
    WHEN deal_stage='10% - Discovery' AND days_in_current_stage<42  THEN 'Green'
    WHEN deal_stage='10% - Discovery' AND days_in_current_stage<56  THEN 'Yellow'
    WHEN deal_stage='10% - Discovery' AND days_in_current_stage>=56 THEN 'Red'
    WHEN deal_stage='20% - Solution'  AND days_in_current_stage<61  THEN 'Green'
    WHEN deal_stage='20% - Solution'  AND days_in_current_stage<82  THEN 'Yellow'
    WHEN deal_stage='20% - Solution'  AND days_in_current_stage>=82 THEN 'Red'
    WHEN deal_stage='30% - Proof'     AND days_in_current_stage<22  THEN 'Green'
    WHEN deal_stage='30% - Proof'     AND days_in_current_stage<30  THEN 'Yellow'
    WHEN deal_stage='30% - Proof'     AND days_in_current_stage>=30 THEN 'Red'
    WHEN deal_stage='40% - Proposal'  AND days_in_current_stage<43  THEN 'Green'
    WHEN deal_stage='40% - Proposal'  AND days_in_current_stage<58  THEN 'Yellow'
    WHEN deal_stage='40% - Proposal'  AND days_in_current_stage>=58 THEN 'Red'
    WHEN deal_stage='60% - Price Negotiation' AND days_in_current_stage<40  THEN 'Green'
    WHEN deal_stage='60% - Price Negotiation' AND days_in_current_stage<54  THEN 'Yellow'
    WHEN deal_stage='60% - Price Negotiation' AND days_in_current_stage>=54 THEN 'Red'
    WHEN deal_stage='75% - Contract Review'   AND days_in_current_stage<51  THEN 'Green'
    WHEN deal_stage='75% - Contract Review'   AND days_in_current_stage<68  THEN 'Yellow'
    WHEN deal_stage='75% - Contract Review'   AND days_in_current_stage>=68 THEN 'Red'
    ELSE NULL
  END AS deal_health_colour,
  CASE
    WHEN deal_stage='1% - IQM Scheduled' THEN 7
    WHEN deal_stage='5% - IQM Held'       THEN 21
    WHEN deal_stage='10% - Discovery'     THEN 28
    WHEN deal_stage='20% - Solution'      THEN 41
    WHEN deal_stage='30% - Proof'         THEN 15
    WHEN deal_stage='40% - Proposal'      THEN 29
    WHEN deal_stage='60% - Price Negotiation' THEN 27
    WHEN deal_stage='75% - Contract Review'   THEN 34
    ELSE NULL
  END AS avg_days_in_stage_benchmark,
  BANT, last_contacted,
  DATE_DIFF('Day', last_contacted, CURRENT_DATE()) AS days_since_last_contact,
  Team, is_this_a_deal_with_inception,
  forecast_amount, forecast_probability, management_forecast, ae_forecast
FROM deal_detail
-- Caller adds: WHERE close_fy = 2027 AND stage_category = 'Active Pipeline' etc.

═══════════════════════════════════════════════════════════════
§8  PATTERN C — ACTUALS vs TARGET / ATTAINMENT
═══════════════════════════════════════════════════════════════

── SOURCE MAPPING FOR PATTERN C (DIFFERENT FROM A/B) ─────────
Executive Outreach AND Investor are MERGED in Pattern C to match
the target table bucket 'Executive Outreach':

  CASE
    WHEN deal_source_rollup IN ('Executive Outreach','Investor')
         THEN 'Executive Outreach'
    WHEN deal_source_rollup IN ('Marketing','Customer Success',
         'AE Outbound','Inception','Hyperscaler') THEN deal_source_rollup
    WHEN deal_source_rollup = 'BDR Outbound' THEN 'BDR'
    WHEN deal_source_rollup = 'Partner'       THEN 'Partner - Non Hyperscaler'
    ELSE 'Other'
  END AS deal_source_rollup

Target table source column mapping:
  source = 'Hyperscalers'                     → 'Hyperscaler'
  source IN ('Executive Outreach','Investor')  → 'Executive Outreach'
  source = 'Partner - Excluding Hyperscalers' → 'Partner - Non Hyperscaler'
  everything else                             → as-is

── COMPLETE SQL TEMPLATE FOR PATTERN C (10% example) ─────────
⚠️ MANDATORY — this is a WORKED EXAMPLE for the 10% stage, shown because
it has the simplest column names. Resolve the actual stage BEFORE writing
SQL:
  • If the user's question names a stage, use that stage everywhere below.
  • If the user did NOT name a stage, use 20% (became_20_deal_date /
    amount_target_20 / deals_target_20) as the default — NOT 10%.
Whichever stage applies, every one of the following must be changed
together — do not copy this template and change only the target column:
  1. became_10_deal_date  → became_<N>_deal_date   (all 4 occurrences:
     create_fy, create_quarter, create_month, and the WHERE filter)
  2. amount_target_10     → amount_target_<N>
  3. deals_target_10      → deals_target_<N>
  4. amount_target_10_l1 / _committed and deals_target_10_l1 / _committed
     → the matching _<N>_l1 / _<N>_committed columns, if that tier applies
Leaving became_10_deal_date in place while only swapping the target
columns silently produces an attainment % for the WRONG stage.

WITH actuals AS (
  SELECT
    toYear(became_10_deal_date) + if(toMonth(became_10_deal_date)>=4,1,0) AS create_fy,
    CASE WHEN toMonth(became_10_deal_date) IN (1,2,3)    THEN 'Q4'
         WHEN toMonth(became_10_deal_date) IN (4,5,6)    THEN 'Q1'
         WHEN toMonth(became_10_deal_date) IN (7,8,9)    THEN 'Q2'
         WHEN toMonth(became_10_deal_date) IN (10,11,12) THEN 'Q3'
    END AS create_quarter,
    LEFT(formatDateTime(became_10_deal_date,'%M'),3) AS create_month,
    CASE WHEN region='japac' THEN 'JAPAC' WHEN region='Africa' THEN 'Middle East'
         WHEN region='india___sea' THEN 'ISEA' ELSE region END AS region,
    CASE
      WHEN deal_source_rollup IN ('Executive Outreach','Investor') THEN 'Executive Outreach'
      WHEN deal_source_rollup IN ('Marketing','Customer Success',
           'AE Outbound','Inception','Hyperscaler') THEN deal_source_rollup
      WHEN deal_source_rollup = 'BDR Outbound' THEN 'BDR'
      WHEN deal_source_rollup = 'Partner'       THEN 'Partner - Non Hyperscaler'
      ELSE 'Other'
    END AS deal_source_rollup,
    SUM(amount)            AS actual_amount,
    countDistinct(deal_id) AS actual_deals
  FROM hs_analytics.deals FINAL
  WHERE pipeline = 'default'
    AND deal_stage <> 'Duplicate Record'
    AND CASE WHEN deal_type IS NULL THEN 'Not Assigned' ELSE deal_type END
        NOT IN ('Partner-Led SMB')
    AND toInt64(deal_id) IN (
      SELECT DISTINCT toInt64(deal_id_hs) FROM kore_ai_hubspot.gs_deal_ids_hs
    )
    AND became_10_deal_date >= '2025-04-01'
    AND (became_10_deal_date != '1900-01-01'
         OR became_20_deal_date != '1900-01-01'
         OR became_30_deal_date != '1900-01-01'
         OR became_40_deal_date != '1900-01-01'
         OR became_60_deal_date != '1900-01-01'
         OR became_75_deal_date != '1900-01-01'
         OR deal_stage IN ('10% - Discovery','20% - Solution','30% - Proof',
                           '40% - Proposal','60% - Price Negotiation',
                           '75% - Contract Review','90% - Deal Desk Review','Closed Won'))
  GROUP BY 1,2,3,4,5
),
targets AS (
  SELECT
    CAST(fy AS INT) AS create_fy,
    quarter         AS create_quarter,
    month           AS create_month,
    region,
    CASE WHEN source = 'Hyperscalers'                        THEN 'Hyperscaler'
         WHEN source IN ('Executive Outreach','Investor')    THEN 'Executive Outreach'
         WHEN source = 'Partner - Excluding Hyperscalers'    THEN 'Partner - Non Hyperscaler'
         ELSE source END AS deal_source_rollup,
    SUM(toFloat64OrZero(amount_target_10))           AS amount_target,
    SUM(toFloat32OrZero(deals_target_10))            AS deals_target,
    SUM(toFloat64OrZero(amount_target_10_l1))        AS amount_target_l1,
    SUM(toFloat32OrZero(deals_target_10_l1))         AS deals_target_l1,
    SUM(toFloat64OrZero(amount_target_10_committed)) AS amount_target_committed,
    SUM(toFloat32OrZero(deals_target_10_committed))  AS deals_target_committed
  FROM kore_ai_hubspot.gs_pipeline_quotas_v1
  GROUP BY 1,2,3,4,5
)
SELECT
  t.create_fy, t.create_quarter, t.create_month,
  t.region, t.deal_source_rollup,
  coalesce(a.actual_amount, 0) AS actual_amount,
  coalesce(a.actual_deals,  0) AS actual_deals,
  t.amount_target, t.deals_target,
  round(coalesce(a.actual_deals, 0)  / nullIf(t.deals_target,  0) * 100, 1) AS deals_attainment_pct,
  round(coalesce(a.actual_amount, 0) / nullIf(t.amount_target, 0) * 100, 1) AS amount_attainment_pct
FROM targets t
LEFT JOIN actuals a
  ON  t.create_fy          = a.create_fy
  AND t.create_quarter     = a.create_quarter
  AND t.create_month       = a.create_month
  AND t.deal_source_rollup = a.deal_source_rollup
  AND t.region             = a.region
-- WHERE t.create_fy = 2027 AND t.create_quarter = 'Q1'

NOTE — FOR 20% ATTAINMENT:
  Replace became_10_deal_date → became_20_deal_date in the actuals CTE.
  Replace amount_target_10 → amount_target_20, deals_target_10 → deals_target_20.

NOTE — FOR PARTNER ATTAINMENT:
  Use kore_ai_hubspot.gs_partner_targets_region_wise with l2_ prefix columns.
  Filter by partner_team_type IN ('Hyperscaler','GSI/SI','Reseller/BPO/TSD') as needed.
  committed_amount_target_10 and _5 do NOT exist in this table.

NOTE — FOR BDR ATTAINMENT:
  Filter actuals WHERE deal_source_rollup = 'BDR' (after source mapping).
  Use the same Pattern C template with that additional source filter.

═══════════════════════════════════════════════════════════════
§9  FISCAL YEAR AND DATE RULES
═══════════════════════════════════════════════════════════════

Kore.ai fiscal year starts April 1:
  FY27 = Apr 2026 – Mar 2027  (current default)
  FY26 = Apr 2025 – Mar 2026
  FY25 = Apr 2024 – Mar 2025

  FY formula: toYear(date) + if(toMonth(date) >= 4, 1, 0)
  Q1 = Apr/May/Jun | Q2 = Jul/Aug/Sep | Q3 = Oct/Nov/Dec | Q4 = Jan/Feb/Mar

Target table stores: fy as 'FY27', quarter as 'Q1', month as 'Apr' etc.
When joining actuals to targets:
  CAST(fy AS INT) must equal the numeric FY (2027, not 'FY27').

Date column standard — ALWAYS:
  CAST(LEFT(coalesce(column_name, '1900-01-01'), 10) AS DATE)

Sentinel '1900-01-01' = date was never set (NULL equivalent).
Check: column != '1900-01-01'

═══════════════════════════════════════════════════════════════
§10  DIMENSION MAPPINGS (consistent across all patterns)
═══════════════════════════════════════════════════════════════

REGION (apply in SELECT, not WHERE):
  'japac'       → 'JAPAC'
  'Africa'      → 'Middle East'
  'india___sea' → 'ISEA'

INDUSTRY:
  ('Financial Services','Banking','Insurance')             → 'Financial Services'
  ('Manufacturing Discreet','Manufacturing Process','CPG') → 'Manufacturing'
  ('Hi-Tech','Telecom / Media / Entertainment')           → 'TMT'
  ('Business Services','Government','Energy & Utilities',
   'Education','Restaurants','null','Energy') or NULL     → 'Other'

ACCOUNT PRIORITY:
  ('P1','P2','P3','P4') → 'P1-P4'
  ('P5','P6','P7')      → 'P5-P7'
  ('P8','P9','P10')     → 'P8-P10'
  NULL → NULL (Pattern A) or 'No Priority' (Pattern B/C detail)

INCEPTION DEALS:
  Only exclude if user explicitly asks. Default: include all deals.

═══════════════════════════════════════════════════════════════
§11  MQL RULES
═══════════════════════════════════════════════════════════════
MQL actuals from hs_analytics.contacts FINAL require ALL THREE:
  1. lifecycle_stage = 'marketingqualifiedlead'
     AND date_entered_marketing_qualified_lead_lifecycle_stage_pipeline IS NOT NULL
  2. company_priority IN ('P1','P2','P3','P4','P5','P6','P7')
  3. lead_status != 'Bad Data'

NEVER omit any of the 3 MQL filters — missing any one inflates counts.

MQL ACTUALS PATTERN:
  SELECT
    region,
    original_source,
    toYYYYMM(toDate(LEFT(date_entered_marketing_qualified_lead_lifecycle_stage_pipeline,10))) AS ym,
    countDistinct(contact_id) AS mql_count
  FROM hs_analytics.contacts FINAL
  WHERE lifecycle_stage = 'marketingqualifiedlead'
    AND date_entered_marketing_qualified_lead_lifecycle_stage_pipeline IS NOT NULL
    AND company_priority IN ('P1','P2','P3','P4','P5','P6','P7')
    AND lead_status != 'Bad Data'
    AND <date range filter on date_entered_... column>
  GROUP BY region, original_source, ym

MQL TARGET PATTERN — always filter by exact quarter, NEVER divide annual by 4:
  SELECT
    region,
    original_source,
    SUM(toFloat64OrZero(mql_target)) AS mql_tgt
  FROM kore_ai_hubspot.gs_marketing_targets
  WHERE fy = 'FY27' AND quarter = 'Q1'
  GROUP BY region, original_source

FILTER CONFIRMATION for MQL queries — include in Filters Applied block:
  - Company Priority: P1–P7 (excludes unranked contacts)
  - Lead Status: excludes 'Bad Data'
  - MQL date range: <the date range used>

═══════════════════════════════════════════════════════════════
§12  DASHBOARD DEFINITIONS
═══════════════════════════════════════════════════════════════
When a user asks about a specific dashboard, apply the correct logic below.
If unclear which dashboard, ask the user.

── DASHBOARD 1: EOP (End-of-Period) ───────────────────────────
PURPOSE: Tracks pipeline health and attainment against EOP targets.
KEY METRICS: EOP Pipeline Value, EOP Target, EOP Attainment %, stage-wise and region-wise EOP.
FILTERS: MANDATORY_BASE_FILTERS + close_date within current quarter end window
         + deal_stage IN active stages (20%–75%)

── DASHBOARD 2: EXEC KPI ──────────────────────────────────────
PURPOSE: Senior leadership view of pipeline performance.
KEY METRICS: Total Active Pipeline ($M), Closed Won ($M), CW Attainment %,
             Win Rate %, Pipeline Coverage, New Logo Count.

── DASHBOARD 3: CS (Customer Success) ─────────────────────────
PURPOSE: Tracks renewals, upsells, expansions, CS team performance.

── DASHBOARD 4: GLOBAL PIPELINE GOVERNANCE ────────────────────
PURPOSE: Executive governance view across all regions, sources, partner types.

── DASHBOARD 5: GLOBAL PIPEGEN ────────────────────────────────
PURPOSE: Org-wide pipeline generation performance and attainment.
KEY METRICS: 5%, 10%, 20% Pipeline Amount and Deal Count, Actual vs Target,
             Attainment %, Pipeline Trend, Funnel Conversion, Region/Source/Industry.
FILTERS: MANDATORY_BASE_FILTERS + stage/date filters + gs_pipeline_quotas_v1 targets.

── DASHBOARD 6: PARTNERSHIP ───────────────────────────────────
PURPOSE: Partner-generated pipeline, target attainment, partner ecosystem performance.
KEY METRICS: Partner 5%/10%/20% Pipeline, Partner Attainment %, PSD Performance,
             Hyperscaler Performance, GSI/SI, Reseller/BPO/TSD, Partner Funnel Conversion.
FILTERS: MANDATORY_BASE_FILTERS + Partner/Hyperscaler sources + partner target tables.

── DASHBOARD 7: MARKETING ─────────────────────────────────────
PURPOSE: MQL actuals vs target, marketing pipeline, marketing attainment.
KEY METRICS: MQL Actual vs Target, Marketing-Sourced Pipeline, Deal Count,
             Source-wise Performance, Region-wise Performance, MQL Conversion Rate.
FILTERS: Mandatory MQL filters (§11) + gs_marketing_targets + selected fiscal period.

── DASHBOARD 8: AE FOCUS ──────────────────────────────────────
PURPOSE: AE pipeline, quota attainment, and sales performance.
KEY METRICS: Active Pipeline, CW ARR, CW Deal Count, Actual vs Quota,
             Attainment %, Win Rate, Pipeline Coverage, Avg Deal Size, Sales Cycle.
FILTERS: MANDATORY_BASE_FILTERS + gs_closed_won_quotas + AE filters.

── DASHBOARD 9: BDR FOCUS ─────────────────────────────────────
PURPOSE: BDR pipeline generation and conversion performance.
KEY METRICS: Meetings Created, Opportunities Generated, 5%/10%/20% PipeGen,
             Stage Conversion, BDR Performance, Region-wise, Target Attainment.
FILTERS: MANDATORY_BASE_FILTERS + BDR ownership filters + selected fiscal period.

═══════════════════════════════════════════════════════════════
§13  SQL GENERATION GUARDRAILS
═══════════════════════════════════════════════════════════════
1.  SELECT / WITH only — no INSERT, UPDATE, DELETE, DROP, ALTER.
2.  FINAL on every hs_analytics.* table.
3.  All MANDATORY BASE FILTERS (§3) on every deals query.
4.  countDistinct(deal_id) — never count() or count(deal_id).
5.  No LIMIT unless user says "top N" or "first N".
6.  All target table numeric columns: SUM(toFloat64OrZero(col)).
7.  Division: always wrap denominator with nullIf(denom, 0).
8.  Date columns: CAST(LEFT(coalesce(col,'1900-01-01'),10) AS DATE).
9.  NEVER compute a quarterly target by dividing annual by 4.
10. State total row count in every answer.

⚠️  IMPORTANT FOR RULES.PY VALIDATOR:
Pattern A uses became_N_deal_date in OR conditions (cumulative counting),
NOT as a cohort anchor with stage exclusions. The rules.py cohort checks
were designed for true cohort funnel queries, not Pattern A.
For Pattern A SQL, add a comment at the top of the query:
  -- Pattern A: cumulative stage counting, not cohort funnel

═══════════════════════════════════════════════════════════════
§14  OUTPUT FORMATTING
═══════════════════════════════════════════════════════════════
- Clean markdown. Tables for data. Bold key numbers.
- Never fabricate numbers. Never run destructive SQL.
- Format dollar amounts: round(sum(amount)/1e6, 1) as $M.
- Always complete full response — never truncate mid-table.

After every DB-backed answer, append:

---
**Filters Applied:**
- Pattern used: [A — Cumulative Funnel / B — Deal Detail / C — Attainment]
- FY anchor column: [which became_N_deal_date or close_date was used]
- [list all active filters: FY, quarter, region, source, stage, etc.]

Please verify these filters match your expectation.
---

═══════════════════════════════════════════════════════════════
§15  VISUAL / CHART GENERATION
═══════════════════════════════════════════════════════════════
Charts are generated automatically by the system after your response.
Do NOT write Chart.js, SVG, or any ```html chart code yourself.

═══════════════════════════════════════════════════════════════
LIVE DATABASE SCHEMA (auto-injected below)
═══════════════════════════════════════════════════════════════
{schema}
"""


def _extract_section(text: str, n: int) -> str:
    """
    Pulls §N (header through body) out of the full prompt text, up to but
    not including the next §-header or the LIVE DATABASE SCHEMA block.
    Extraction is driven by the section markers already in the text, not
    hand-copied, so splitting the prompt can't silently drop or duplicate
    rule content that exists in only one place today.
    """
    pattern = rf'═+\n§{n}\b.*?(?=\n═+\n§\d|\n═+\nLIVE DATABASE SCHEMA|\Z)'
    m = re.search(pattern, text, re.S)
    return m.group(0).strip() if m else ""


def _split_section_2(sec2: str) -> Tuple[str, str, str]:
    """
    §2 (TOOL USAGE) bundles three logically separate things: which tool to
    call and when (SQL Agent's job), the export-intent marker (a
    conversational routing decision, not a SQL decision), and the
    chart-spec instruction (now the Narrator Agent's job, since
    choose_chart_spec moved to the Narrator's tool set). Splits it into
    (tool_usage_only, export_intent_block, chart_spec_block).
    """
    export_m = re.search(r'EXPORT INTENT RULE:.*?(?=\nCHART SPEC RULE:|\Z)', sec2, re.S)
    chart_m = re.search(r'CHART SPEC RULE:.*', sec2, re.S)
    tool_only = sec2
    if export_m:
        tool_only = tool_only.replace(export_m.group(0), '')
    if chart_m:
        tool_only = tool_only.replace(chart_m.group(0), '')
    return (
        tool_only.strip(),
        export_m.group(0).strip() if export_m else "",
        chart_m.group(0).strip() if chart_m else "",
    )


def _build_sql_agent_prompt() -> str:
    """
    SQL Agent's prompt: everything needed to turn a question into
    validated SQL and retrieve data, plus the greeting/export fast paths
    (§1, export-intent) so a turn that needs no data at all — "hi", "export
    that as PDF" — can be answered directly without ever reaching the
    Narrator Agent. Includes the live schema; the Narrator Agent does not
    need it (it only ever sees already-labeled result rows, never raw
    columns it has to interpret against a schema).
    """
    full = _build_system_prompt()
    sec1 = _extract_section(full, 1)
    sec2 = _extract_section(full, 2)
    tool_only, export_block, _chart_block = _split_section_2(sec2)
    sections = "\n\n".join(
        s for s in (
            tool_only,
            export_block,
            *[_extract_section(full, n) for n in range(3, 14)],
        ) if s
    )
    schema_start = full.find("LIVE DATABASE SCHEMA")
    schema_block = full[schema_start - 66:] if schema_start != -1 else ""

    patches_block = ""
    if _SQL_AGENT_PATCHES:
        patches_block = (
            "\n\n═══════════════════════════════════════════════════════════════\n"
            "§16  HUMAN-APPROVED CLARIFICATIONS (see /debug/proposals for origin)\n"
            "═══════════════════════════════════════════════════════════════\n"
            + "\n".join(f"- {p}" for p in _SQL_AGENT_PATCHES)
        )

    return f"""
You are DIUD's SQL Agent — responsible for understanding a business question,
writing correct ClickHouse SQL, and retrieving validated data. You never
fabricate numbers and never run destructive SQL. If the question needs no
data at all (a greeting, an export request on already-fetched data, a
question about your own capabilities), handle it directly per the rules
below without calling query_clickhouse.

Once you have a clean, rule-compliant result, STOP — do not write the final
narrative answer yourself. A separate Narrator Agent turns your retrieved
data into the response the user sees.

RULE PRIORITY ORDER (highest → lowest):
  1. Greeting Rule (§1)
  2. Safety — SELECT/WITH only, no destructive SQL ever
  3. MANDATORY_BASE_FILTERS (§3)
  4. Tool Usage (§2)
  5. Business & SQL Rules (§4–§13)
  6. Human-approved clarifications (§16), if any — same authority as §4–§13,
     added later after a real recurring mistake was reviewed and confirmed

{sec1}

{sections}

{schema_block}
{patches_block}
"""


def _build_narrator_agent_prompt() -> str:
    """
    Narrator Agent's prompt: only what's needed to turn already-verified
    data into a clear written answer and (optionally) pick a chart. No
    schema, no SQL rules, no table names — by the time this agent runs,
    all of that ambiguity has already been resolved by the SQL Agent and
    the rule engine. This is deliberately the smaller of the two prompts.
    """
    full = _build_system_prompt()
    sec1 = _extract_section(full, 1)
    sec2 = _extract_section(full, 2)
    _tool_only, _export_block, chart_block = _split_section_2(sec2)
    sec14 = _extract_section(full, 14)
    sec15 = _extract_section(full, 15)

    patches_block = ""
    if _NARRATOR_AGENT_PATCHES:
        patches_block = (
            "\n\n═══════════════════════════════════════════════════════════════\n"
            "HUMAN-APPROVED CLARIFICATIONS (see /debug/proposals for origin)\n"
            "═══════════════════════════════════════════════════════════════\n"
            + "\n".join(f"- {p}" for p in _NARRATOR_AGENT_PATCHES)
        )

    return f"""
You are DIUD's Narrator Agent — responsible for turning already-verified
data into a clear, executive-grade written answer. You do not have database
access and cannot write SQL; the data handed to you has already been
retrieved and checked by a separate SQL Agent. Never invent, adjust, or
recompute any number — use only what appears in the data block you're
given, in a [VERIFIED TOTALS] block, or from a compute_verified_metric
result already included in that data.

RULE PRIORITY ORDER (highest → lowest):
  1. Greeting Rule (§1) — only relevant if you're ever invoked with no data
  2. Formatting & Chart Rules (§14–§15)
  3. Chart spec selection (below)
  4. Human-approved clarifications, if any — added later after a real
     recurring mistake was reviewed and confirmed

{sec1}

{chart_block}

{sec14}

{sec15}
{patches_block}
"""


_SQL_AGENT_PROMPT = _build_sql_agent_prompt()
_NARRATOR_AGENT_PROMPT = _build_narrator_agent_prompt()


def _cached_sql_agent_prompt() -> list:
    """
    Prompt-caching wrapper for the SQL Agent's prompt. Reads the CURRENT
    _SQL_AGENT_PROMPT global each call (not a stale copy), so it stays
    correct after refresh_schema() or on_startup() reassign it.
    """
    return [
        {
            "type": "text",
            "text": _SQL_AGENT_PROMPT,
            "cache_control": {"type": "ephemeral"},
        }
    ]


def _cached_narrator_agent_prompt() -> list:
    """
    Prompt-caching wrapper for the Narrator Agent's prompt. Much smaller
    than the SQL Agent's — no schema, no SQL rules — so this was already
    cheap before caching; caching still helps across the many narration
    calls within a single busy session.
    """
    return [
        {
            "type": "text",
            "text": _NARRATOR_AGENT_PROMPT,
            "cache_control": {"type": "ephemeral"},
        }
    ]

# =============================================================================
# Data Normalizer — Stage 2 of the target architecture.
#
# Takes the raw JSON payload from the ClickHouse proxy and turns it into the
# one normalized shape every downstream consumer relies on: (columns, rows),
# where rows is always List[dict]. This used to be inline logic buried
# inside run_clickhouse_query() (steps: figure out if the proxy returned a
# list or a dict, find the actual row array under whatever key it used,
# reconcile column names from a separate columns/meta field if rows came
# back as bare arrays instead of dicts). Extracted here so it's independently
# testable and so the shape guarantee it provides — dicts, real column
# names, matching every row — is explicit rather than implicit.
#
# Both Chart Builder (charts.py:build_chart_html) and everything downstream
# of Fact Binder consume the QueryResult this produces — there is no
# separate normalization path for either.
# =============================================================================
def normalize_query_response(payload) -> Tuple[Optional[List[str]], Optional[List[dict]], Optional[str]]:
    """
    Returns (columns, rows, error). Exactly one of (columns, rows) or
    error will be set. `error` is a human-readable string suitable for
    returning directly as the tool_result when normalization isn't
    possible (e.g. the proxy returned something with no row data).
    """
    if isinstance(payload, list):
        rows = payload
        api_columns = None
    elif isinstance(payload, dict):
        rows = payload.get("data") or payload.get("rows") or payload.get("result") or payload.get("results")
        api_columns = payload.get("columns") or payload.get("meta") or payload.get("column_names")
        if rows is None:
            return None, None, json.dumps(payload, indent=2, default=str)[:3000]
    else:
        return None, None, f"Unexpected response type: {type(payload)}"

    if not rows:
        return None, None, "Query returned 0 rows."

    if isinstance(rows[0], dict):
        columns = list(rows[0].keys())
        norm_rows = rows
    else:
        if api_columns and len(api_columns) == len(rows[0]):
            columns = [c["name"] if isinstance(c, dict) else c for c in api_columns]
        else:
            columns = [f"col_{i}" for i in range(len(rows[0]))]
        norm_rows = [dict(zip(columns, r)) for r in rows]

    return columns, norm_rows, None


# =============================================================================
# ClickHouse query runner
# =============================================================================
def run_clickhouse_query(sql: str, session_id: Optional[str] = None) -> str:
    base_url = _base_url()
    token    = _token()

    if not base_url:
        return "DATABASE CONNECTION FAILED: CLICKHOUSE_API_URL is not set."
    if not token:
        return "DATABASE CONNECTION FAILED: CLICKHOUSE_API_TOKEN is not set."
        
    cleaned = re.sub(r'^\s*--[^\n]*\n', '', sql.strip(), flags=re.MULTILINE).strip()
    stripped = cleaned.upper()
    if not (stripped.startswith("SELECT") or stripped.startswith("WITH")):
        return "ERROR: Only SELECT/WITH queries are permitted."
    for kw in FORBIDDEN_KEYWORDS:
        if re.search(rf'\b{kw}\b', stripped):
            return f"ERROR: Forbidden keyword: {kw}"
    

    print(f"🔍 SQL (session={session_id}) → {sql[:200]}")

    try:
        resp = httpx.post(
            f"{base_url}/query",
            headers=_auth_headers(),
            json={"query": sql},
            timeout=60,
        )

        if resp.status_code == 401:
            return "DATABASE CONNECTION FAILED: 401 Unauthorized."
        if resp.status_code == 403:
            return "DATABASE CONNECTION FAILED: 403 Forbidden."
        if resp.status_code == 422:
            return f"ERROR: Proxy rejected query (422): {resp.text[:400]}"
        if resp.status_code == 500:
            return f"DATABASE ERROR: HTTP 500: {resp.text[:400]}"
        if resp.status_code != 200:
            return f"DATABASE ERROR: HTTP {resp.status_code} — {resp.text[:300]}"

        payload = resp.json()

        columns, norm_rows, norm_error = normalize_query_response(payload)
        if norm_error:
            return norm_error

        total_rows = len(norm_rows)

        if session_id:
            existing = _SESSION_STORE.get(session_id)
            # Always store the most recent query result for chart/export purposes.
            # Use a SEPARATE key (chat_session_id + ":latest") for the most recent
            # result specifically, while preserving the old "largest result wins"
            # behavior under the original session_id for export/CSV purposes.
            _store_result(f"{session_id}:latest", QueryResult(
                sql             = sql,
                columns         = columns,
                rows            = norm_rows,
                total_rows      = total_rows,
                captured_at     = datetime.utcnow().isoformat() + "Z",
                filters_applied = _extract_filters_from_sql(sql),
            ))
            # Only overwrite the EXPORT-facing session result if this is larger.
            # This prevents a trivial follow-up query from replacing a large
            # deal-list result meant for CSV/PDF export with a tiny one.
            if not existing or total_rows > existing.total_rows:
                _store_result(session_id, QueryResult(
                    sql             = sql,
                    columns         = columns,
                    rows            = norm_rows,
                    total_rows      = total_rows,
                    captured_at     = datetime.utcnow().isoformat() + "Z",
                    filters_applied = _extract_filters_from_sql(sql),
                ))
            else:
                print(f"   ⏭️  Skipping EXPORT session store update — {total_rows} row(s) won't overwrite existing {existing.total_rows} rows.")

        CHAT_DISPLAY_LIMIT = 100
        header = " | ".join(columns)
        lines  = [header, "-" * min(len(header), 140)]
        for row in norm_rows[:CHAT_DISPLAY_LIMIT]:
            lines.append(" | ".join(str(row.get(c, "")) for c in columns))

        if total_rows > CHAT_DISPLAY_LIMIT:
            lines.append(
                f"\n📊 **Showing {CHAT_DISPLAY_LIMIT} of {total_rows} rows.** "
                f"Say **\"export these deals as CSV\"** or **\"export as PDF\"** "
                f"to download all {total_rows} records."
            )

        result = "\n".join(lines)
        print(f"   ✅ {total_rows} rows returned. Session store updated.")
        return result

    except httpx.ConnectError as e:
        return f"DATABASE CONNECTION FAILED: Could not reach {base_url}. {e}"
    except httpx.TimeoutException:
        return "DATABASE CONNECTION FAILED: Query timed out after 60 seconds."
    except Exception as exc:
        traceback.print_exc()
        return f"DATABASE CONNECTION FAILED: {type(exc).__name__}: {exc}"


def _extract_filters_from_sql(sql: str) -> str:
    sql_upper = sql.upper()
    filters = []
    if "PIPELINE = 'DEFAULT'" in sql_upper:
        filters.append("Pipeline: default")
    
    # A genuine cohort query anchors on exactly ONE became_X_deal_date
    # stage plus a deal_stage NOT IN (...) exclusion of earlier stages.
    # A normal Pattern A funnel query legitimately references ALL 7
    # became_X_deal_date columns as part of ordinary cumulative OR-chain
    # stage counting — that is NOT cohort filtering, and labeling every
    # one of those 7 columns "Cohort: N% qualified deals" (as the old
    # code did) produced a garbled, factually wrong subtitle like
    # "Cohort: 5% qualified deals; Cohort: 10% qualified deals; ..." on
    # every ordinary funnel chart.
    referenced_stages = sorted(set(re.findall(r'BECAME_(\d+)_DEAL_DATE', sql_upper)), key=int)
    if len(referenced_stages) == 1 and re.search(r'DEAL_STAGE\s+NOT\s+IN', sql_upper):
        filters.append(f"Cohort: {referenced_stages[0]}% starting stage")
    elif referenced_stages:
        # Report the TRUE FY anchor (the column actually used inside
        # toYear(...)), not just "however many stage columns happen to be
        # referenced" — matches the same anchor-detection logic used by
        # rules.py's pattern_a_stage_anchor / pattern_c_stage_anchor checks.
        anchor_match = re.search(r'TOYEAR\(\s*BECAME_(\d+)_DEAL_DATE\s*\)', sql_upper)
        anchor_stage = anchor_match.group(1) if anchor_match else referenced_stages[0]
        filters.append(f"FY anchor stage: {anchor_stage}%")
    
    if "CLOSE_DATE" in sql_upper and "2026-04-01" in sql:
        filters.append("FY27 active pipeline")
    if "DEAL_STAGE" in sql_upper:
        m = re.search(r"deal_stage\s+IN\s*\(([^)]+)\)", sql, re.IGNORECASE)
        if m:
            filters.append(f"Stage filter: {m.group(1)[:60]}")
    # Only treat `region = '...'` as a REAL filter when it appears in a WHERE/AND/OR
    # boolean-predicate position. The normalized-region CASE expression present in
    # nearly every base query (e.g. "CASE WHEN region='japac' THEN 'JAPAC' ...")
    # also matches a naive `region\s*=\s*'...'` search and was being misread as a
    # user-applied filter on every query, producing a spurious "Region: japac"
    # subtitle even when no region filter was applied.
    region_filter_match = re.search(
        r"(?:WHERE|AND|OR)\s+region\s*=\s*'([^']+)'", sql, re.IGNORECASE
    )
    if region_filter_match:
        filters.append(f"Region: {region_filter_match.group(1)}")
    if "COMPANY_PRIORITY" in sql_upper:
        filters.append("Company Priority: P1–P7")
    if "LEAD_STATUS" in sql_upper and "BAD DATA" in sql_upper:
        filters.append("Lead Status: excludes 'Bad Data'")
    if "LIFECYCLE_STAGE" in sql_upper and "MARKETINGQUALIFIEDLEAD" in sql_upper:
        filters.append("Lifecycle: MQL only")
    return "; ".join(filters) if filters else "Standard base filters applied"

def _is_numericish(v) -> bool:
    try:
        float(v)
        return True
    except (TypeError, ValueError):
        return False


def _to_float(v) -> float:
    try:
        return float(v)
    except (TypeError, ValueError):
        return 0.0


def _compute_deterministic_headline(stored: Optional["QueryResult"]) -> Optional[str]:
    """For small aggregate results, compute totals/averages in Python so Claude
    never has to do the arithmetic itself when writing the summary."""
    if not stored or not stored.rows or len(stored.rows) > 20:
        return None

    numeric_cols = [
        c for c in stored.columns
        if any(_is_numericish(r.get(c)) for r in stored.rows)
    ]
    if not numeric_cols:
        return None

    parts = []
    for col in numeric_cols:
        total = sum(_to_float(r.get(col)) for r in stored.rows)
        parts.append(f"{col}_total={total:,.2f}")
        if len(stored.rows) > 1:
            avg = total / len(stored.rows)
            parts.append(f"{col}_avg={avg:,.2f}")
    return " | ".join(parts)


_METRIC_FNS: Dict[str, Callable[[List[dict]], float]] = {
    "win_rate": lambda rows: (
        sum(1 for r in rows if r.get("deal_stage") == "Closed Won")
        / max(sum(1 for r in rows if r.get("deal_stage") in ("Closed Won", "Closed Lost")), 1)
        * 100
    ),
    "active_pipeline_total": lambda rows: sum(
        _to_float(r.get("amount")) for r in rows
        if r.get("deal_stage") in (
            "20% - Solution", "30% - Proof", "40% - Proposal",
            "60% - Price Negotiation", "75% - Contract Review",
        )
    ) / 1_000_000,
    "avg_deal_size": lambda rows: (
        sum(_to_float(r.get("amount")) for r in rows) / max(len(rows), 1)
    ),
    "closed_won_total": lambda rows: sum(
        _to_float(r.get("amount")) for r in rows if r.get("deal_stage") == "Closed Won"
    ) / 1_000_000,
    "mql_count": lambda rows: float(len(rows)),
    "attainment_pct": lambda rows: (
        sum(_to_float(r.get("achieved_m")) for r in rows)
        / max(sum(_to_float(r.get("target_m")) for r in rows), 0.0001) * 100
    ),
}


def compute_verified_metric(metric_name: str, session_id: Optional[str]) -> str:
    # This metric must be computed from the query Claude just ran THIS
    # turn, not whichever past query in the session happened to return
    # the most rows (the old largest-wins cache).
    stored = (
        (_get_result(f"{session_id}:latest") or _get_result(session_id))
        if session_id else None
    )
    if not stored or not stored.rows:
        return "ERROR: No stored query result for this session — run a query first."
    fn = _METRIC_FNS.get(metric_name)
    if not fn:
        return f"ERROR: Unknown metric '{metric_name}'."
    try:
        value = fn(stored.rows)
        return f"VERIFIED {metric_name} = {value:,.2f}"
    except Exception as e:
        return f"ERROR computing {metric_name}: {e}"
        
        
        
def _parse_rows_for_validation(query_result: str, session_id: Optional[str]) -> List[dict]:
    """
    run_clickhouse_query() returns a human-readable pipe-delimited text table,
    not JSON, so we can't parse query_result directly. Instead, pull the
    structured rows it already stored in _SESSION_STORE for this session.
    """
    if not session_id:
        return []
    stored = _SESSION_STORE.get(session_id)
    if not stored:
        return []
    return stored.rows


# =============================================================================
# Startup
# =============================================================================
@app.on_event("startup")
async def on_startup():
    global _SQL_AGENT_PROMPT, _NARRATOR_AGENT_PROMPT
    try:
        discover_schema()
        _SQL_AGENT_PROMPT = _build_sql_agent_prompt()
        _NARRATOR_AGENT_PROMPT = _build_narrator_agent_prompt()
    except Exception as e:
        print(f"⚠️  Schema discovery failed: {e}")
    _cleanup_sessions()
    print("🚀 DIUD v5 started — SQL Agent / Narrator Agent split, session-store export enabled.")


# =============================================================================
# Claude tool definition
# =============================================================================
_QUERY_TOOL = {
    "name": "query_clickhouse",
    "description": (
        "Execute a SELECT query against ClickHouse. "
        "Use for deal pipeline, AE performance, win/loss, regions, stages, MQL metrics. "
        "ALWAYS use fully-qualified table names. "
        "Relay DATABASE CONNECTION FAILED errors directly to the user."
    ),
    "input_schema": {
        "type": "object",
        "properties": {
            "sql": {
                "type": "string",
                "description": (
                    "Valid ClickHouse SELECT or WITH query. "
                    "For deal LIST queries: NO LIMIT unless user asks for 'top N'. "
                    "Return all matching rows — the system displays them safely."
                ),
            }
        },
        "required": ["sql"],
    },
}

_RULEBOOK_TOOL = {
    "name": "lookup_business_rule",
    "description": (
        "Look up the exact business rule for a specific metric/topic BEFORE "
        "writing SQL for it. ALWAYS call this before generating SQL for MQL, "
        "active pipeline, cohort funnels, attainment/targets, closed won, "
        "partner targets, or dashboard-specific questions."
    ),
    "input_schema": {
        "type": "object",
        "properties": {
            "topic": {
                "type": "string",
                "enum": [
                    "mql", "active_pipeline", "cohort_funnel", "attainment",
                    "closed_won", "partner_targets", "dashboard_definitions",
                ],
            }
        },
        "required": ["topic"],
    },
}

_COMPUTE_METRIC_TOOL = {
    "name": "compute_verified_metric",
    "description": (
        "Compute a verified metric (win_rate, active_pipeline_total, "
        "avg_deal_size, closed_won_total, mql_count, attainment_pct) from the "
        "current session's query results. ALWAYS use this instead of computing "
        "percentages or sums yourself from the raw row dump."
    ),
    "input_schema": {
        "type": "object",
        "properties": {
            "metric_name": {
                "type": "string",
                "enum": [
                    "win_rate", "active_pipeline_total", "avg_deal_size",
                    "closed_won_total", "mql_count", "attainment_pct",
                ],
            },
        },
        "required": ["metric_name"],
    },
}

_CHART_SPEC_TOOL = {
    "name": "choose_chart_spec",
    "description": (
        "Call this AFTER query_clickhouse returns a result you're about to summarize, "
        "if a chart would help the user (funnel/stage breakdowns, actual-vs-target, "
        "top-N comparisons, breakdowns by region/source/competitor/etc). Do NOT call "
        "this for raw deal lists or single-number answers. "
        "You choose WHAT to chart — the chart type and which real columns to use. "
        "You do NOT choose any numbers, percentages, or pixel widths — those are always "
        "computed independently from the actual query result, never from what you write here. "
        "label_column and value_column MUST be exact column names from the query result you "
        "just got back — do not invent or rename columns."
    ),
    "input_schema": {
        "type": "object",
        "properties": {
            "chart_type": {
                "type": "string",
                "enum": ["funnel", "attainment", "donut", "bar_h"],
                "description": (
                    "funnel: cumulative stage counts (Pattern A). "
                    "attainment: actual vs target columns present (Pattern C). "
                    "donut: a breakdown into 6 or fewer categories. "
                    "bar_h: a ranked breakdown into 7+ categories, or a top-N comparison."
                ),
            },
            "label_column": {
                "type": "string",
                "description": "Exact column name to use as category labels (ignored for funnel/attainment, which use fixed columns).",
            },
            "value_column": {
                "type": "string",
                "description": "Exact column name to use as the chart values (ignored for funnel/attainment).",
            },
            "exclude_values": {
                "type": "array",
                "items": {"type": "string"},
                "description": "Category values to exclude from the chart scale (e.g. 'N/A', 'Unknown') so real categories stay readable. Excluded rows are footnoted, not dropped from the underlying data.",
            },
            "title": {
                "type": "string",
                "description": "Short chart title, e.g. 'Lost Deals by Competitor'.",
            },
            "insight_note": {
                "type": "string",
                "description": "Optional one-line callout under the chart, e.g. 'Microsoft is our top competitive loss this quarter'. Must be a fact already visible in the query result — do not introduce new numbers here.",
            },
        },
        "required": ["chart_type"],
    },
}

# Tool sets per agent — the SQL Agent never sees choose_chart_spec (it isn't
# the one deciding what to chart), and the Narrator Agent never sees the
# data-fetching tools (it has no database access and shouldn't be able to
# query anything itself — it only narrates what it's handed).
_SQL_AGENT_TOOLS = [_QUERY_TOOL, _RULEBOOK_TOOL, _COMPUTE_METRIC_TOOL]
_NARRATOR_AGENT_TOOLS = [_CHART_SPEC_TOOL]

# =============================================================================
# Pydantic models
# =============================================================================
class ChatMessage(BaseModel):
    role: Literal["user", "assistant"]
    content: str

class FeedbackRequest(BaseModel):
    session_id: Optional[str] = None
    user_question: str
    assistant_reply: str
    rating: Literal["up", "down"]
    issue_type: Optional[str] = None
    comment: Optional[str] = None

class ChatRequest(BaseModel):
    message: str
    history: List[ChatMessage] = []
    session_id: Optional[str] = None
    model: str = "sonnet"

class ExportPreviewRequest(BaseModel):
    conversation: List[ChatMessage] = []
    title: str = "Pipeline Intelligence Report"
    export_type: Literal["pdf", "pptx"] = "pdf"
    detail_level: Literal["summary", "detailed"] = "detailed"
    session_id: Optional[str] = None

class ExportDownloadRequest(BaseModel):
    format: Literal["pdf", "pptx", "csv"]
    content: Optional[str] = None
    title: str = "Pipeline Intelligence Report"
    session_id: Optional[str] = None


# =============================================================================
# Claude tool loop
# FIX 1: max_tokens raised to CHAT_MAX_TOKENS (4096) for chat
# =============================================================================
def _extract_text(content_blocks) -> str:
    return "\n".join(
        b.text for b in content_blocks if hasattr(b, "text") and b.text
    ).strip()


# =============================================================================
# Response Builder — Stage 3 of the target architecture.
#
# Takes Claude's already fact-checked reply text, this turn's validated
# query result (if any), and Claude's optional chart spec, and produces the
# final string sent to the browser. Two jobs, kept explicit and separate
# from the tool-loop orchestration in _call_claude():
#   1. Attach a chart above the reply — but only when there's real,
#      validated data behind it, and only when the reply isn't itself a
#      failure/unfinished message (a chart glued under an apology is worse
#      than no chart).
#   2. If Claude never produced any text at all (tool loop exhausted its
#      round budget without a final answer), synthesize a clean,
#      user-facing fallback message — distinguishing a pre-execution rule
#      rejection (not a database error) from an actual DB/timeout failure,
#      so internal validator text never leaks to the end user verbatim.
# =============================================================================
def build_final_response(
    reply: str,
    latest_validated_result: Optional["QueryResult"],
    pending_chart_spec: Optional[dict],
    last_error: Optional[str],
    last_error_is_rule_violation: bool,
    session_id: Optional[str] = None,
) -> str:
    _INVESTIGATION_MARKERS = ("let me check", "let me look", "let me verify",
                               "i'll check", "i'll look", "checking what",
                               "let me fetch", "let me re-run", "let me query")
    reply_looks_unfinished = (
        any(m in reply.lower() for m in _INVESTIGATION_MARKERS)
        and "filters applied" not in reply.lower()
    )
    # Don't attach a chart to a reply that's actually reporting a failure —
    # e.g. the rule-violation / DB-error fallback text, or an apology that
    # slipped through with no real data behind it.
    _FAILURE_MARKERS = ("couldn't complete this", "wasn't able to", "rule violation",
                         "database error", "i wasn't able to build a query",
                         "could you rephrase")
    reply_looks_failed = any(m in reply.lower() for m in _FAILURE_MARKERS)

    # Chart is built HERE, once, from latest_validated_result (this turn's
    # own data — never re-fetched from the shared session store) and
    # pending_chart_spec (Claude's optional choose_chart_spec call, already
    # validated against this exact result's real columns when it was
    # accepted). If Claude never called choose_chart_spec, spec is None and
    # build_chart_html falls back to the original auto-detection — same
    # behavior as before this feature existed. Either way, every number in
    # the output comes from latest_validated_result.rows, computed in
    # charts.py — Claude's spec can only ever pick type/columns/title.
    if (reply and latest_validated_result and latest_validated_result.rows
            and not reply_already_has_chart(reply)
            and not reply_looks_unfinished and not reply_looks_failed):
        chart_html = build_chart_html(
            latest_validated_result.columns,
            latest_validated_result.rows,
            latest_validated_result.filters_applied,
            spec=pending_chart_spec,
        )
        if chart_html:
            reply = chart_html + "\n\n" + reply.lstrip()

    if not reply:
        if last_error and last_error_is_rule_violation:
            # This was rejected by our own SQL guardrails before ever
            # reaching the database — it is NOT a database error, and the
            # internal violation text (fed to the model to self-correct)
            # is not meant for the end user. Log the full detail
            # server-side and show a clean, honest message instead.
            print(f"[RULE-VIOLATION-EXHAUSTED] session={session_id} last_error={last_error}")
            reply = (
                "⚠️ I wasn't able to build a query that satisfies all of our data-quality "
                "rules for this question after several attempts. This is usually a sign the "
                "question needs to be broken down or phrased more specifically (e.g. naming "
                "a single stage or metric). Could you try rephrasing, or check **/debug/db** "
                "if this keeps happening?"
            )
        elif last_error:
            reply = (
                "⚠️ I couldn't complete this query. The last database error was:\n\n"
                f"`{last_error[:400]}`\n\n"
                "Could you rephrase the question, or check **/debug/db** if this persists?"
            )
        else:
            reply = (
                "⚠️ I wasn't able to finish answering this in time — it may need "
                "a more specific or simpler question. Could you try rephrasing it "
                "(e.g. break it into smaller asks)?"
            )
    return reply


def _format_data_envelope(result: "QueryResult") -> str:
    """
    Formats a validated QueryResult into the plain-text block handed to
    the Narrator Agent — same pipe-table style already used for the SQL
    Agent's own tool results (run_clickhouse_query), so both agents see
    data in a consistent, familiar shape. Includes the deterministic
    headline (if one was computed) so the Narrator never has to derive a
    total/average itself from raw rows.
    """
    CHAT_DISPLAY_LIMIT = 100
    header = " | ".join(result.columns)
    lines = [header, "-" * min(len(header), 140)]
    for row in result.rows[:CHAT_DISPLAY_LIMIT]:
        lines.append(" | ".join(str(row.get(c, "")) for c in result.columns))
    if result.total_rows > CHAT_DISPLAY_LIMIT:
        lines.append(f"\n(Showing {CHAT_DISPLAY_LIMIT} of {result.total_rows} rows.)")
    if result.filters_applied:
        lines.append(f"\nFilters applied: {result.filters_applied}")
    headline = _compute_deterministic_headline(result)
    if headline:
        lines.append(f"\n[VERIFIED TOTALS — use these exact figures, do not recompute]: {headline}")
    return "\n".join(lines)


def _dispatch_sql_tool(tool_block, session_id: Optional[str], original_user_message: str,
                        latest_validated_result: Optional["QueryResult"]):
    """
    Handles one tool_use block for the SQL Agent (lookup_business_rule,
    compute_verified_metric, or a query_clickhouse SQL execution — the
    SQL Agent no longer has choose_chart_spec at all, so that branch is
    gone entirely, not just unreachable).

    Returns (content: str, is_error: bool, updated_latest_validated_result,
    last_error: Optional[str], last_error_is_rule_violation: bool).
    The last two are only ever non-None/True when this specific call was
    the one that failed — the caller is responsible for tracking the most
    recent one across the whole loop.
    """
    if tool_block.name == "lookup_business_rule":
        topic = tool_block.input.get("topic", "")
        return get_rulebook_entry(topic), False, latest_validated_result, None, False

    if tool_block.name == "compute_verified_metric":
        metric_name = tool_block.input.get("metric_name", "")
        return compute_verified_metric(metric_name, session_id), False, latest_validated_result, None, False

    sql = tool_block.input.get("sql", "")

    # --- pre-execution rule gate -------------------------------
    sql_violations = validate_sql_against_rules(sql, original_user_message)
    _log_rule_audit(session_id, sql, sql_violations, "pre_execute", original_user_message)

    if sql_violations:
        cohort_violations = [v for v in sql_violations if "cohort" in v.lower() or "8b" in v.lower()]

        if cohort_violations:
            query_result = (
                "RULE VIOLATION — cohort funnel SQL rejected (§8b). NOT executed.\n\n"
                "Violations:\n"
                + "\n".join(f"- {v}" for v in sql_violations)
                + """

⚠️ REQUIRED: Rewrite from scratch using this exact skeleton:

WITH cohort AS (
  SELECT deal_id, deal_stage, amount
  FROM hs_analytics.deals FINAL
  WHERE became_<N>_deal_date IS NOT NULL
    AND deal_stage NOT IN (
        '1% - Prospect', '<stages before N%>'
    )
    AND pipeline = 'default'
    AND CASE WHEN deal_type IS NULL THEN 'Not Assigned' ELSE deal_type END
        NOT IN ('Partner-Led SMB')
    AND toInt64(deal_id) IN (
        SELECT DISTINCT toInt64(deal_id_hs) FROM kore_ai_hubspot.gs_deal_ids_hs
    )
    AND toDate(LEFT(coalesce(became_<N>_deal_date, '1900-01-01'), 10))
        BETWEEN '<start_date>' AND '<end_date>'
)
SELECT
    deal_stage,
    countDistinct(deal_id)       AS deal_count,
    round(SUM(amount) / 1e6, 1) AS pipeline_m
FROM cohort
GROUP BY deal_stage
ORDER BY deal_count DESC

Fill in <N> and the stage exclusion list from §8b, then resubmit.
"""
            )
        else:
            query_result = (
                "RULE VIOLATION — query rejected, NOT executed against the database:\n"
                + "\n".join(f"- {v}" for v in sql_violations)
                + "\n\nRewrite the SQL to satisfy these rules, then call the tool again."
            )
        return query_result, True, latest_validated_result, query_result, True

    query_result = run_clickhouse_query(sql, session_id=session_id)
    is_error = any(query_result.startswith(p) for p in [
        "DATABASE CONNECTION FAILED", "ERROR:", "DATABASE ERROR:"
    ])

    if is_error:
        return query_result, True, latest_validated_result, query_result, False

    # --- post-execution result-level rule check ------------
    parsed_rows = _parse_rows_for_validation(query_result, session_id)
    result_violations = validate_result_against_rules(parsed_rows, original_user_message, sql)
    _log_rule_audit(session_id, sql, result_violations, "post_execute", original_user_message)
    if result_violations:
        query_result += (
            "\n\n⚠️ RESULT RULE VIOLATION:\n"
            + "\n".join(f"- {v}" for v in result_violations)
            + "\nThis result is likely wrong — re-derive using the single-CTE "
              "cohort pattern from §8b and call the tool again."
        )
        return query_result, True, latest_validated_result, query_result, False

    # --- deterministic headline injection ------
    this_result = _get_result(f"{session_id}:latest") if session_id else None
    headline = _compute_deterministic_headline(this_result)
    if headline:
        query_result += (
            f"\n\n[VERIFIED TOTALS — use these exact figures, "
            f"do not recompute]: {headline}"
        )
    if this_result and this_result.rows:
        latest_validated_result = this_result

    return query_result, False, latest_validated_result, None, False


def _run_sql_agent(safe_messages: list, original_user_message: str, session_id: Optional[str],
                    selected_model: str, max_tokens: int):
    """
    Runs the SQL Agent's tool loop to completion: drafts SQL, gets it
    checked by the deterministic rule engine (pre- and post-execution),
    retries on violation, up to MAX_ROUNDS. Never writes the user-facing
    answer — that's the Narrator Agent's job. If the loop ends with no
    validated result and no tool call was ever needed at all (a greeting,
    an export request, a capability question), sql_agent_text carries
    that direct reply so the caller can use it as-is without invoking the
    Narrator Agent for something that was never about data in the first
    place.

    Returns: (sql_agent_text, latest_validated_result, last_error,
    last_error_is_rule_violation, updated_safe_messages)
    """
    response = _ai_client.messages.create(
        model=selected_model,
        system=_cached_sql_agent_prompt(),
        messages=safe_messages,
        tools=_SQL_AGENT_TOOLS,
        temperature=0,
        max_tokens=max_tokens,
    )

    latest_validated_result: Optional["QueryResult"] = None
    last_error = None
    last_error_is_rule_violation = False
    MAX_ROUNDS = 8

    for round_num in range(MAX_ROUNDS):
        if response.stop_reason != "tool_use":
            break

        tool_blocks = [b for b in response.content if b.type == "tool_use"]
        if not tool_blocks:
            break

        tool_result_blocks = []
        for tool_block in tool_blocks:
            content, is_error, latest_validated_result, err, err_is_violation = _dispatch_sql_tool(
                tool_block, session_id, original_user_message, latest_validated_result
            )
            if err is not None:
                last_error = err
                last_error_is_rule_violation = err_is_violation
            tool_result_blocks.append({
                "type": "tool_result", "tool_use_id": tool_block.id,
                "content": content, "is_error": is_error,
            })

        safe_messages = safe_messages + [
            {"role": "assistant", "content": response.content},
            {"role": "user", "content": tool_result_blocks},
        ]
        is_last_round = (round_num == MAX_ROUNDS - 1)
        response = _ai_client.messages.create(
            model=selected_model,
            system=_cached_sql_agent_prompt(),
            messages=safe_messages,
            tools=[] if is_last_round else _SQL_AGENT_TOOLS,
            temperature=0,
            max_tokens=max_tokens,
        )

    sql_agent_text = _extract_text(response.content)

    # ── INCOMPLETE-TURN GUARD ────────────────────────────────────────
    # If the final round was cut off by max_tokens, the SQL Agent may just
    # be narrating its NEXT planned step without having executed it. Force
    # one continuation so the investigation actually finishes.
    if response.stop_reason == "max_tokens":
        continue_messages = safe_messages + [
            {"role": "assistant", "content": sql_agent_text},
            {"role": "user", "content": (
                "Continue. You were cut off mid-response. If you said you "
                "were going to check or fetch something, actually call the "
                "appropriate tool now and complete the investigation. Do "
                "not repeat or re-explain your plan — execute it."
            )},
        ]
        try:
            continue_response = _ai_client.messages.create(
                model=selected_model,
                system=_cached_sql_agent_prompt(),
                messages=continue_messages,
                tools=_SQL_AGENT_TOOLS,
                temperature=0,
                max_tokens=max_tokens,
            )
            for _ in range(3):
                if continue_response.stop_reason != "tool_use":
                    break
                tool_blocks = [b for b in continue_response.content if b.type == "tool_use"]
                if not tool_blocks:
                    break
                tool_result_blocks = []
                for tb in tool_blocks:
                    content, is_error, latest_validated_result, err, err_is_violation = _dispatch_sql_tool(
                        tb, session_id, original_user_message, latest_validated_result
                    )
                    if err is not None:
                        last_error = err
                        last_error_is_rule_violation = err_is_violation
                    tool_result_blocks.append({
                        "type": "tool_result", "tool_use_id": tb.id,
                        "content": content, "is_error": is_error,
                    })
                continue_messages = continue_messages + [
                    {"role": "assistant", "content": continue_response.content},
                    {"role": "user", "content": tool_result_blocks},
                ]
                continue_response = _ai_client.messages.create(
                    model=selected_model, system=_cached_sql_agent_prompt(),
                    messages=continue_messages,
                    tools=_SQL_AGENT_TOOLS,
                    temperature=0, max_tokens=max_tokens,
                )
            completed = _extract_text(continue_response.content)
            if completed:
                sql_agent_text = completed
                safe_messages = continue_messages
        except Exception as e:
            print(f"⚠️ SQL Agent continuation after max_tokens cutoff failed: {e}")

    return sql_agent_text, latest_validated_result, last_error, last_error_is_rule_violation, safe_messages


def _run_narrator_agent(original_user_message: str, latest_validated_result: "QueryResult",
                         session_id: Optional[str], selected_model: str, max_tokens: int):
    """
    Runs the Narrator Agent on ALREADY-VERIFIED data only. Gets a fresh,
    minimal message list — not the SQL Agent's tool-calling history — just
    the question and the validated rows. No database access, no SQL
    rulebook in its prompt; its only tool is choose_chart_spec. This is
    the actual multi-agent boundary: everything upstream of this point is
    the SQL Agent's concern, everything from here on is the Narrator's.

    Returns: (reply, pending_chart_spec, narrator_messages) — the last is
    needed so the fact-binding verifier's regeneration call, if triggered,
    continues from the same context this agent actually used.
    """
    envelope = _format_data_envelope(latest_validated_result) if latest_validated_result else ""
    narrator_messages = [{
        "role": "user",
        "content": (
            f"Question: {original_user_message}\n\n"
            f"Verified data (already checked against every business rule — "
            f"use only these numbers):\n{envelope}"
        ),
    }]

    response = _ai_client.messages.create(
        model=selected_model,
        system=_cached_narrator_agent_prompt(),
        messages=narrator_messages,
        tools=_NARRATOR_AGENT_TOOLS,
        temperature=0,
        max_tokens=max_tokens,
    )

    pending_chart_spec: Optional[dict] = None

    # The Narrator only ever has one tool available, so this loop is
    # intentionally short — a couple of rounds is enough for "call
    # choose_chart_spec, then write the answer".
    for _ in range(3):
        if response.stop_reason != "tool_use":
            break
        tool_blocks = [b for b in response.content if b.type == "tool_use"]
        if not tool_blocks:
            break
        tool_result_blocks = []
        for tb in tool_blocks:
            clean_spec, spec_error = _validate_chart_spec(tb.input, latest_validated_result.columns)
            if spec_error:
                tool_result_blocks.append({
                    "type": "tool_result", "tool_use_id": tb.id,
                    "content": f"Chart spec rejected: {spec_error}. Real columns available: {latest_validated_result.columns}",
                    "is_error": True,
                })
                continue
            pending_chart_spec = clean_spec
            tool_result_blocks.append({
                "type": "tool_result", "tool_use_id": tb.id,
                "content": "Chart spec accepted.", "is_error": False,
            })
        narrator_messages = narrator_messages + [
            {"role": "assistant", "content": response.content},
            {"role": "user", "content": tool_result_blocks},
        ]
        response = _ai_client.messages.create(
            model=selected_model,
            system=_cached_narrator_agent_prompt(),
            messages=narrator_messages,
            tools=[],  # chart spec is a one-shot choice, not re-offered
            temperature=0,
            max_tokens=max_tokens,
        )

    reply = _extract_text(response.content)
    return reply, pending_chart_spec, narrator_messages


def _call_claude(messages: list, max_tokens: int = CHAT_MAX_TOKENS,
                 session_id: Optional[str] = None, model: str = "sonnet") -> str:
    """
    Orchestrator: runs the SQL Agent, and — only if it actually produced
    validated data — hands off to the Narrator Agent to write the answer.
    If the turn never needed data (a greeting, an export request, a
    capabilities question), the SQL Agent's own direct reply is used
    as-is and the Narrator Agent is never invoked at all.
    """
    selected_model = ALLOWED_MODELS.get(model, ALLOWED_MODELS["sonnet"])

    original_user_message = next(
        (m.get("content", "") for m in reversed(messages)
         if m.get("role") == "user" and isinstance(m.get("content", ""), str)),
        ""
    )

    safe_messages = []
    for m in messages:
        content = m.get("content", "")
        if isinstance(content, list):
            text_parts = [
                b.get("text", "") if isinstance(b, dict) else (b.text if hasattr(b, "text") else "")
                for b in content if (isinstance(b, dict) and b.get("type") == "text")
                   or (hasattr(b, "type") and b.type == "text")
            ]
            text = "\n".join(t for t in text_parts if t).strip()
            if text:
                safe_messages.append({"role": m["role"], "content": text})
        else:
            safe_messages.append({"role": m["role"], "content": content})

    sql_agent_text, latest_validated_result, last_error, last_error_is_rule_violation, safe_messages = \
        _run_sql_agent(safe_messages, original_user_message, session_id, selected_model, max_tokens)

    pending_chart_spec: Optional[dict] = None

    if latest_validated_result and latest_validated_result.rows:
        # There's real data — hand off to the Narrator Agent to write the
        # actual answer. sql_agent_text (whatever the SQL Agent said while
        # fetching data, e.g. "let me check that") is deliberately discarded.
        reply, pending_chart_spec, narrator_messages = _run_narrator_agent(
            original_user_message, latest_validated_result, session_id, selected_model, max_tokens
        )
        fact_check_messages = narrator_messages
        fact_check_system = _cached_narrator_agent_prompt()
    else:
        # No data was ever fetched — either this was a no-data turn
        # (greeting, export marker, capability question) that the SQL
        # Agent already answered directly, or SQL Agent exhausted its
        # retries. Either way, there's nothing for the Narrator Agent to
        # narrate.
        reply = sql_agent_text
        fact_check_messages = safe_messages
        fact_check_system = _cached_sql_agent_prompt()

    # ── FACT-BINDING VERIFIER ───────────────────────────────────────
    if reply and session_id:
        stored = _get_result(f"{session_id}:latest") or _get_result(session_id)
        if stored and stored.rows:
            violations = validate_summary_against_facts(reply, stored.rows)
            if violations:
                _log_rule_audit(session_id, "summary_check", violations, "post_summary", original_user_message)
                print(f"⚠️ Unverified numbers in summary: {violations}")
                retry_messages = fact_check_messages + [
                    {"role": "assistant", "content": reply},
                    {"role": "user", "content": (
                        "Rewrite your last response as a clean, standalone answer "
                        "using ONLY numbers that literally appear in the returned "
                        "rows, the [VERIFIED TOTALS] block, or a "
                        "compute_verified_metric result. IMPORTANT: This rewrite is "
                        "the ONLY version the user will ever see — they have not "
                        "seen your previous draft. Do NOT reference 'your previous "
                        "response', do NOT apologize, do NOT say 'you're right' — "
                        "just give the answer directly, as if this is the first "
                        "and only time you're answering."
                    )},
                ]
                try:
                    retry_response = _ai_client.messages.create(
                        model=selected_model,
                        system=fact_check_system,
                        messages=retry_messages,
                        tools=[],
                        temperature=0,
                        max_tokens=max_tokens,
                    )
                    reply = _extract_text(retry_response.content) or reply
                except Exception as e:
                    print(f"⚠️ Fact-binding regeneration failed: {e}")
                    # keep original reply rather than losing the response entirely

    # ── RESPONSE BUILDER ──────────────────────────────────────────────
    return build_final_response(
        reply, latest_validated_result, pending_chart_spec,
        last_error, last_error_is_rule_violation, session_id,
    )

# =============================================================================
# Routes — chat
# =============================================================================
@app.get("/", response_class=HTMLResponse)
def root():
    with open("chat.html", "r") as f:
        return HTMLResponse(content=f.read())

@app.get("/logo.png")
def serve_logo():
    return FileResponse("logo.png", media_type="image/png")


@app.get("/debug/db")
def debug_db():
    base_url = _base_url()
    token    = _token()
    config = {
        "CLICKHOUSE_API_URL":   base_url or "❌ NOT SET",
        "CLICKHOUSE_API_TOKEN": f"✅ set ({len(token)} chars)" if token else "❌ NOT SET",
    }
    if not base_url or not token:
        return {"status": "MISCONFIGURED", "config": config}

    tests = {}
    try:
        r = httpx.get(base_url, timeout=10)
        tests["GET /"] = {"status": r.status_code}
    except Exception as e:
        tests["GET /"] = {"error": str(e)}

    ping = run_clickhouse_query("SELECT 1 AS ping")
    query_ok = not any(ping.startswith(p) for p in ["DATABASE CONNECTION FAILED", "ERROR:", "DATABASE ERROR:"])
    tests["SELECT 1"] = {"ok": query_ok, "result": ping[:100]}

    return {
        "status":            "OK" if query_ok else "FAILED",
        "config":            config,
        "discovered_tables": list(_LIVE_SCHEMA.keys()),
        "active_sessions":   len(_SESSION_STORE),
        "tests":             tests,
    }


@app.get("/debug/metrics")
def debug_metrics(_admin: None = Depends(_require_admin)):
    """
    Success rate and failure-pattern breakdown over the current in-memory
    audit window (last 200 rule-check events). Previously this data was
    logged but never aggregated anywhere — this is the first place it's
    actually queryable as a number instead of a scrollback of print()
    lines. NOTE: resets on server restart along with the rest of
    _RULE_AUDIT_LOG — see the recommendation to move this to external
    storage if that resolution becomes a real gap in practice.

    Developer-only: requires the X-Admin-Token header.
    """
    return get_audit_metrics()


@app.post("/feedback")
def submit_feedback(req: FeedbackRequest):
    """
    Records a thumbs up/down from the chat UI. This is Explicit Learning's
    raw material (a human directly saying whether an answer was good) —
    kept as a separate signal from the rule-violation audit log (which
    only tells you whether DIUD's OWN rules were satisfied, not whether a
    person actually found the answer useful or correct). Deliberately
    does not change any behavior itself; see /debug/feedback.
    """
    _log_feedback(req.session_id, req.user_question, req.assistant_reply,
                  req.rating, req.issue_type, req.comment)
    return {"status": "recorded"}


@app.get("/debug/feedback")
def debug_feedback(_admin: None = Depends(_require_admin)):
    """
    Thumbs up/down rate, broken down by question pattern, plus the most
    recent free-text comments. Sorted worst-first (by pattern) so the
    question type users are least happy with surfaces at the top rather
    than needing to be found by eye.

    Developer-only: requires the X-Admin-Token header. (POST /feedback,
    where users actually submit a rating, stays open to everyone —
    only VIEWING the aggregated results and comments is gated.)
    """
    return get_feedback_metrics()


@app.get("/debug/alerts")
def debug_alerts(_admin: None = Depends(_require_admin)):
    """
    The actual answer to "does DIUD stop repeating its mistakes": not
    automatically, but this is where a repeated mistake becomes
    impossible to miss instead of requiring someone to notice it by eye
    across two separate endpoints. Cross-references /debug/metrics and
    /debug/feedback, flags anything crossing a repeat-failure threshold,
    and — this is the important part — never acts on any of it itself.
    Sorted worst-first (critical > high > medium). Developer-only.
    """
    return get_flagged_issues()


@app.get("/debug/proposals")
def list_proposals(_admin: None = Depends(_require_admin)):
    """All pending proposals, plus the last 50 resolved (approved/rejected). Developer-only."""
    return {"pending": _PENDING_PROPOSALS, "history": _PROPOSAL_HISTORY[:50]}


@app.post("/debug/proposals/draft")
def draft_proposals(_admin: None = Depends(_require_admin)):
    """
    Drafts a fix proposal for each currently-flagged issue that doesn't
    already have one pending or previously resolved. Each draft costs one
    Claude call — this is an explicit action a person triggers, never run
    automatically by /debug/alerts itself, so checking alerts never
    silently spends API calls. Developer-only.
    """
    flags = get_flagged_issues()["flags"]
    already_seen = {p["flag"]["subject"] for p in _PENDING_PROPOSALS + _PROPOSAL_HISTORY}
    drafted = []
    for flag in flags:
        if flag["subject"] in already_seen:
            continue
        try:
            proposal = _draft_fix_proposal(flag)
            _PENDING_PROPOSALS.append(proposal)
            drafted.append(proposal)
        except Exception as e:
            print(f"⚠️ Failed to draft proposal for {flag['subject']}: {e}")
    return {"drafted": len(drafted), "proposals": drafted}


@app.post("/debug/proposals/{proposal_id}/approve")
def approve_proposal(proposal_id: str, _admin: None = Depends(_require_admin)):
    """
    The ONLY endpoint in the entire system that changes a live agent
    prompt as a result of the audit/feedback/flag pipeline — and it only
    ever does so for the exact proposal a human just explicitly approved
    by ID. Appends the proposed clarification to the target agent's
    patch list and rebuilds that agent's cached prompt immediately.
    Developer-only — the single most consequential endpoint here.
    """
    global _SQL_AGENT_PROMPT, _NARRATOR_AGENT_PROMPT
    proposal = next((p for p in _PENDING_PROPOSALS if p["id"] == proposal_id), None)
    if not proposal:
        raise HTTPException(status_code=404, detail="Proposal not found or already resolved.")
    if not proposal["proposed_addition"]:
        raise HTTPException(status_code=400, detail="This proposal has no safe fix to apply.")

    if proposal["target_prompt"] == "narrator_agent":
        _NARRATOR_AGENT_PATCHES.append(proposal["proposed_addition"])
        _NARRATOR_AGENT_PROMPT = _build_narrator_agent_prompt()
    else:
        _SQL_AGENT_PATCHES.append(proposal["proposed_addition"])
        _SQL_AGENT_PROMPT = _build_sql_agent_prompt()

    proposal["status"] = "approved"
    proposal["resolved_at"] = datetime.utcnow().isoformat()
    _PENDING_PROPOSALS.remove(proposal)
    _PROPOSAL_HISTORY.insert(0, proposal)
    print(f"[PROPOSAL-APPROVED] {proposal_id} applied to {proposal['target_prompt']}")
    return {"status": "approved", "applied_to": proposal["target_prompt"]}


@app.post("/debug/proposals/{proposal_id}/reject")
def reject_proposal(proposal_id: str, _admin: None = Depends(_require_admin)):
    """Discards a proposal without applying anything. No prompt is touched."""
    proposal = next((p for p in _PENDING_PROPOSALS if p["id"] == proposal_id), None)
    if not proposal:
        raise HTTPException(status_code=404, detail="Proposal not found or already resolved.")
    proposal["status"] = "rejected"
    proposal["resolved_at"] = datetime.utcnow().isoformat()
    _PENDING_PROPOSALS.remove(proposal)
    _PROPOSAL_HISTORY.insert(0, proposal)
    return {"status": "rejected"}


@app.post("/refresh-schema")
def refresh_schema():
    global _SQL_AGENT_PROMPT, _NARRATOR_AGENT_PROMPT
    schema = discover_schema()
    _SQL_AGENT_PROMPT = _build_sql_agent_prompt()
    _NARRATOR_AGENT_PROMPT = _build_narrator_agent_prompt()
    return {"status": "refreshed", "tables": list(_LIVE_SCHEMA.keys())}


@app.post("/chat")
def chat(payload: ChatRequest):
    messages = [{"role": m.role, "content": m.content} for m in payload.history]
    messages.append({"role": "user", "content": payload.message})
    print(f"💬 [chat] session={payload.session_id} msg={payload.message[:80]}")

    try:
        # FIX 1: pass CHAT_MAX_TOKENS explicitly
        reply = _call_claude(messages, max_tokens=CHAT_MAX_TOKENS,
                             session_id=payload.session_id, model=payload.model)
    except Exception as exc:
        traceback.print_exc()
        raise HTTPException(status_code=502, detail=f"Claude error: {exc}")

    has_dataset = payload.session_id is not None and payload.session_id in _SESSION_STORE
    stored = _SESSION_STORE.get(payload.session_id) if payload.session_id else None

    return {
        "reply":        reply,
        "has_dataset":  has_dataset,
        "dataset_rows": stored.total_rows if stored else 0,
        "export_intent": "__EXPORT_INTENT__" in reply,
    }


# =============================================================================
# Retry
# =============================================================================
class RetryRequest(BaseModel):
    history:    List[ChatMessage] = []
    session_id: Optional[str] = None
    model:      str = "sonnet"


@app.post("/chat/retry")
def chat_retry(payload: RetryRequest):
    if not payload.history:
        raise HTTPException(status_code=400, detail="history must not be empty for retry.")

    clean_history = list(payload.history)
    while clean_history and clean_history[-1].role == "assistant":
        clean_history.pop()

    if not clean_history or clean_history[-1].role != "user":
        raise HTTPException(
            status_code=400,
            detail="No user message found to retry. history must end with a user turn."
        )

    last_user_msg = clean_history[-1].content
    prior_history = clean_history[:-1]

    print(f"🔄 [retry] session={payload.session_id} retrying: {last_user_msg[:80]}")

    messages = [{"role": m.role, "content": m.content} for m in prior_history]
    messages.append({"role": "user", "content": last_user_msg})

    try:
        reply = _call_claude(messages, max_tokens=CHAT_MAX_TOKENS, session_id=payload.session_id)
    except Exception as exc:
        traceback.print_exc()
        raise HTTPException(status_code=502, detail=f"Claude error on retry: {exc}")

    has_dataset = payload.session_id is not None and payload.session_id in _SESSION_STORE
    stored      = _SESSION_STORE.get(payload.session_id) if payload.session_id else None

    return {
        "reply":         reply,
        "has_dataset":   has_dataset,
        "dataset_rows":  stored.total_rows if stored else 0,
        "export_intent": "__EXPORT_INTENT__" in reply,
        "retried":       True,
    }


# =============================================================================
# Session info
# =============================================================================
@app.get("/session/{session_id}/dataset-info")
def session_dataset_info(session_id: str):
    # Prefer the most recent result over the largest-wins export cache —
    # "what can I export" should reflect what the user just saw.
    result = _get_result(f"{session_id}:latest") or _get_result(session_id)
    if not result:
        return {"has_dataset": False}
    return {
        "has_dataset":     True,
        "total_rows":      result.total_rows,
        "columns":         result.columns,
        "captured_at":     result.captured_at,
        "filters_applied": result.filters_applied,
        "sql_preview":     result.sql[:300],
    }


# =============================================================================
# Export preview
# FIX 1: uses EXPORT_MAX_TOKENS (6144) for long reports
# =============================================================================
@app.post("/export/preview")
async def export_preview(req: ExportPreviewRequest):
    if not req.conversation:
        raise HTTPException(status_code=400, detail="No conversation to export.")

    print(f"📄 [export/preview] session={req.session_id} type={req.export_type}")

    # "Export this" should reflect the most recent result, not whichever
    # past query in this session happened to return the most rows.
    stored = (
        (_get_result(f"{req.session_id}:latest") or _get_result(req.session_id))
        if req.session_id else None
    )

    try:
        ai_content = _generate_export_content(
            conversation   = req.conversation,
            title          = req.title,
            export_type    = req.export_type,
            detail_level   = req.detail_level,
            stored_dataset = stored,
        )
    except Exception as exc:
        traceback.print_exc()
        raise HTTPException(status_code=502, detail=f"Content generation error: {exc}")

    return {
        "content":       ai_content,
        "title":         req.title,
        "export_type":   req.export_type,
        "word_count":    len(ai_content.split()),
        "generated_at":  date.today().isoformat(),
        "total_rows":    stored.total_rows if stored else 0,
        "filters":       stored.filters_applied if stored else "",
    }


# =============================================================================
# Export download
# =============================================================================
@app.post("/export/download")
async def export_download(req: ExportDownloadRequest):
    print(f"⬇️  [export/download] format={req.format} session={req.session_id}")

    if req.format in ("csv", "xlsx"):
        # Same lookup pattern as export_preview / session_dataset_info.
        stored = (
            (_get_result(f"{req.session_id}:latest") or _get_result(req.session_id))
            if req.session_id else None
        )
        if not stored:
            raise HTTPException(
                status_code=404,
                detail=(
                    "No query result found for this session. "
                    "Ask a deal-list question first, then export."
                ),
            )

        if req.format == "csv":
            file_bytes = _build_csv(stored)
            media_type = "text/csv"
            ext = "csv"
        else:
            file_bytes = _build_xlsx(stored, req.title)
            media_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            ext = "xlsx"

        return StreamingResponse(
            io.BytesIO(file_bytes),
            media_type=media_type,
            headers={
                "Content-Disposition": f'attachment; filename="{_safe_filename(req.title)}.{ext}"',
                "X-Total-Rows": str(stored.total_rows),
            },
        )

    if not req.content:
        raise HTTPException(status_code=400, detail="content is required for PDF/PPTX export.")

    try:
        if req.format == "pdf":
            file_bytes = _build_pdf(req.title, req.content)
            media_type = "application/pdf"
            ext        = "pdf"
        else:
            file_bytes = _build_pptx(req.title, req.content)
            media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            ext        = "pptx"

        return StreamingResponse(
            io.BytesIO(file_bytes),
            media_type=media_type,
            headers={"Content-Disposition": f'attachment; filename="{_safe_filename(req.title)}.{ext}"'},
        )
    except Exception as exc:
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"File generation error: {exc}")


# =============================================================================
# Helpers
# =============================================================================
def _safe_filename(title: str) -> str:
    return re.sub(r'[^\w\-]', '_', title)[:60]


def _strip_md(t: str) -> str:
    t = re.sub(r'\*{1,3}(.*?)\*{1,3}', r'\1', t)
    t = re.sub(r'#{1,6}\s*', '', t)
    t = re.sub(r'^[\-\*•]\s*', '', t, flags=re.M)
    t = re.sub(r'`(.*?)`', r'\1', t)
    t = t.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
    return t.strip()


# =============================================================================
# CSV builder
# =============================================================================
def _build_csv(stored: QueryResult) -> bytes:
    buf = io.StringIO()
    buf.write(f"# Title: {stored.sql[:80]}\n")
    buf.write(f"# Generated: {date.today().isoformat()}\n")
    buf.write(f"# Total Records: {stored.total_rows}\n")
    buf.write(f"# Filters: {stored.filters_applied}\n")
    buf.write(f"# Captured at: {stored.captured_at}\n")
    buf.write("#\n")

    writer = csv.DictWriter(
        buf,
        fieldnames     = stored.columns,
        extrasaction   = "ignore",
        lineterminator = "\n",
    )
    writer.writeheader()
    for row in stored.rows:
        writer.writerow(row)

    return buf.getvalue().encode("utf-8")


_CURRENCY_COLUMN_HINTS = (
    "amount", "value", "revenue", "target", "pipeline", "arr", "acv",
    "price", "cost", "spend", "budget", "quota",
)


def _looks_like_currency_column(col_name: str) -> bool:
    lower = col_name.lower()
    return any(hint in lower for hint in _CURRENCY_COLUMN_HINTS)


def _build_xlsx(stored: QueryResult, title: str) -> bytes:
    """
    Styled Excel export — bold white-on-navy header row, currency-formatted
    amount/value/target-style columns, auto-sized columns, frozen header
    row for scrolling through long deal lists. Mirrors the CSV export's
    data (same stored.columns / stored.rows, so the two formats can never
    silently disagree on row count) with real spreadsheet formatting on
    top, rather than a second, separately-maintained code path.

    IMPORTANT — completeness guarantee: this always writes every row in
    stored.rows, which is the full result ClickHouse returned (never the
    100-row CHAT_DISPLAY_LIMIT used for the on-screen preview). If that
    ever doesn't match stored.total_rows, something upstream truncated
    the result silently — we'd rather surface that loudly here than ship
    a spreadsheet that quietly has fewer rows than it claims.
    """
    wb = Workbook()
    ws = wb.active
    ws.title = "Deals"[:31]

    header_font = Font(name="Arial", size=11, bold=True, color="FFFFFF")
    header_fill = PatternFill(start_color="1E3A5F", end_color="1E3A5F", fill_type="solid")
    body_font = Font(name="Arial", size=10.5)
    header_align = Alignment(horizontal="center", vertical="center")

    # Row 1: header
    for col_idx, col_name in enumerate(stored.columns, start=1):
        cell = ws.cell(row=1, column=col_idx, value=col_name)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = header_align
    ws.freeze_panes = "A2"

    currency_cols = {c for c in stored.columns if _looks_like_currency_column(c)}

    # Data rows
    for row_idx, row in enumerate(stored.rows, start=2):
        for col_idx, col_name in enumerate(stored.columns, start=1):
            raw_val = row.get(col_name, "")
            cell = ws.cell(row=row_idx, column=col_idx)
            cell.font = body_font
            if col_name in currency_cols:
                try:
                    cell.value = float(raw_val)
                    cell.number_format = '$#,##0'
                    continue
                except (TypeError, ValueError):
                    pass  # fall through — not actually numeric for this row, write as-is
            cell.value = raw_val

    # Auto-width: based on the longer of header text or a sample of data
    for col_idx, col_name in enumerate(stored.columns, start=1):
        sample_lens = [len(str(r.get(col_name, ""))) for r in stored.rows[:200]]
        width = max([len(col_name)] + sample_lens) + 3
        ws.column_dimensions[get_column_letter(col_idx)].width = min(max(width, 10), 60)

    # A compact metadata strip below the table — same facts as the CSV's
    # comment header, kept out of the data rows themselves so it never
    # gets mistaken for a record or breaks a downstream pivot/filter.
    meta_row = len(stored.rows) + 3
    completeness_note = (
        f"{len(stored.rows)} of {stored.total_rows} total records included"
        if len(stored.rows) != stored.total_rows
        else f"All {stored.total_rows} matching records included"
    )
    for i, line in enumerate([
        f"Generated: {date.today().isoformat()}",
        f"Filters: {stored.filters_applied or 'none'}",
        completeness_note,
    ]):
        c = ws.cell(row=meta_row + i, column=1, value=line)
        c.font = Font(name="Arial", size=9, italic=True, color="666666")

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# =============================================================================
# Export content generation
# FIX 1: EXPORT_MAX_TOKENS used here for long document generation
# =============================================================================

def _render_html_to_png(html_content: str, width: int = 800) -> Optional[bytes]:
    try:
        from playwright.sync_api import sync_playwright
    except ImportError:
        print("⚠️  Playwright not installed.")
        return None

    full_html = f"""<!DOCTYPE html><html><head>
<meta charset="UTF-8">
<link href="https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700;800&display=swap" rel="stylesheet">
<style>
  * {{ box-sizing: border-box; margin: 0; padding: 0; }}
  body {{ background: #ffffff; font-family: 'Inter', Arial, sans-serif; }}
</style>
</head><body>{html_content}
</body></html>"""

    try:
        with sync_playwright() as p:
            browser = p.chromium.launch(
                args=[
                    "--no-sandbox",
                    "--disable-setuid-sandbox",
                    "--disable-dev-shm-usage",
                    "--disable-gpu",
                    "--single-process",          # ← critical for Render
                    "--no-zygote",               # ← critical for Render
                ]
            )
            page = browser.new_page(viewport={"width": width, "height": 800})
            page.set_content(full_html, wait_until="networkidle", timeout=20000)
            page.wait_for_timeout(2000)
            height = page.evaluate("document.body.scrollHeight")
            page.set_viewport_size({"width": width, "height": max(height, 100)})
            png_bytes = page.screenshot(full_page=True)
            browser.close()
            return png_bytes
    except Exception as e:
        print(f"⚠️  Playwright screenshot failed: {e}")
        return None


def _extract_html_blocks(text: str):
    """
    Split a message into alternating text/html segments.
    Returns list of ("text"|"html", content) tuples.
    """
    segments = []
    last_end = 0
    for m in re.finditer(r'```html\s*\n(.*?)```', text, flags=re.DOTALL | re.IGNORECASE):
        segments.append(("text", text[last_end:m.start()]))
        segments.append(("html", m.group(1).strip()))
        last_end = m.end()
    segments.append(("text", text[last_end:]))
    return segments


def _generate_export_content(
    conversation:    List[ChatMessage],
    title:           str,
    export_type:     str,
    detail_level:    str = "detailed",
    stored_dataset:  Optional[QueryResult] = None,
) -> str:

    def _clean_text_segment(text: str) -> str:
        text = re.sub(r'```[\w]*\n.*?```', '', text, flags=re.DOTALL)
        text = re.sub(r'^\s*--\s*$', '', text, flags=re.MULTILINE)
        text = re.sub(r'^\s*`+\s*$', '', text, flags=re.MULTILINE)
        text = re.sub(r'\n{3,}', '\n\n', text)
        text = text.replace('■', '').replace('🟦', '')
        return text.strip()

    # ── VERBATIM MODE ────────────────────────────────────────────────────
    if detail_level == "detailed":
        if export_type == "pptx":
            lines = [f"# {title}", f"*Exported: {date.today().strftime('%B %d, %Y')}*", ""]
            for i, m in enumerate(conversation):
                role = "User" if m.role == "user" else "DIUD Agent"
                segments = _extract_html_blocks(m.content)
                lines.append(f"SLIDE: {role} — Turn {i + 1}")
                for seg_type, seg_content in segments:
                    if seg_type == "html":
                        lines.append("- [Interactive chart — see live DIUD interface]")
                    else:
                        cleaned = _clean_text_segment(seg_content)
                        for part in cleaned.split("\n"):
                            part = part.strip()
                            if part:
                                lines.append(f"- {part}")
            return "\n".join(lines)
        else:
            # PDF: encode segments as JSON with sentinel so _build_pdf can embed chart images
            pdf_segments = []
            for i, m in enumerate(conversation):
                role = "User" if m.role == "user" else "DIUD Agent"
                pdf_segments.append({"type": "turn_header", "turn": i + 1, "role": role})
                for seg_type, seg_content in _extract_html_blocks(m.content):
                    if seg_type == "html":
                        pdf_segments.append({"type": "html_chart", "html": seg_content})
                    else:
                        cleaned = _clean_text_segment(seg_content)
                        if cleaned:
                            pdf_segments.append({"type": "text", "content": cleaned})
            return "__VERBATIM_SEGMENTS__" + json.dumps(pdf_segments)

    # ── SUMMARY MODE ─────────────────────────────────────────────────────
    selected_model = ALLOWED_MODELS["sonnet"]
    conv_text = "\n\n".join(
        f"{'USER' if m.role == 'user' else 'DIUD AGENT'}: {m.content}"
        for m in conversation
    )
    format_hint = (
        "Format as a PowerPoint: use SLIDE: <title> for each slide, then bullet points."
        if export_type == "pptx"
        else "Format as a professional PDF report: ## section headers, narrative prose, tables."
    )
    dataset_hint = ""
    if stored_dataset:
        dataset_hint = (
            f"\n\nDATASET CONTEXT: The query returned {stored_dataset.total_rows} records "
            f"with columns: {', '.join(stored_dataset.columns[:12])}. "
            f"Filters: {stored_dataset.filters_applied}."
        )
    prompt = f"""You are preparing a professional {export_type.upper()} summary report.

CONVERSATION:
{conv_text}
{dataset_hint}

TASK: Create a concise executive summary titled "{title}"

{format_hint}

REQUIREMENTS:
- Executive summary at the start with key numbers only
- Logical sections: summary, key metrics, insights, recommendations
- Bold key numbers; clean professional tone
- Today: {date.today().strftime('%B %d, %Y')}
- Generate the COMPLETE document — do not truncate

Generate the summary report now:"""

    response = _ai_client.messages.create(
        model=selected_model,
        system="You are a professional document formatter. Follow the instructions exactly. Never add content not requested.",
        messages=[{"role": "user", "content": prompt}],
        temperature=0,
        max_tokens=EXPORT_MAX_TOKENS,
    )
    ai_text = _extract_text(response.content)

    if stored_dataset and stored_dataset.total_rows > 0:
        table_md  = _rows_to_markdown_table(stored_dataset)
        meta_line = (
            f"**Total records:** {stored_dataset.total_rows:,} | "
            f"**Filters:** {stored_dataset.filters_applied} | "
            f"**Exported:** {date.today().strftime('%B %d, %Y')}"
        )
        full_section = f"\n\n## Data Export ({stored_dataset.total_rows:,} records)\n\n{meta_line}\n\n{table_md}"
        if "[DEAL_TABLE_PLACEHOLDER]" in ai_text:
            ai_text = ai_text.replace("[DEAL_TABLE_PLACEHOLDER]", full_section)
        else:
            ai_text = ai_text.rstrip() + full_section

    return ai_text

def _rows_to_markdown_table(stored: QueryResult) -> str:
    if not stored.rows:
        return "_No data._"

    cols   = stored.columns
    header = "| " + " | ".join(cols) + " |"
    sep    = "| " + " | ".join("---" for _ in cols) + " |"
    lines  = [header, sep]

    for row in stored.rows:
        cells = [str(row.get(c, "")).replace("|", "\\|") for c in cols]
        lines.append("| " + " | ".join(cells) + " |")

    return "\n".join(lines)


# =============================================================================
# PDF Builder
# =============================================================================
_C_NAVY  = colors.HexColor("#0D1B3E")
_C_BLUE  = colors.HexColor("#1565C0")
_C_WHITE = colors.white
_C_BG    = colors.HexColor("#F7F9FC")
_C_TXT   = colors.HexColor("#1E293B")
_C_DIM   = colors.HexColor("#94A3B8")

_SECTION_COLORS = {
    "executive": colors.HexColor("#0D1B3E"),
    "pipeline":  colors.HexColor("#1565C0"),
    "metric":    colors.HexColor("#004D40"),
    "deal":      colors.HexColor("#1565C0"),
    "regional":  colors.HexColor("#BF360C"),
    "win":       colors.HexColor("#B71C1C"),
    "loss":      colors.HexColor("#B71C1C"),
    "recommend": colors.HexColor("#1B5E20"),
    "summary":   colors.HexColor("#0D1B3E"),
    "analysis":  colors.HexColor("#1565C0"),
    "overview":  colors.HexColor("#004D40"),
    "insight":   colors.HexColor("#1B5E20"),
}

PW, PH = A4
_ML = _MR = 0.6 * inch
_MT = 0.45 * inch
_MB = 0.40 * inch
_HDR_H = 44
_FTR_H = 20
_CW    = PW - _ML - _MR


def _pdf_styles():
    return {
        "Cover_Title": ParagraphStyle("Cover_Title", fontSize=26, leading=32,
            textColor=_C_WHITE, fontName="Helvetica-Bold", spaceAfter=8),
        "Cover_Sub": ParagraphStyle("Cover_Sub", fontSize=13, leading=18,
            textColor=colors.HexColor("#B0BEC5"), fontName="Helvetica"),
        "Section_H": ParagraphStyle("Section_H", fontSize=11, leading=15,
            textColor=_C_WHITE, fontName="Helvetica-Bold"),
        "Body": ParagraphStyle("Body", fontSize=9, leading=14, textColor=_C_TXT,
            fontName="Helvetica", spaceAfter=4),
        "Bullet": ParagraphStyle("Bullet", fontSize=9, leading=14, textColor=_C_TXT,
            fontName="Helvetica", leftIndent=12, firstLineIndent=-8, spaceAfter=3),
        "H2": ParagraphStyle("H2", fontSize=11, leading=15, textColor=_C_NAVY,
            fontName="Helvetica-Bold", spaceBefore=10, spaceAfter=4),
        "H3": ParagraphStyle("H3", fontSize=9, leading=13, textColor=_C_BLUE,
            fontName="Helvetica-Bold", spaceBefore=6, spaceAfter=2),
        "TH": ParagraphStyle("TH", fontSize=7, leading=9, textColor=_C_WHITE,
            fontName="Helvetica-Bold"),
        "TD": ParagraphStyle("TD", fontSize=7, leading=9, textColor=_C_TXT,
            fontName="Helvetica"),
    }


def _parse_sections(text: str):
    parts = re.split(r'^##\s+', text, flags=re.MULTILINE)
    return [
        (lines[0].strip(), lines[1].strip() if len(lines) > 1 else "")
        for part in parts if part.strip()
        for lines in [part.strip().split("\n", 1)]
    ]


def _build_pdf(title: str, report_text: str) -> bytes:
    from reportlab.platypus import Image as RLImage
    buf    = io.BytesIO()
    styles = _pdf_styles()

    def _on_page(canvas, doc):
        canvas.saveState()
        canvas.setFillColor(_C_NAVY)
        canvas.rect(0, PH - _HDR_H - _MT, PW, _HDR_H + _MT, fill=1, stroke=0)
        canvas.setFillColor(_C_WHITE)
        canvas.setFont("Helvetica-Bold", 10)
        canvas.drawString(_ML, PH - _MT - 28, title)
        canvas.setFillColor(_C_BG)
        canvas.rect(0, 0, PW, _FTR_H + _MB, fill=1, stroke=0)
        canvas.setFillColor(_C_DIM)
        canvas.setFont("Helvetica", 7)
        canvas.drawCentredString(
            PW / 2, _MB + 5,
            f"DIUD Report  |  AI-Generated  |  CONFIDENTIAL  |  {date.today().strftime('%B %Y')}"
        )
        canvas.drawRightString(PW - _MR, _MB + 5, f"Page {canvas.getPageNumber()}")
        canvas.restoreState()

    frame    = Frame(_ML, _MB + _FTR_H, _CW, PH - _HDR_H - _MT - _MB - _FTR_H, id="main")
    template = PageTemplate(id="main", frames=[frame], onPage=_on_page)
    doc = BaseDocTemplate(buf, pagesize=A4, leftMargin=_ML, rightMargin=_MR,
                          topMargin=_MT + _HDR_H, bottomMargin=_MB + _FTR_H)
    doc.addPageTemplates([template])

    cover_story = [
        Spacer(1, 1.0 * inch),
        Paragraph(title, styles["Cover_Title"]),
        Paragraph(f"Generated {date.today().strftime('%B %d, %Y')}", styles["Cover_Sub"]),
        PageBreak(),
    ]

    # ── VERBATIM SEGMENTS PATH ────────────────────────────────────────────
    if report_text.startswith("__VERBATIM_SEGMENTS__"):
        segments = json.loads(report_text[len("__VERBATIM_SEGMENTS__"):])
        story = cover_story[:]

        for seg in segments:
            if seg["type"] == "turn_header":
                role_label = "User" if seg["role"] == "User" else "DIUD Agent"
                bar_color  = _C_BLUE if seg["role"] == "User" else _C_NAVY
                story.append(Spacer(1, 8))
                story.append(Table(
                    [[Paragraph(f"Turn {seg['turn']} — {role_label}", styles["Section_H"])]],
                    colWidths=[_CW],
                    style=TableStyle([
                        ("BACKGROUND",    (0,0),(-1,-1), bar_color),
                        ("TOPPADDING",    (0,0),(-1,-1), 6),
                        ("BOTTOMPADDING", (0,0),(-1,-1), 6),
                        ("LEFTPADDING",   (0,0),(-1,-1), 10),
                    ])
                ))
                story.append(Spacer(1, 6))

            elif seg["type"] == "html_chart":
                png_bytes = _render_html_to_png(seg["html"], width=760)
                if png_bytes:
                    img = RLImage(io.BytesIO(png_bytes), width=_CW, height=_CW * 0.55)
                    img.hAlign = "CENTER"
                    story.append(img)
                    story.append(Spacer(1, 10))
                else:
                    story.append(Paragraph(
                        "[Chart could not be rendered — view in live DIUD interface]",
                        styles["Body"]
                    ))

            elif seg["type"] == "text":
                lines = seg["content"].split("\n")
                i = 0
                while i < len(lines):
                    line = lines[i].strip()
                    if line.startswith("|") and i + 1 < len(lines) and "---" in lines[i + 1]:
                        table_lines = []
                        while i < len(lines) and lines[i].strip().startswith("|"):
                            table_lines.append(lines[i].strip())
                            i += 1
                        story.append(_md_table_to_rl(table_lines, styles))
                        story.append(Spacer(1, 6))
                        continue
                    if not line:
                        story.append(Spacer(1, 3))
                    elif line.startswith("### "):
                        story.append(Paragraph(_strip_md(line), styles["H3"]))
                    elif line.startswith("## "):
                        story.append(Paragraph(_strip_md(line), styles["H2"]))
                    elif line.startswith(("- ", "* ", "• ")):
                        story.append(Paragraph("• " + _strip_md(line[2:]), styles["Bullet"]))
                    else:
                        story.append(Paragraph(_strip_md(line), styles["Body"]))
                    i += 1

        doc.build(story)
        return buf.getvalue()

    # ── SUMMARY PATH (existing logic, unchanged) ──────────────────────────
    sections = _parse_sections(report_text)
    story = cover_story[:]

    if not sections:
        # strip residual HTML before fallback rendering
        clean_text = re.sub(r'```html.*?```', '[Chart omitted]', report_text, flags=re.DOTALL | re.IGNORECASE)
        clean_text = re.sub(r'```.*?```', '', clean_text, flags=re.DOTALL)
        clean_text = re.sub(r'^\s*--\s*$', '', clean_text, flags=re.MULTILINE)
        for line in clean_text.split("\n"):
            line = line.strip()
            if line:
                story.append(Paragraph(_strip_md(line), styles["Body"]))
    else:
        for sec_title, sec_body in sections:
            color_key = next((k for k in _SECTION_COLORS if k in sec_title.lower()), None)
            bar_color = _SECTION_COLORS.get(color_key, _C_BLUE)
            story.append(Table(
                [[Paragraph(sec_title.upper(), styles["Section_H"])]],
                colWidths=[_CW],
                style=TableStyle([
                    ("BACKGROUND",    (0, 0), (-1, -1), bar_color),
                    ("TOPPADDING",    (0, 0), (-1, -1), 8),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
                    ("LEFTPADDING",   (0, 0), (-1, -1), 10),
                ])
            ))
            story.append(Spacer(1, 6))

            lines = sec_body.split("\n")
            i = 0
            while i < len(lines):
                line = lines[i].strip()
                if line.startswith("|") and i + 1 < len(lines) and "---" in lines[i + 1]:
                    table_lines = []
                    while i < len(lines) and lines[i].strip().startswith("|"):
                        table_lines.append(lines[i].strip())
                        i += 1
                    story.append(_md_table_to_rl(table_lines, styles))
                    story.append(Spacer(1, 6))
                    continue
                if not line:
                    story.append(Spacer(1, 3))
                elif line.startswith("### "):
                    story.append(Paragraph(_strip_md(line), styles["H3"]))
                elif line.startswith("## "):
                    story.append(Paragraph(_strip_md(line), styles["H2"]))
                elif line.startswith(("- ", "* ", "• ")):
                    story.append(Paragraph("• " + _strip_md(line[2:]), styles["Bullet"]))
                else:
                    story.append(Paragraph(_strip_md(line), styles["Body"]))
                i += 1

            story.extend([Spacer(1, 12), PageBreak()])

    doc.build(story)
    return buf.getvalue()


def _md_table_to_rl(table_lines: list, styles: dict):
    data = []
    for idx, line in enumerate(table_lines):
        if "---" in line:
            continue
        cells = [c.strip().replace("\\|", "|") for c in line.strip("|").split("|")]
        if idx == 0:
            row = [Paragraph(_strip_md(c), styles["TH"]) for c in cells]
        else:
            row = [Paragraph(_strip_md(c), styles["TD"]) for c in cells]
        data.append(row)

    if not data:
        return Spacer(1, 1)

    num_cols = max(len(r) for r in data)
    col_w    = _CW / max(num_cols, 1)

    tbl = Table(data, colWidths=[col_w] * num_cols, repeatRows=1)
    tbl.setStyle(TableStyle([
        ("BACKGROUND",    (0, 0), (-1, 0),  _C_NAVY),
        ("TEXTCOLOR",     (0, 0), (-1, 0),  _C_WHITE),
        ("FONTNAME",      (0, 0), (-1, 0),  "Helvetica-Bold"),
        ("FONTSIZE",      (0, 0), (-1, -1), 7),
        ("ROWBACKGROUNDS",(0, 1), (-1, -1), [colors.white, colors.HexColor("#F8FAFC")]),
        ("GRID",          (0, 0), (-1, -1), 0.3, colors.HexColor("#E2E8F0")),
        ("TOPPADDING",    (0, 0), (-1, -1), 3),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
        ("LEFTPADDING",   (0, 0), (-1, -1), 4),
        ("RIGHTPADDING",  (0, 0), (-1, -1), 4),
        ("VALIGN",        (0, 0), (-1, -1), "MIDDLE"),
    ]))
    return tbl


# =============================================================================
# PPTX Builder
# =============================================================================
_C_NAVY_P  = RGBColor(0x0D, 0x1B, 0x3E)
_C_DNAV_P  = RGBColor(0x0A, 0x11, 0x28)
_C_BLUE_P  = RGBColor(0x1E, 0x88, 0xE5)
_C_WHITE_P = RGBColor(0xFF, 0xFF, 0xFF)
_C_LTBG_P  = RGBColor(0xF5, 0xF7, 0xFA)
_C_TXT_P   = RGBColor(0x1A, 0x1A, 0x2E)
_C_DIM_P   = RGBColor(0x88, 0x99, 0xAA)

_SLIDE_ACCENT = {
    "overview":  RGBColor(0x1E, 0x88, 0xE5),
    "pipeline":  RGBColor(0x00, 0x89, 0x7B),
    "deal":      RGBColor(0x1E, 0x88, 0xE5),
    "metric":    RGBColor(0x2E, 0x7D, 0x32),
    "regional":  RGBColor(0xBF, 0x36, 0x0C),
    "win":       RGBColor(0x2E, 0x7D, 0x32),
    "loss":      RGBColor(0xC6, 0x28, 0x28),
    "recommend": RGBColor(0x1B, 0x5E, 0x20),
    "summary":   RGBColor(0x1E, 0x88, 0xE5),
    "analysis":  RGBColor(0x00, 0x89, 0x7B),
    "insight":   RGBColor(0x1B, 0x5E, 0x20),
    "executive": RGBColor(0x0D, 0x1B, 0x3E),
}


def _pptx_bg(slide, color):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def _pptx_rect(slide, l, t, w, h, color):
    shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = color; shp.line.fill.background()
    return shp

def _pptx_txt(slide, text, l, t, w, h, bold=False, size=18, color=None, align=PP_ALIGN.LEFT):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color or _C_TXT_P
    return txb

def _parse_slides(text: str):
    slides, cur_title, cur_bullets = [], None, []
    for line in text.split("\n"):
        line = line.rstrip()
        if line.startswith("SLIDE:"):
            if cur_title is not None:
                slides.append((cur_title, cur_bullets))
            cur_title, cur_bullets = line[6:].strip(), []
        elif line.startswith("- ") and cur_title:
            cur_bullets.append(line[2:].strip())
    if cur_title is not None:
        slides.append((cur_title, cur_bullets))
    return slides

def _build_pptx(title: str, slide_text: str) -> bytes:
    slides_data = _parse_slides(slide_text) or [(title, [slide_text[:400]])]
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    footer_text = f"DIUD  |  AI-Generated  |  CONFIDENTIAL  |  {date.today().strftime('%B %Y')}"
    blank = prs.slide_layouts[6]

    def _footer(s):
        _pptx_rect(s, 0, 7.1, 13.33, 0.4, _C_DNAV_P)
        _pptx_txt(s, footer_text, 0.3, 7.12, 12, 0.35, size=7, color=_C_DIM_P, align=PP_ALIGN.CENTER)

    def _accent(t):
        for k, c in _SLIDE_ACCENT.items():
            if k in t: return c
        return _C_BLUE_P

    cover = prs.slides.add_slide(blank)
    _pptx_bg(cover, _C_NAVY_P)
    _pptx_rect(cover, 0, 3.2, 13.33, 0.06, _C_BLUE_P)
    _pptx_txt(cover, title, 0.8, 1.6, 11.5, 1.4, bold=True, size=34, color=_C_WHITE_P)
    _pptx_txt(cover, "Deals Intelligence Report", 0.8, 3.0, 8, 0.6,
              size=15, color=RGBColor(0xB0, 0xBE, 0xC5))
    _pptx_txt(cover, f"Generated: {date.today().strftime('%B %d, %Y')}", 0.8, 3.6, 6, 0.45,
              size=12, color=RGBColor(0x78, 0x90, 0x9C))

    for i, (s_title, bullets) in enumerate(slides_data):
        slide = prs.slides.add_slide(blank)
        ac = _accent(s_title.lower())
        _pptx_bg(slide, _C_LTBG_P)
        _pptx_rect(slide, 0, 0, 13.33, 0.9, ac)
        _pptx_txt(slide, s_title.upper(), 0.35, 0.1, 12.5, 0.7, bold=True, size=18, color=_C_WHITE_P)
        _pptx_txt(slide, str(i + 1), 12.5, 0.12, 0.6, 0.6, size=11, color=_C_WHITE_P, align=PP_ALIGN.RIGHT)
        _pptx_rect(slide, 0.3, 1.0, 12.73, 5.9, _C_WHITE_P)
        if bullets:
            txb = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.6))
            txb.word_wrap = True
            tf = txb.text_frame; tf.word_wrap = True
            for j, bullet in enumerate(bullets[:12]):
                p = tf.add_paragraph() if j > 0 else tf.paragraphs[0]
                p.space_before = Pt(4)
                dot = p.add_run(); dot.text = "●  "; dot.font.size = Pt(8); dot.font.color.rgb = ac
                run = p.add_run(); run.text = bullet; run.font.size = Pt(12); run.font.color.rgb = _C_TXT_P
        else:
            _pptx_txt(slide, "No data available.", 0.5, 1.2, 12, 0.5, size=11, color=_C_DIM_P)
        _footer(slide)

    buf = io.BytesIO(); prs.save(buf)
    return buf.getvalue()
